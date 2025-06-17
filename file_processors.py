"""
File processing module for Cyber Assessment Reviewer
Handles extraction of text from various file formats (PDF, DOCX, XLSX, PPTX)
with enhanced cybersecurity-focused text processing and evidence categorization
"""

import logging
import pandas as pd
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import PyPDF2
from docx import Document
from pptx import Presentation

from models import EvidenceFile, ControlData
from config import Config

logger = logging.getLogger(__name__)

class CybersecurityTextProcessor:
    """Enhanced text processing for cybersecurity evidence analysis"""

    # Cybersecurity-relevant keywords and patterns
    SECURITY_KEYWORDS = {
        'access_control': [
            'access control', 'authentication', 'authorization', 'rbac', 'role-based',
            'multi-factor', 'mfa', '2fa', 'single sign-on', 'sso', 'privileged access',
            'least privilege', 'segregation of duties', 'password policy'
        ],
        'encryption': [
            'encryption', 'encrypted', 'cryptography', 'ssl', 'tls', 'aes', 'rsa',
            'certificate', 'key management', 'pki', 'hash', 'digital signature'
        ],
        'monitoring': [
            'monitoring', 'logging', 'audit trail', 'siem', 'log analysis',
            'intrusion detection', 'ids', 'ips', 'security monitoring', 'alerting'
        ],
        'vulnerability': [
            'vulnerability', 'patch management', 'security updates', 'penetration test',
            'vulnerability scan', 'security assessment', 'risk assessment'
        ],
        'incident_response': [
            'incident response', 'incident management', 'security incident',
            'breach response', 'forensics', 'containment', 'recovery'
        ],
        'compliance': [
            'compliance', 'audit', 'policy', 'procedure', 'standard', 'guideline',
            'framework', 'control', 'requirement', 'documentation'
        ],
        'network_security': [
            'firewall', 'network security', 'segmentation', 'dmz', 'vpn',
            'network access control', 'nac', 'perimeter security'
        ],
        'data_protection': [
            'data protection', 'data classification', 'data loss prevention', 'dlp',
            'backup', 'data retention', 'data disposal', 'privacy'
        ]
    }

    @classmethod
    def categorize_text_by_security_domain(cls, text: str) -> Dict[str, List[str]]:
        """Categorize text content by cybersecurity domains"""
        categorized = {domain: [] for domain in cls.SECURITY_KEYWORDS.keys()}

        # Split text into sentences for better context
        sentences = re.split(r'[.!?]+', text)

        for sentence in sentences:
            sentence = sentence.strip().lower()
            if len(sentence) < 10:  # Skip very short sentences
                continue

            for domain, keywords in cls.SECURITY_KEYWORDS.items():
                if any(keyword in sentence for keyword in keywords):
                    categorized[domain].append(sentence.strip())

        # Remove empty categories
        return {k: v for k, v in categorized.items() if v}

    @classmethod
    def extract_technical_details(cls, text: str) -> List[str]:
        """Extract technical details and configurations from text"""
        technical_patterns = [
            r'(?i)(configured?|implemented?|enabled?|disabled?)\s+[^.]{10,100}',
            r'(?i)(version|v\d+\.\d+|build\s+\d+)[^.]{5,50}',
            r'(?i)(port\s+\d+|protocol\s+\w+|algorithm\s+\w+)',
            r'(?i)(certificate|key\s+length|bit\s+encryption)',
            r'(?i)(frequency|interval|threshold|timeout)[^.]{5,50}',
            r'(?i)(server|system|application|database)[^.]{10,100}'
        ]

        technical_details = []
        for pattern in technical_patterns:
            matches = re.findall(pattern, text)
            technical_details.extend([match if isinstance(match, str) else ' '.join(match)
                                    for match in matches])

        return list(set(technical_details))  # Remove duplicates

    @classmethod
    def identify_evidence_types(cls, text: str) -> List[str]:
        """Identify types of evidence present in the text"""
        evidence_patterns = {
            'policy_document': [r'(?i)policy', r'(?i)procedure', r'(?i)standard', r'(?i)guideline'],
            'configuration': [r'(?i)config', r'(?i)setting', r'(?i)parameter', r'(?i)option'],
            'log_data': [r'(?i)log', r'(?i)event', r'(?i)timestamp', r'(?i)audit trail'],
            'test_results': [r'(?i)test', r'(?i)scan', r'(?i)assessment', r'(?i)validation'],
            'training_material': [r'(?i)training', r'(?i)awareness', r'(?i)education', r'(?i)course'],
            'certificate': [r'(?i)certificate', r'(?i)certification', r'(?i)accreditation'],
            'screenshot': [r'(?i)screenshot', r'(?i)image', r'(?i)figure', r'(?i)diagram'],
            'report': [r'(?i)report', r'(?i)analysis', r'(?i)findings', r'(?i)summary']
        }

        identified_types = []
        text_lower = text.lower()

        for evidence_type, patterns in evidence_patterns.items():
            if any(re.search(pattern, text_lower) for pattern in patterns):
                identified_types.append(evidence_type)

        return identified_types

class FileProcessor:
    """Base class for file processors"""
    
    def __init__(self, config: Config = None):
        self.config = config or Config()
    
    def process_file(self, filepath: str) -> Dict[str, str]:
        """Process file and return extracted text sections"""
        raise NotImplementedError("Subclasses must implement process_file method")

class PDFProcessor(FileProcessor):
    """Enhanced processor for PDF files with cybersecurity-focused analysis"""

    def process_file(self, filepath: str) -> Dict[str, str]:
        """Extract text from PDF file with enhanced cybersecurity analysis"""
        processed_sections = {}
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = min(len(pdf_reader.pages), self.config.MAX_PAGES_PDF)

                all_text = ""
                page_texts = {}

                # Extract text from all pages
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        page_texts[f"Page {page_num + 1}"] = text
                        all_text += f"\n{text}"

                # Categorize content by security domains
                security_categories = CybersecurityTextProcessor.categorize_text_by_security_domain(all_text)

                # Add categorized sections
                for domain, sentences in security_categories.items():
                    if sentences:
                        domain_text = '. '.join(sentences[:10])  # Limit to top 10 sentences
                        processed_sections[f"Security Domain - {domain.replace('_', ' ').title()}"] = domain_text[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

                # Add technical details section
                technical_details = CybersecurityTextProcessor.extract_technical_details(all_text)
                if technical_details:
                    processed_sections["Technical Details"] = '. '.join(technical_details[:5])[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

                # Add evidence type identification
                evidence_types = CybersecurityTextProcessor.identify_evidence_types(all_text)
                if evidence_types:
                    processed_sections["Evidence Types"] = f"Identified evidence types: {', '.join(evidence_types)}"

                # Add original page-based sections (limited)
                for page_name, page_text in list(page_texts.items())[:3]:  # Limit to first 3 pages
                    processed_sections[f"Original - {page_name}"] = page_text[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

        except Exception as e:
            logger.error(f"Error reading PDF {filepath}: {e}")
            return {"Error": f"Failed to read PDF: {str(e)}"}

        return processed_sections

class DOCXProcessor(FileProcessor):
    """Enhanced processor for DOCX files with cybersecurity-focused analysis"""

    def process_file(self, filepath: str) -> Dict[str, str]:
        """Extract text from DOCX file with enhanced cybersecurity analysis"""
        processed_sections = {}
        try:
            doc = Document(filepath)

            # Extract all text content
            full_text = []
            headings = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Check if it's a heading
                    if paragraph.style.name.startswith('Heading'):
                        headings.append(paragraph.text.strip())
                    full_text.append(paragraph.text)

            all_text = '\n'.join(full_text)

            # Add document structure information
            if headings:
                processed_sections["Document Structure"] = f"Document headings: {'; '.join(headings[:10])}"

            # Categorize content by security domains
            security_categories = CybersecurityTextProcessor.categorize_text_by_security_domain(all_text)

            # Add categorized sections
            for domain, sentences in security_categories.items():
                if sentences:
                    domain_text = '. '.join(sentences[:8])  # Limit to top 8 sentences
                    processed_sections[f"Security Domain - {domain.replace('_', ' ').title()}"] = domain_text[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

            # Add technical details section
            technical_details = CybersecurityTextProcessor.extract_technical_details(all_text)
            if technical_details:
                processed_sections["Technical Details"] = '. '.join(technical_details[:5])[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

            # Add evidence type identification
            evidence_types = CybersecurityTextProcessor.identify_evidence_types(all_text)
            if evidence_types:
                processed_sections["Evidence Types"] = f"Identified evidence types: {', '.join(evidence_types)}"

            # Add original content sections (limited)
            if full_text:
                section_size = max(1, len(full_text) // 3)  # Divide into ~3 sections
                for i in range(0, min(len(full_text), section_size * 3), section_size):
                    section_text = '\n'.join(full_text[i:i+section_size])
                    processed_sections[f"Original Section {i//section_size + 1}"] = section_text[:self.config.MAX_TEXT_LENGTH_PER_SECTION]

        except Exception as e:
            logger.error(f"Error reading DOCX {filepath}: {e}")
            return {"Error": f"Failed to read DOCX: {str(e)}"}

        return processed_sections

class PPTXProcessor(FileProcessor):
    """Processor for PPTX files"""
    
    def process_file(self, filepath: str) -> Dict[str, str]:
        """Extract text from PPTX file with slide references"""
        slides_text = {}
        try:
            prs = Presentation(filepath)
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text[f"Slide {i + 1}"] = '\n'.join(slide_text)[:self.config.MAX_TEXT_LENGTH_PER_SECTION]
                    
        except Exception as e:
            logger.error(f"Error reading PPTX {filepath}: {e}")
            return {"Error": f"Failed to read PPTX: {str(e)}"}
        
        return slides_text

class XLSXProcessor(FileProcessor):
    """Processor for XLSX files"""
    
    def process_file(self, filepath: str) -> Dict[str, str]:
        """Extract text from XLSX file with sheet references"""
        sheets_text = {}
        try:
            xlsx_file = pd.ExcelFile(filepath)
            
            for sheet_name in xlsx_file.sheet_names[:5]:  # Limit to first 5 sheets
                df = pd.read_excel(xlsx_file, sheet_name=sheet_name)
                
                # Handle NaN values by replacing with empty strings
                df = df.fillna('')
                
                # Convert to string representation
                text_content = f"Sheet: {sheet_name}\n"
                text_content += f"Columns: {', '.join(str(col) for col in df.columns)}\n"
                text_content += f"Rows: {len(df)}\n\n"
                
                # Add sample data
                if len(df) > 0:
                    text_content += df.head(10).to_string()
                    
                sheets_text[sheet_name] = text_content[:self.config.MAX_TEXT_LENGTH_PER_SECTION]
                
        except Exception as e:
            logger.error(f"Error reading XLSX {filepath}: {e}")
            return {"Error": f"Failed to read XLSX: {str(e)}"}
        
        return sheets_text

class AssessmentFileProcessor:
    """Specialized processor for assessment files (typically XLSX)"""
    
    def __init__(self, config: Config = None):
        self.config = config or Config()
    
    def process_assessment_file(self, filepath: str) -> List[ControlData]:
        """Process assessment file and extract control data"""
        try:
            df = pd.read_excel(filepath)
            
            # Standardize column names using flexible mapping
            self._standardize_columns(df)
            
            # Handle NaN values before converting to dict
            df = df.fillna('')
            
            # Convert to ControlData objects
            controls = []
            for _, row in df.iterrows():
                try:
                    control = ControlData.from_dict(row.to_dict())
                    controls.append(control)
                except Exception as e:
                    logger.warning(f"Error processing control row: {e}")
                    continue
            
            logger.info(f"Processed {len(controls)} controls from assessment file")
            return controls
            
        except Exception as e:
            logger.error(f"Error processing assessment file {filepath}: {e}")
            raise
    
    def _standardize_columns(self, df: pd.DataFrame) -> None:
        """Standardize column names using flexible mapping"""
        for standard_name, possible_names in self.config.COLUMN_MAPPING.items():
            for col in df.columns:
                if col.lower().strip() in [name.lower() for name in possible_names]:
                    df.rename(columns={col: standard_name}, inplace=True)
                    break

class FileProcessorFactory:
    """Factory class for creating appropriate file processors"""
    
    _processors = {
        '.pdf': PDFProcessor,
        '.docx': DOCXProcessor,
        '.pptx': PPTXProcessor,
        '.ppt': PPTXProcessor,
        '.xlsx': XLSXProcessor
    }
    
    @classmethod
    def get_processor(cls, file_extension: str, config: Config = None) -> Optional[FileProcessor]:
        """Get appropriate processor for file extension"""
        processor_class = cls._processors.get(file_extension.lower())
        if processor_class:
            return processor_class(config)
        return None
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Get list of supported file extensions"""
        return list(cls._processors.keys())

class EvidenceProcessor:
    """Main processor for evidence files"""
    
    def __init__(self, config: Config = None):
        self.config = config or Config()
    
    def process_evidence_file(self, filepath: str) -> EvidenceFile:
        """Process evidence file and return EvidenceFile object"""
        file_path = Path(filepath)
        file_extension = file_path.suffix.lower()
        
        # Get appropriate processor
        processor = FileProcessorFactory.get_processor(file_extension, self.config)
        if not processor:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        # Process file
        sections = processor.process_file(str(file_path))
        
        # Create EvidenceFile object
        evidence_file = EvidenceFile(
            filename=file_path.name,
            filepath=str(file_path),
            file_type=file_extension,
            size_bytes=file_path.stat().st_size,
            sections=sections
        )
        
        logger.info(f"Processed evidence file: {file_path.name} ({len(sections)} sections)")
        return evidence_file
    
    def process_multiple_evidence_files(self, filepaths: List[str]) -> Dict[str, EvidenceFile]:
        """Process multiple evidence files"""
        evidence_files = {}
        
        for filepath in filepaths:
            try:
                evidence_file = self.process_evidence_file(filepath)
                evidence_files[evidence_file.filename] = evidence_file
            except Exception as e:
                logger.error(f"Error processing evidence file {filepath}: {e}")
                continue
        
        return evidence_files
    
    def extract_combined_evidence_text(self, evidence_files: Dict[str, EvidenceFile], 
                                     max_sections: int = None) -> Dict[str, Dict[str, str]]:
        """Extract combined evidence text in the format expected by the AI backend"""
        max_sections = max_sections or self.config.MAX_EVIDENCE_SECTIONS
        
        combined_evidence = {}
        for filename, evidence_file in evidence_files.items():
            # Filter out error sections and limit sections
            valid_sections = {
                section: text for section, text in evidence_file.sections.items()
                if "Error" not in section and text.strip()
            }
            
            # Limit number of sections
            if len(valid_sections) > max_sections:
                valid_sections = dict(list(valid_sections.items())[:max_sections])
            
            combined_evidence[filename] = valid_sections
        
        return combined_evidence

def is_supported_file_type(filename: str) -> bool:
    """Check if file type is supported"""
    file_extension = Path(filename).suffix.lower()
    return file_extension in FileProcessorFactory.get_supported_extensions()
