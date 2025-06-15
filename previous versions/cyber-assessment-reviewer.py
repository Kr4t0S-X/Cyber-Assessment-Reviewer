# Cyber Assessment Review System with Local LLM
# Main application file: app.py

import os
import json
import torch
import pandas as pd
from flask import Flask, render_template, request, jsonify, session
from werkzeug.utils import secure_filename
from datetime import datetime
import uuid
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
import re
from typing import List, Dict, Any

# Initialize Flask app
app = Flask(__name__)
app.secret_key = 'your-secret-key-here'  # Change this in production
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'xlsx', 'pdf', 'ppt', 'pptx', 'docx'}

# Create necessary directories
Path("uploads").mkdir(exist_ok=True)
Path("sessions").mkdir(exist_ok=True)
Path("static").mkdir(exist_ok=True)
Path("templates").mkdir(exist_ok=True)

class CyberAssessmentReviewer:
    def __init__(self, model_name="microsoft/phi-2"):
        """
        Initialize the LLM for cyber assessment review.
        Using Phi-2 as it's lightweight and can run on a T14S laptop.
        """
        print("Loading model... This may take a few minutes on first run.")
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        
        # Load tokenizer and model
        self.tokenizer = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
        self.model = AutoModelForCausalLM.from_pretrained(
            model_name,
            torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
            device_map="auto",
            trust_remote_code=True
        )
        
        # Set padding token
        self.tokenizer.pad_token = self.tokenizer.eos_token
        
        # Create pipeline
        self.pipe = pipeline(
            "text-generation",
            model=self.model,
            tokenizer=self.tokenizer,
            max_new_tokens=512,
            temperature=0.7,
            do_sample=True,
            top_p=0.9
        )
        
    def extract_text_from_pdf(self, filepath: str) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error reading PDF: {e}")
        return text[:5000]  # Limit to first 5000 chars for performance
    
    def extract_text_from_docx(self, filepath: str) -> str:
        """Extract text from DOCX file"""
        text = ""
        try:
            doc = Document(filepath)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            print(f"Error reading DOCX: {e}")
        return text[:5000]
    
    def extract_text_from_pptx(self, filepath: str) -> str:
        """Extract text from PPTX file"""
        text = ""
        try:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error reading PPTX: {e}")
        return text[:5000]
    
    def extract_text_from_xlsx(self, filepath: str) -> str:
        """Extract text from XLSX file"""
        text = ""
        try:
            df = pd.read_excel(filepath, sheet_name=None)
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string() + "\n\n"
        except Exception as e:
            print(f"Error reading XLSX: {e}")
        return text[:5000]
    
    def extract_evidence_text(self, filepath: str) -> str:
        """Extract text from evidence file based on extension"""
        ext = Path(filepath).suffix.lower()
        if ext == '.pdf':
            return self.extract_text_from_pdf(filepath)
        elif ext == '.docx':
            return self.extract_text_from_docx(filepath)
        elif ext in ['.ppt', '.pptx']:
            return self.extract_text_from_pptx(filepath)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(filepath)
        return ""
    
    def analyze_control_point(self, control: Dict, evidence_text: str) -> Dict:
        """Analyze a single control point with its evidence"""
        prompt = f"""You are a cybersecurity assessment expert. Analyze this control point:

Control: {control.get('control', 'N/A')}
Requirement: {control.get('requirement', 'N/A')}
Supplier Answer: {control.get('answer', 'N/A')}
Evidence Summary: {evidence_text[:1000]}

Please analyze:
1. Does the evidence support the supplier's answer?
2. What are the key findings and risks?
3. What remediation is needed?

Provide a structured response."""

        try:
            response = self.pipe(prompt)[0]['generated_text']
            # Extract the response after the prompt
            analysis = response[len(prompt):].strip()
            
            # Parse the analysis into structured format
            return {
                'control': control.get('control', 'N/A'),
                'evidence_validity': self._extract_validity(analysis),
                'key_findings': self._extract_findings(analysis),
                'risks': self._extract_risks(analysis),
                'remediation': self._extract_remediation(analysis),
                'full_analysis': analysis
            }
        except Exception as e:
            print(f"Error analyzing control: {e}")
            return {
                'control': control.get('control', 'N/A'),
                'evidence_validity': 'Error in analysis',
                'key_findings': [],
                'risks': [],
                'remediation': [],
                'full_analysis': f'Error: {str(e)}'
            }
    
    def _extract_validity(self, text: str) -> str:
        """Extract evidence validity assessment"""
        if any(word in text.lower() for word in ['valid', 'support', 'adequate', 'sufficient']):
            return 'Valid'
        elif any(word in text.lower() for word in ['invalid', 'insufficient', 'inadequate']):
            return 'Invalid'
        return 'Unclear'
    
    def _extract_findings(self, text: str) -> List[str]:
        """Extract key findings from analysis"""
        findings = []
        lines = text.split('\n')
        for line in lines:
            if any(word in line.lower() for word in ['finding', 'found', 'identified', 'discovered']):
                findings.append(line.strip())
        return findings[:3]  # Limit to top 3
    
    def _extract_risks(self, text: str) -> List[str]:
        """Extract risks from analysis"""
        risks = []
        lines = text.split('\n')
        for line in lines:
            if any(word in line.lower() for word in ['risk', 'threat', 'vulnerability', 'exposure']):
                risks.append(line.strip())
        return risks[:3]
    
    def _extract_remediation(self, text: str) -> List[str]:
        """Extract remediation recommendations"""
        remediation = []
        lines = text.split('\n')
        for line in lines:
            if any(word in line.lower() for word in ['recommend', 'should', 'must', 'remediate', 'implement']):
                remediation.append(line.strip())
        return remediation[:3]

# Initialize the reviewer
reviewer = CyberAssessmentReviewer()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/')
def index():
    """Render the main page"""
    return render_template('index.html')

@app.route('/upload_assessment', methods=['POST'])
def upload_assessment():
    """Handle assessment file upload"""
    if 'assessment' not in request.files:
        return jsonify({'error': 'No assessment file provided'}), 400
    
    file = request.files['assessment']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        # Create session ID
        session_id = str(uuid.uuid4())
        session['session_id'] = session_id
        session_path = Path(f"sessions/{session_id}")
        session_path.mkdir(exist_ok=True)
        
        # Save assessment file
        filename = secure_filename(file.filename)
        filepath = session_path / f"assessment_{filename}"
        file.save(str(filepath))
        
        # Parse assessment
        try:
            df = pd.read_excel(filepath)
            controls = df.to_dict('records')
            
            # Save controls to session
            with open(session_path / 'controls.json', 'w') as f:
                json.dump(controls, f)
            
            return jsonify({
                'success': True,
                'session_id': session_id,
                'controls_count': len(controls),
                'columns': list(df.columns)
            })
        except Exception as e:
            return jsonify({'error': f'Error parsing assessment: {str(e)}'}), 400
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/upload_evidence', methods=['POST'])
def upload_evidence():
    """Handle evidence files upload"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    evidence_path = session_path / 'evidence'
    evidence_path.mkdir(exist_ok=True)
    
    uploaded_files = []
    for key in request.files:
        file = request.files[key]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = evidence_path / filename
            file.save(str(filepath))
            uploaded_files.append(filename)
    
    return jsonify({
        'success': True,
        'uploaded_files': uploaded_files
    })

@app.route('/analyze', methods=['POST'])
def analyze():
    """Run the analysis on uploaded assessment and evidence"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    
    # Load controls
    with open(session_path / 'controls.json', 'r') as f:
        controls = json.load(f)
    
    # Extract text from all evidence files
    evidence_texts = {}
    evidence_path = session_path / 'evidence'
    if evidence_path.exists():
        for evidence_file in evidence_path.iterdir():
            if evidence_file.is_file():
                text = reviewer.extract_evidence_text(str(evidence_file))
                evidence_texts[evidence_file.name] = text
    
    # Combine all evidence texts
    combined_evidence = "\n\n".join(evidence_texts.values())
    
    # Analyze each control
    results = []
    for control in controls[:10]:  # Limit to first 10 for demo
        analysis = reviewer.analyze_control_point(control, combined_evidence)
        results.append(analysis)
    
    # Generate summary report
    summary = {
        'total_controls': len(controls),
        'analyzed_controls': len(results),
        'valid_evidence': sum(1 for r in results if r['evidence_validity'] == 'Valid'),
        'invalid_evidence': sum(1 for r in results if r['evidence_validity'] == 'Invalid'),
        'unclear_evidence': sum(1 for r in results if r['evidence_validity'] == 'Unclear'),
        'high_risk_findings': [],
        'key_remediations': []
    }
    
    # Save results
    with open(session_path / 'results.json', 'w') as f:
        json.dump({
            'results': results,
            'summary': summary,
            'timestamp': datetime.now().isoformat()
        }, f, indent=2)
    
    return jsonify({
        'success': True,
        'results': results,
        'summary': summary
    })

@app.route('/download_report/<session_id>')
def download_report(session_id):
    """Download the analysis report"""
    session_path = Path(f"sessions/{session_id}")
    results_path = session_path / 'results.json'
    
    if not results_path.exists():
        return jsonify({'error': 'Report not found'}), 404
    
    with open(results_path, 'r') as f:
        data = json.load(f)
    
    # Create Excel report
    report_path = session_path / 'cyber_assessment_report.xlsx'
    
    # Convert results to DataFrame
    df = pd.DataFrame(data['results'])
    
    # Create Excel writer
    with pd.ExcelWriter(report_path, engine='openpyxl') as writer:
        # Summary sheet
        summary_df = pd.DataFrame([data['summary']])
        summary_df.to_excel(writer, sheet_name='Summary', index=False)
        
        # Detailed results sheet
        df.to_excel(writer, sheet_name='Detailed Analysis', index=False)
    
    return send_file(str(report_path), as_attachment=True)

# Create templates/index.html
html_template = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cyber Assessment Review System</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: #f5f5f5;
            color: #333;
            line-height: 1.6;
        }
        
        .container {
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
        }
        
        header {
            background: #1a73e8;
            color: white;
            padding: 30px 0;
            text-align: center;
            margin-bottom: 40px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }
        
        h1 {
            font-size: 2.5em;
            margin-bottom: 10px;
        }
        
        .subtitle {
            font-size: 1.2em;
            opacity: 0.9;
        }
        
        .upload-section {
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            margin-bottom: 30px;
        }
        
        .upload-section h2 {
            color: #1a73e8;
            margin-bottom: 20px;
        }
        
        .file-input-wrapper {
            position: relative;
            margin-bottom: 20px;
        }
        
        .file-input {
            position: absolute;
            opacity: 0;
            width: 100%;
            height: 100%;
            cursor: pointer;
        }
        
        .file-input-label {
            display: block;
            padding: 12px 20px;
            background: #f0f0f0;
            border: 2px dashed #ccc;
            border-radius: 4px;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s;
        }
        
        .file-input-label:hover {
            background: #e0e0e0;
            border-color: #1a73e8;
        }
        
        .file-list {
            margin-top: 10px;
            padding: 10px;
            background: #f9f9f9;
            border-radius: 4px;
            display: none;
        }
        
        .file-item {
            padding: 5px;
            margin: 5px 0;
            background: white;
            border-radius: 3px;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }
        
        .btn {
            padding: 12px 30px;
            background: #1a73e8;
            color: white;
            border: none;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
            transition: background 0.3s;
        }
        
        .btn:hover {
            background: #1557b0;
        }
        
        .btn:disabled {
            background: #ccc;
            cursor: not-allowed;
        }
        
        .results-section {
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            display: none;
        }
        
        .summary-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }
        
        .summary-card {
            background: #f9f9f9;
            padding: 20px;
            border-radius: 8px;
            text-align: center;
        }
        
        .summary-card h3 {
            color: #666;
            font-size: 14px;
            margin-bottom: 10px;
        }
        
        .summary-card .value {
            font-size: 36px;
            font-weight: bold;
            color: #1a73e8;
        }
        
        .control-result {
            background: #f9f9f9;
            padding: 20px;
            margin: 15px 0;
            border-radius: 8px;
            border-left: 4px solid #1a73e8;
        }
        
        .control-result h4 {
            color: #333;
            margin-bottom: 10px;
        }
        
        .validity-badge {
            display: inline-block;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 14px;
            font-weight: bold;
            margin-bottom: 10px;
        }
        
        .validity-badge.valid {
            background: #4caf50;
            color: white;
        }
        
        .validity-badge.invalid {
            background: #f44336;
            color: white;
        }
        
        .validity-badge.unclear {
            background: #ff9800;
            color: white;
        }
        
        .finding-section {
            margin: 15px 0;
        }
        
        .finding-section h5 {
            color: #666;
            margin-bottom: 5px;
        }
        
        .finding-list {
            list-style: none;
            padding-left: 20px;
        }
        
        .finding-list li {
            margin: 5px 0;
            position: relative;
        }
        
        .finding-list li:before {
            content: "•";
            position: absolute;
            left: -15px;
            color: #1a73e8;
        }
        
        .loading {
            display: none;
            text-align: center;
            padding: 40px;
        }
        
        .spinner {
            border: 4px solid #f3f3f3;
            border-top: 4px solid #1a73e8;
            border-radius: 50%;
            width: 40px;
            height: 40px;
            animation: spin 2s linear infinite;
            margin: 0 auto 20px;
        }
        
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        
        .error-message {
            background: #fee;
            color: #c33;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
        }
        
        .success-message {
            background: #efe;
            color: #3c3;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
        }
    </style>
</head>
<body>
    <header>
        <h1>Cyber Assessment Review System</h1>
        <p class="subtitle">AI-Powered Security Control Validation</p>
    </header>
    
    <div class="container">
        <!-- Upload Assessment Section -->
        <div class="upload-section">
            <h2>Step 1: Upload Assessment File</h2>
            <div class="file-input-wrapper">
                <input type="file" id="assessmentFile" class="file-input" accept=".xlsx">
                <label for="assessmentFile" class="file-input-label">
                    Click to select assessment file (.xlsx)
                </label>
            </div>
            <div id="assessmentFileList" class="file-list"></div>
        </div>
        
        <!-- Upload Evidence Section -->
        <div class="upload-section">
            <h2>Step 2: Upload Evidence Files</h2>
            <div class="file-input-wrapper">
                <input type="file" id="evidenceFiles" class="file-input" multiple accept=".pdf,.ppt,.pptx,.xlsx,.docx">
                <label for="evidenceFiles" class="file-input-label">
                    Click to select evidence files (PDF, PPT, XLSX, DOCX)
                </label>
            </div>
            <div id="evidenceFileList" class="file-list"></div>
        </div>
        
        <!-- Analyze Button -->
        <div style="text-align: center; margin: 30px 0;">
            <button id="analyzeBtn" class="btn" disabled>Analyze Assessment</button>
        </div>
        
        <!-- Loading Section -->
        <div id="loadingSection" class="loading">
            <div class="spinner"></div>
            <p>Analyzing assessment and evidence files...</p>
            <p style="color: #666; font-size: 14px; margin-top: 10px;">This may take a few minutes</p>
        </div>
        
        <!-- Error/Success Messages -->
        <div id="errorMessage" class="error-message"></div>
        <div id="successMessage" class="success-message"></div>
        
        <!-- Results Section -->
        <div id="resultsSection" class="results-section">
            <h2>Analysis Results</h2>
            
            <div id="summaryGrid" class="summary-grid"></div>
            
            <h3 style="margin: 30px 0 20px;">Detailed Control Analysis</h3>
            <div id="detailedResults"></div>
            
            <div style="text-align: center; margin-top: 30px;">
                <button id="downloadReportBtn" class="btn">Download Full Report</button>
            </div>
        </div>
    </div>
    
    <script>
        let sessionId = null;
        let assessmentUploaded = false;
        let evidenceUploaded = false;
        
        // File upload handlers
        document.getElementById('assessmentFile').addEventListener('change', async (e) => {
            const file = e.target.files[0];
            if (!file) return;
            
            const formData = new FormData();
            formData.append('assessment', file);
            
            try {
                const response = await fetch('/upload_assessment', {
                    method: 'POST',
                    body: formData
                });
                
                const data = await response.json();
                
                if (data.success) {
                    sessionId = data.session_id;
                    assessmentUploaded = true;
                    
                    document.getElementById('assessmentFileList').style.display = 'block';
                    document.getElementById('assessmentFileList').innerHTML = `
                        <div class="file-item">
                            <span>${file.name}</span>
                            <span style="color: #4caf50;">✓ ${data.controls_count} controls found</span>
                        </div>
                    `;
                    
                    checkAnalyzeButton();
                } else {
                    showError(data.error || 'Failed to upload assessment');
                }
            } catch (error) {
                showError('Error uploading assessment: ' + error.message);
            }
        });
        
        document.getElementById('evidenceFiles').addEventListener('change', async (e) => {
            const files = e.target.files;
            if (!files.length) return;
            
            if (!sessionId) {
                showError('Please upload assessment file first');
                return;
            }
            
            const formData = new FormData();
            for (let i = 0; i < files.length; i++) {
                formData.append(`evidence_${i}`, files[i]);
            }
            
            try {
                const response = await fetch('/upload_evidence', {
                    method: 'POST',
                    body: formData
                });
                
                const data = await response.json();
                
                if (data.success) {
                    evidenceUploaded = true;
                    
                    const fileListHtml = data.uploaded_files.map(filename => 
                        `<div class="file-item">
                            <span>${filename}</span>
                            <span style="color: #4caf50;">✓</span>
                        </div>`
                    ).join('');
                    
                    document.getElementById('evidenceFileList').style.display = 'block';
                    document.getElementById('evidenceFileList').innerHTML = fileListHtml;
                    
                    checkAnalyzeButton();
                } else {
                    showError(data.error || 'Failed to upload evidence files');
                }
            } catch (error) {
                showError('Error uploading evidence: ' + error.message);
            }
        });
        
        // Analyze button handler
        document.getElementById('analyzeBtn').addEventListener('click', async () => {
            document.getElementById('loadingSection').style.display = 'block';
            document.getElementById('resultsSection').style.display = 'none';
            hideMessages();
            
            try {
                const response = await fetch('/analyze', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    }
                });
                
                const data = await response.json();
                
                if (data.success) {
                    displayResults(data);
                    showSuccess('Analysis completed successfully!');
                } else {
                    showError(data.error || 'Analysis failed');
                }
            } catch (error) {
                showError('Error during analysis: ' + error.message);
            } finally {
                document.getElementById('loadingSection').style.display = 'none';
            }
        });
        
        // Download report handler
        document.getElementById('downloadReportBtn').addEventListener('click', () => {
            if (sessionId) {
                window.location.href = `/download_report/${sessionId}`;
            }
        });
        
        // Helper functions
        function checkAnalyzeButton() {
            const btn = document.getElementById('analyzeBtn');
            btn.disabled = !(assessmentUploaded && evidenceUploaded);
        }
        
        function showError(message) {
            const errorDiv = document.getElementById('errorMessage');
            errorDiv.textContent = message;
            errorDiv.style.display = 'block';
            setTimeout(() => errorDiv.style.display = 'none', 5000);
        }
        
        function showSuccess(message) {
            const successDiv = document.getElementById('successMessage');
            successDiv.textContent = message;
            successDiv.style.display = 'block';
            setTimeout(() => successDiv.style.display = 'none', 5000);
        }
        
        function hideMessages() {
            document.getElementById('errorMessage').style.display = 'none';
            document.getElementById('successMessage').style.display = 'none';
        }
        
        function displayResults(data) {
            // Display summary
            const summaryHtml = `
                <div class="summary-card">
                    <h3>Total Controls</h3>
                    <div class="value">${data.summary.total_controls}</div>
                </div>
                <div class="summary-card">
                    <h3>Analyzed</h3>
                    <div class="value">${data.summary.analyzed_controls}</div>
                </div>
                <div class="summary-card">
                    <h3>Valid Evidence</h3>
                    <div class="value" style="color: #4caf50;">${data.summary.valid_evidence}</div>
                </div>
                <div class="summary-card">
                    <h3>Invalid Evidence</h3>
                    <div class="value" style="color: #f44336;">${data.summary.invalid_evidence}</div>
                </div>
            `;
            document.getElementById('summaryGrid').innerHTML = summaryHtml;
            
            // Display detailed results
            const resultsHtml = data.results.map(result => `
                <div class="control-result">
                    <h4>${result.control}</h4>
                    <span class="validity-badge ${result.evidence_validity.toLowerCase()}">${result.evidence_validity}</span>
                    
                    ${result.key_findings.length > 0 ? `
                        <div class="finding-section">
                            <h5>Key Findings:</h5>
                            <ul class="finding-list">
                                ${result.key_findings.map(f => `<li>${f}</li>`).join('')}
                            </ul>
                        </div>
                    ` : ''}
                    
                    ${result.risks.length > 0 ? `
                        <div class="finding-section">
                            <h5>Identified Risks:</h5>
                            <ul class="finding-list">
                                ${result.risks.map(r => `<li>${r}</li>`).join('')}
                            </ul>
                        </div>
                    ` : ''}
                    
                    ${result.remediation.length > 0 ? `
                        <div class="finding-section">
                            <h5>Remediation Recommendations:</h5>
                            <ul class="finding-list">
                                ${result.remediation.map(r => `<li>${r}</li>`).join('')}
                            </ul>
                        </div>
                    ` : ''}
                </div>
            `).join('');
            
            document.getElementById('detailedResults').innerHTML = resultsHtml;
            document.getElementById('resultsSection').style.display = 'block';
        }
    </script>
</body>
</html>
"""

# Save the HTML template
with open('templates/index.html', 'w') as f:
    f.write(html_template)

# Requirements file content
requirements_content = """flask==3.0.0
torch==2.1.0
transformers==4.36.0
pandas==2.1.4
openpyxl==3.1.2
PyPDF2==3.0.1
python-docx==1.1.0
python-pptx==0.6.23
werkzeug==3.0.1
accelerate==0.25.0
sentencepiece==0.1.99
"""

# Save requirements.txt
with open('requirements.txt', 'w') as f:
    f.write(requirements_content)

# Installation and run instructions
instructions = """
# Cyber Assessment Review System - Setup Instructions

## Installation Steps:

1. Create a virtual environment:
   ```bash
   python -m venv venv
   source venv/bin/activate  # On Windows: venv\\Scripts\\activate
   ```

2. Install dependencies:
   ```bash
   pip install -r requirements.txt
   ```

3. Run the application:
   ```bash
   python app.py
   ```

4. Open your browser and go to: http://localhost:5000

## Usage:

1. Upload your cyber assessment Excel file (.xlsx format)
2. Upload evidence files (PDF, PPT, DOCX, XLSX)
3. Click "Analyze Assessment" to run the AI review
4. View results and download the detailed report

## System Requirements:

- Python 3.8+
- At least 8GB RAM (16GB recommended)
- ~3GB disk space for the model
- CPU or GPU (GPU recommended for faster processing)

## Notes:

- First run will download the model (~1.5GB)
- The system uses Microsoft Phi-2 model, optimized for laptops
- Analysis is limited to first 10 controls for demo purposes
- Modify the limit in the analyze() function for full assessment

## Customization:

To use a different model, change the model_name parameter in CyberAssessmentReviewer:
- For smaller model: "microsoft/phi-1_5"
- For larger model: "mistralai/Mistral-7B-Instruct-v0.2" (requires more RAM)
"""

# Save instructions
with open('SETUP_INSTRUCTIONS.md', 'w') as f:
    f.write(instructions)

if __name__ == '__main__':
    print("Starting Cyber Assessment Review System...")
    print("This may take a few minutes on first run to download the AI model.")
    app.run(debug=True, host='0.0.0.0', port=5000)
                    