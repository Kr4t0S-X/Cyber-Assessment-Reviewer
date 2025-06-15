# Cyber Assessment Review System with Mistral 7B
# Main application file: app.py

import sys
import subprocess
import importlib.util
import os

# Check and install dependencies
def check_and_install_dependencies():
    """Check for required dependencies and install missing ones"""
    
    # Check if we should skip dependency checking
    if os.environ.get('SKIP_DEP_CHECK', '').lower() in ['1', 'true', 'yes']:
        print("Skipping dependency check (SKIP_DEP_CHECK is set)")
        return
    
    print("Checking dependencies...")
    
    # Core dependencies that are always needed
    core_dependencies = {
        'flask': 'flask==3.0.0',
        'pandas': 'pandas==2.1.4',
        'openpyxl': 'openpyxl==3.1.2',
        'PyPDF2': 'PyPDF2==3.0.1',
        'docx': 'python-docx==1.1.0',
        'pptx': 'python-pptx==0.6.23',
        'werkzeug': 'werkzeug==3.0.1',
        'requests': 'requests==2.31.0',
        'numpy': 'numpy==1.24.3'
    }
    
    # Optional dependencies for Transformers mode
    transformers_dependencies = {
        'torch': 'torch==2.1.0',
        'transformers': 'transformers==4.36.0',
        'accelerate': 'accelerate==0.25.0',
        'sentencepiece': 'sentencepiece==0.1.99',
        'bitsandbytes': 'bitsandbytes==0.41.3',
        'scipy': 'scipy==1.11.4',
        'safetensors': 'safetensors==0.4.1'
    }
    
    missing_core = []
    missing_transformers = []
    
    # Check core dependencies
    for module, package in core_dependencies.items():
        if importlib.util.find_spec(module) is None:
            missing_core.append(package)
    
    # Check if Ollama is available
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        ollama_available = response.status_code == 200
    except:
        ollama_available = False
    
    # If Ollama is not available, check Transformers dependencies
    if not ollama_available:
        print("Ollama not detected. Checking Transformers dependencies...")
        for module, package in transformers_dependencies.items():
            if importlib.util.find_spec(module) is None:
                missing_transformers.append(package)
    
    # Install missing packages
    if missing_core or missing_transformers:
        print("\nMissing dependencies detected. Installing...")
        
        # Upgrade pip first
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--upgrade", "pip"])
        
        # Install missing core dependencies
        if missing_core:
            print(f"\nInstalling core dependencies: {', '.join(missing_core)}")
            for package in missing_core:
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                    print(f"✓ Installed {package}")
                except subprocess.CalledProcessError as e:
                    print(f"✗ Failed to install {package}: {e}")
                    sys.exit(1)
        
        # Install missing Transformers dependencies if needed
        if missing_transformers and not ollama_available:
            print(f"\nInstalling Transformers dependencies: {', '.join(missing_transformers)}")
            
            # Special handling for PyTorch
            if 'torch' in [pkg.split('==')[0] for pkg in missing_transformers]:
                print("\nInstalling PyTorch... This may take a few minutes.")
                try:
                    # Try to install PyTorch with CUDA support
                    if sys.platform.startswith('linux') or sys.platform == 'win32':
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio", 
                            "--index-url", "https://download.pytorch.org/whl/cu118"
                        ])
                    else:
                        # macOS or other platforms
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio"
                        ])
                    print("✓ Installed PyTorch")
                except subprocess.CalledProcessError:
                    print("✗ Failed to install PyTorch with CUDA, trying CPU-only version...")
                    try:
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio"
                        ])
                        print("✓ Installed PyTorch (CPU-only)")
                    except subprocess.CalledProcessError as e:
                        print(f"✗ Failed to install PyTorch: {e}")
                        sys.exit(1)
            
            # Install other Transformers dependencies
            for package in missing_transformers:
                if not package.startswith('torch'):
                    try:
                        subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                        print(f"✓ Installed {package}")
                    except subprocess.CalledProcessError as e:
                        print(f"✗ Failed to install {package}: {e}")
                        # Don't exit for optional dependencies
        
        print("\n✅ All required dependencies installed!")
        print("Please restart the application for changes to take effect.")
        sys.exit(0)
    else:
        if ollama_available:
            print("✅ All core dependencies satisfied. Using Ollama backend.")
        else:
            print("✅ All dependencies satisfied.")

# Run dependency check before imports
check_and_install_dependencies()

# Now import everything else
import json
import pandas as pd
from flask import Flask, render_template, request, jsonify, session, send_file
from werkzeug.utils import secure_filename
from datetime import datetime
import uuid
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import re
from typing import List, Dict, Any, Optional
import logging
from dataclasses import dataclass
import numpy as np
import requests
import subprocess
import time

# Optional imports for Transformers mode
try:
    import torch
    from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig, pipeline
    TRANSFORMERS_AVAILABLE = True
except ImportError:
    TRANSFORMERS_AVAILABLE = False
    print("Transformers not available. Using Ollama mode only.")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)
app.secret_key = 'your-secret-key-here'  # Change this in production
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'xlsx', 'pdf', 'ppt', 'pptx', 'docx'}

# Create necessary directories
for dir_name in ['uploads', 'sessions', 'static', 'models']:
    Path(dir_name).mkdir(exist_ok=True)

@dataclass
class AssessmentResult:
    """Data class for assessment results"""
    control_id: str
    control_name: str
    requirement: str
    supplier_answer: str
    evidence_validity: str
    compliance_status: str
    risk_level: str
    key_findings: List[str]
    risks: List[str]
    remediation: List[str]
    evidence_references: List[str]
    confidence_score: float

class CyberAssessmentReviewer:
    def __init__(self, model_name="mistral:7b-instruct", use_ollama=True):
        """
        Initialize with Mistral 7B for optimal cyber assessment analysis
        Supports both Ollama (recommended) and Transformers backends
        """
        logger.info(f"Initializing Cyber Assessment Reviewer")
        self.use_ollama = use_ollama
        self.model_name = model_name
        
        if self.use_ollama:
            # Check if Ollama is running
            if self.check_ollama_running():
                logger.info("Ollama detected and running")
                # Check if model is available
                if not self.check_ollama_model():
                    logger.info(f"Model {model_name} not found. Pulling from Ollama...")
                    self.pull_ollama_model()
                logger.info(f"Using Ollama with model: {model_name}")
            else:
                logger.warning("Ollama not running. Falling back to Transformers mode.")
                self.use_ollama = False
        
        if not self.use_ollama:
            if not TRANSFORMERS_AVAILABLE:
                raise RuntimeError("Neither Ollama nor Transformers is available. Please install Ollama or transformers library.")
            
            # Transformers initialization (original code)
            self.device = "cuda" if (TRANSFORMERS_AVAILABLE and torch.cuda.is_available()) else "cpu"
            logger.info(f"Using device: {self.device}")
            
            # Use HuggingFace model name for Transformers
            hf_model_name = "mistralai/Mistral-7B-Instruct-v0.3"
            
            # Configure quantization for efficient memory usage
            if self.device == "cuda":
                self.quantization_config = BitsAndBytesConfig(
                    load_in_4bit=True,
                    bnb_4bit_compute_dtype=torch.float16,
                    bnb_4bit_quant_type="nf4",
                    bnb_4bit_use_double_quant=True
                )
                logger.info("Using 4-bit quantization for efficient memory usage")
            else:
                self.quantization_config = None
            
            # Load tokenizer
            self.tokenizer = AutoTokenizer.from_pretrained(hf_model_name)
            if not self.tokenizer.pad_token:
                self.tokenizer.pad_token = self.tokenizer.eos_token
            
            # Load model with optimizations
            logger.info("Loading Mistral 7B model... This may take a few minutes.")
            self.model = AutoModelForCausalLM.from_pretrained(
                hf_model_name,
                quantization_config=self.quantization_config,
                torch_dtype=torch.float16 if self.device == "cuda" else torch.float32,
                device_map="auto",
                trust_remote_code=True,
                low_cpu_mem_usage=True
            )
            
            # Create pipeline with optimized settings
            self.pipe = pipeline(
                "text-generation",
                model=self.model,
                tokenizer=self.tokenizer,
                max_new_tokens=1024,
                temperature=0.3,
                do_sample=True,
                top_p=0.9,
                repetition_penalty=1.1
            )
        
        # Cyber security frameworks for context
        self.frameworks = {
            "NIST": "NIST Cybersecurity Framework",
            "ISO27001": "ISO/IEC 27001:2022",
            "SOC2": "SOC 2 Type II",
            "CIS": "CIS Controls v8",
            "PCI-DSS": "PCI-DSS v4.0"
        }
        
        logger.info("Model initialization complete")
    
    def check_ollama_running(self) -> bool:
        """Check if Ollama server is running"""
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            return response.status_code == 200
        except:
            return False
    
    def check_ollama_model(self) -> bool:
        """Check if the model is available in Ollama"""
        try:
            response = requests.get("http://localhost:11434/api/tags")
            if response.status_code == 200:
                models = response.json().get("models", [])
                return any(model["name"] == self.model_name for model in models)
        except:
            pass
        return False
    
    def pull_ollama_model(self):
        """Pull model using Ollama"""
        try:
            # Try to pull the model
            logger.info(f"Pulling {self.model_name} - this may take several minutes...")
            response = requests.post(
                "http://localhost:11434/api/pull",
                json={"name": self.model_name},
                stream=True
            )
            
            for line in response.iter_lines():
                if line:
                    data = json.loads(line)
                    if "status" in data:
                        logger.info(f"Ollama: {data['status']}")
                        
            logger.info(f"Model {self.model_name} pulled successfully")
        except Exception as e:
            logger.error(f"Failed to pull model: {e}")
            raise RuntimeError(f"Could not pull model {self.model_name}. Please run: ollama pull {self.model_name}")
    
    def generate_with_ollama(self, prompt: str, temperature: float = 0.3) -> str:
        """Generate response using Ollama API"""
        try:
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": self.model_name,
                    "prompt": prompt,
                    "temperature": temperature,
                    "stream": False,
                    "options": {
                        "num_predict": 1024,
                        "top_p": 0.9,
                        "repeat_penalty": 1.1
                    }
                }
            )
            
            if response.status_code == 200:
                return response.json()["response"]
            else:
                logger.error(f"Ollama API error: {response.status_code}")
                return ""
        except Exception as e:
            logger.error(f"Error generating with Ollama: {e}")
            return ""
    
    def extract_text_from_pdf(self, filepath: str, max_pages: int = 50) -> Dict[str, str]:
        """Extract text from PDF file with page references"""
        pages_text = {}
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = min(len(pdf_reader.pages), max_pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        pages_text[f"Page {page_num + 1}"] = text[:2000]  # Limit per page
        except Exception as e:
            logger.error(f"Error reading PDF {filepath}: {e}")
            return {"Error": f"Failed to read PDF: {str(e)}"}
        
        return pages_text
    
    def extract_text_from_docx(self, filepath: str) -> Dict[str, str]:
        """Extract text from DOCX file with section references"""
        sections_text = {}
        try:
            doc = Document(filepath)
            
            # Extract paragraphs
            full_text = []
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
                    
            # Group into sections
            if full_text:
                section_size = max(1, len(full_text) // 5)  # Divide into ~5 sections
                for i in range(0, len(full_text), section_size):
                    section_text = '\n'.join(full_text[i:i+section_size])
                    sections_text[f"Section {i//section_size + 1}"] = section_text[:2000]
                    
        except Exception as e:
            logger.error(f"Error reading DOCX {filepath}: {e}")
            return {"Error": f"Failed to read DOCX: {str(e)}"}
        
        return sections_text
    
    def extract_text_from_pptx(self, filepath: str) -> Dict[str, str]:
        """Extract text from PPTX file with slide references"""
        slides_text = {}
        try:
            prs = Presentation(filepath)
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text[f"Slide {i + 1}"] = '\n'.join(slide_text)[:2000]
                    
        except Exception as e:
            logger.error(f"Error reading PPTX {filepath}: {e}")
            return {"Error": f"Failed to read PPTX: {str(e)}"}
        
        return slides_text
    
    def extract_text_from_xlsx(self, filepath: str) -> Dict[str, str]:
        """Extract text from XLSX file with sheet references"""
        sheets_text = {}
        try:
            xlsx_file = pd.ExcelFile(filepath)
            
            for sheet_name in xlsx_file.sheet_names[:5]:  # Limit to first 5 sheets
                df = pd.read_excel(xlsx_file, sheet_name=sheet_name)
                
                # Convert to string representation
                text_content = f"Sheet: {sheet_name}\n"
                text_content += f"Columns: {', '.join(df.columns)}\n"
                text_content += f"Rows: {len(df)}\n\n"
                
                # Add sample data
                if len(df) > 0:
                    text_content += df.head(10).to_string()
                    
                sheets_text[sheet_name] = text_content[:2000]
                
        except Exception as e:
            logger.error(f"Error reading XLSX {filepath}: {e}")
            return {"Error": f"Failed to read XLSX: {str(e)}"}
        
        return sheets_text
    
    def extract_evidence_text(self, filepath: str) -> Dict[str, str]:
        """Extract text from evidence file with source references"""
        ext = Path(filepath).suffix.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(filepath)
        elif ext == '.docx':
            return self.extract_text_from_docx(filepath)
        elif ext in ['.ppt', '.pptx']:
            return self.extract_text_from_pptx(filepath)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(filepath)
        
        return {"Error": f"Unsupported file type: {ext}"}
    
    def create_cyber_prompt(self, control: Dict, evidence_text: str, framework: str = "NIST") -> str:
        """Create specialized prompt for cyber assessment analysis"""
        prompt = f"""You are a senior cybersecurity compliance auditor reviewing control assessments. 
Analyze the following control against the provided evidence using the {self.frameworks.get(framework, framework)} framework.

CONTROL INFORMATION:
- Control ID: {control.get('control_id', control.get('control', 'N/A'))}
- Control Name: {control.get('control_name', control.get('name', 'N/A'))}
- Requirement: {control.get('requirement', control.get('description', 'N/A'))}
- Supplier Answer: {control.get('answer', control.get('response', 'N/A'))}
- Implementation Status: {control.get('status', 'Unknown')}

EVIDENCE PROVIDED:
{evidence_text[:3000]}

Provide your analysis in the following JSON format:
{{
    "evidence_validity": "Valid|Partially Valid|Invalid|No Evidence",
    "compliance_status": "Compliant|Partially Compliant|Non-Compliant",
    "risk_level": "Critical|High|Medium|Low",
    "confidence_score": 0.0-1.0,
    "key_findings": ["finding1", "finding2", "finding3"],
    "identified_risks": ["risk1", "risk2", "risk3"],
    "remediation_steps": ["step1", "step2", "step3"],
    "evidence_gaps": ["gap1", "gap2"]
}}

Focus on:
1. Whether evidence adequately demonstrates the control implementation
2. Specific security vulnerabilities or compliance gaps
3. Practical, prioritized remediation recommendations
4. Risk impact on the organization

Provide only the JSON response without additional text."""

        return prompt
    
    def parse_llm_response(self, response: str) -> Dict[str, Any]:
        """Parse and validate LLM JSON response"""
        try:
            # Extract JSON from response
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                json_str = json_match.group()
                result = json.loads(json_str)
                
                # Ensure all required fields exist
                default_result = {
                    "evidence_validity": "Invalid",
                    "compliance_status": "Non-Compliant",
                    "risk_level": "High",
                    "confidence_score": 0.5,
                    "key_findings": [],
                    "identified_risks": [],
                    "remediation_steps": [],
                    "evidence_gaps": []
                }
                
                # Merge with defaults
                for key, default_value in default_result.items():
                    if key not in result:
                        result[key] = default_value
                
                return result
            else:
                raise ValueError("No JSON found in response")
                
        except Exception as e:
            logger.error(f"Error parsing LLM response: {e}")
            logger.debug(f"Response was: {response}")
            
            # Return structured default response
            return {
                "evidence_validity": "Error",
                "compliance_status": "Unknown",
                "risk_level": "Unknown",
                "confidence_score": 0.0,
                "key_findings": ["Error parsing AI response"],
                "identified_risks": ["Unable to assess risks due to parsing error"],
                "remediation_steps": ["Review control manually"],
                "evidence_gaps": ["Technical error in assessment"]
            }
    
    def analyze_control_point(self, control: Dict, evidence_texts: Dict[str, Dict[str, str]], 
                            framework: str = "NIST") -> AssessmentResult:
        """Analyze a single control point with its evidence"""
        # Combine evidence from all files
        combined_evidence = []
        evidence_references = []
        
        for filename, file_sections in evidence_texts.items():
            for section, text in file_sections.items():
                if "Error" not in section and text.strip():
                    combined_evidence.append(f"[{filename} - {section}]:\n{text}")
                    evidence_references.append(f"{filename} - {section}")
        
        evidence_text = "\n\n".join(combined_evidence[:5])  # Limit to top 5 sections
        
        # Create prompt
        prompt = self.create_cyber_prompt(control, evidence_text, framework)
        
        try:
            # Generate response based on backend
            if self.use_ollama:
                generated_text = self.generate_with_ollama(prompt, temperature=0.3)
            else:
                # Use Transformers pipeline
                response = self.pipe(prompt, max_new_tokens=800, temperature=0.3)[0]['generated_text']
                # Extract only the generated part
                generated_text = response[len(prompt):].strip()
            
            # Parse response
            parsed = self.parse_llm_response(generated_text)
            
            # Create AssessmentResult
            result = AssessmentResult(
                control_id=control.get('control_id', control.get('control', 'N/A')),
                control_name=control.get('control_name', control.get('name', 'N/A')),
                requirement=control.get('requirement', control.get('description', 'N/A')),
                supplier_answer=control.get('answer', control.get('response', 'N/A')),
                evidence_validity=parsed['evidence_validity'],
                compliance_status=parsed['compliance_status'],
                risk_level=parsed['risk_level'],
                key_findings=parsed.get('key_findings', [])[:3],
                risks=parsed.get('identified_risks', [])[:3],
                remediation=parsed.get('remediation_steps', [])[:3],
                evidence_references=evidence_references[:3],
                confidence_score=float(parsed.get('confidence_score', 0.5))
            )
            
            return result
            
        except Exception as e:
            logger.error(f"Error analyzing control {control.get('control_id', 'Unknown')}: {e}")
            
            # Return error result
            return AssessmentResult(
                control_id=control.get('control_id', 'Unknown'),
                control_name=control.get('control_name', 'Unknown'),
                requirement=control.get('requirement', 'Unknown'),
                supplier_answer=control.get('answer', 'Unknown'),
                evidence_validity="Error",
                compliance_status="Unknown",
                risk_level="Unknown",
                key_findings=[f"Analysis error: {str(e)}"],
                risks=["Unable to assess due to error"],
                remediation=["Manual review required"],
                evidence_references=[],
                confidence_score=0.0
            )
    
    def calculate_risk_score(self, results: List[AssessmentResult]) -> Dict[str, Any]:
        """Calculate overall risk metrics from assessment results"""
        if not results:
            return {
                "overall_risk_score": 0,
                "risk_distribution": {},
                "compliance_percentage": 0,
                "high_priority_controls": []
            }
        
        # Count by risk level
        risk_counts = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0, "Unknown": 0}
        compliance_counts = {"Compliant": 0, "Partially Compliant": 0, "Non-Compliant": 0}
        
        for result in results:
            risk_counts[result.risk_level] = risk_counts.get(result.risk_level, 0) + 1
            compliance_counts[result.compliance_status] = compliance_counts.get(result.compliance_status, 0) + 1
        
        # Calculate risk score (0-100, higher is worse)
        risk_weights = {"Critical": 10, "High": 7, "Medium": 4, "Low": 1, "Unknown": 5}
        total_risk = sum(risk_counts[level] * risk_weights.get(level, 0) for level in risk_counts)
        max_risk = len(results) * 10  # If all were critical
        risk_score = (total_risk / max_risk * 100) if max_risk > 0 else 0
        
        # Calculate compliance percentage
        compliant = compliance_counts.get("Compliant", 0)
        total = len(results)
        compliance_percentage = (compliant / total * 100) if total > 0 else 0
        
        # Identify high priority controls
        high_priority = [
            {
                "control_id": r.control_id,
                "control_name": r.control_name,
                "risk_level": r.risk_level,
                "compliance_status": r.compliance_status
            }
            for r in results 
            if r.risk_level in ["Critical", "High"] and r.compliance_status == "Non-Compliant"
        ][:10]  # Top 10
        
        return {
            "overall_risk_score": round(risk_score, 1),
            "risk_distribution": risk_counts,
            "compliance_distribution": compliance_counts,
            "compliance_percentage": round(compliance_percentage, 1),
            "high_priority_controls": high_priority,
            "total_controls_assessed": len(results)
        }

# Initialize the reviewer
try:
    # Try Ollama first (recommended)
    reviewer = CyberAssessmentReviewer(use_ollama=True)
except Exception as e:
    logger.warning(f"Failed to initialize with Ollama: {e}")
    logger.info("Falling back to Transformers mode...")
    try:
        reviewer = CyberAssessmentReviewer(use_ollama=False)
    except Exception as e2:
        logger.error(f"Failed to initialize reviewer: {e2}")
        logger.error("Please install Ollama (recommended) or required Python packages")
        raise

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/')
def index():
    """Render the main page"""
    backend = "Ollama" if reviewer.use_ollama else "Transformers"
    return render_template('index.html', backend=backend)

@app.route('/upload_assessment', methods=['POST'])
def upload_assessment():
    """Handle assessment file upload"""
    if 'assessment' not in request.files:
        return jsonify({'error': 'No assessment file provided'}), 400
    
    file = request.files['assessment']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        # Create session ID
        session_id = str(uuid.uuid4())
        session['session_id'] = session_id
        session_path = Path(f"sessions/{session_id}")
        session_path.mkdir(exist_ok=True)
        
        # Save assessment file
        filename = secure_filename(file.filename)
        filepath = session_path / f"assessment_{filename}"
        file.save(str(filepath))
        
        # Parse assessment
        try:
            df = pd.read_excel(filepath)
            
            # Detect column names (flexible mapping)
            column_mapping = {
                'control': ['control', 'control_id', 'id', 'control id', 'control_number'],
                'name': ['name', 'control_name', 'title', 'control name', 'description'],
                'requirement': ['requirement', 'requirements', 'description', 'control_description'],
                'answer': ['answer', 'response', 'supplier_answer', 'supplier response', 'implementation'],
                'status': ['status', 'implementation_status', 'compliance_status']
            }
            
            # Standardize column names
            for standard_name, possible_names in column_mapping.items():
                for col in df.columns:
                    if col.lower().strip() in [name.lower() for name in possible_names]:
                        df.rename(columns={col: standard_name}, inplace=True)
                        break
            
            controls = df.to_dict('records')
            
            # Save controls to session
            with open(session_path / 'controls.json', 'w') as f:
                json.dump(controls, f)
            
            # Save column info
            session['columns'] = list(df.columns)
            session['framework'] = request.form.get('framework', 'NIST')
            
            return jsonify({
                'success': True,
                'session_id': session_id,
                'controls_count': len(controls),
                'columns': list(df.columns),
                'sample_control': controls[0] if controls else None
            })
        except Exception as e:
            logger.error(f"Error parsing assessment: {e}")
            return jsonify({'error': f'Error parsing assessment: {str(e)}'}), 400
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/upload_evidence', methods=['POST'])
def upload_evidence():
    """Handle evidence files upload"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    evidence_path = session_path / 'evidence'
    evidence_path.mkdir(exist_ok=True)
    
    uploaded_files = []
    for key in request.files:
        file = request.files[key]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = evidence_path / filename
            file.save(str(filepath))
            uploaded_files.append({
                'filename': filename,
                'size': os.path.getsize(filepath)
            })
    
    return jsonify({
        'success': True,
        'uploaded_files': uploaded_files,
        'total_files': len(uploaded_files)
    })

@app.route('/analyze', methods=['POST'])
def analyze():
    """Run the analysis on uploaded assessment and evidence"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    
    # Load controls
    try:
        with open(session_path / 'controls.json', 'r') as f:
            controls = json.load(f)
    except Exception as e:
        return jsonify({'error': f'Failed to load controls: {str(e)}'}), 500
    
    # Extract text from all evidence files
    evidence_texts = {}
    evidence_path = session_path / 'evidence'
    if evidence_path.exists():
        for evidence_file in evidence_path.iterdir():
            if evidence_file.is_file():
                logger.info(f"Processing evidence file: {evidence_file.name}")
                file_texts = reviewer.extract_evidence_text(str(evidence_file))
                evidence_texts[evidence_file.name] = file_texts
    
    # Get framework
    framework = session.get('framework', 'NIST')
    
    # Analyze controls
    results = []
    max_controls = min(len(controls), int(request.json.get('max_controls', 20)))
    
    for i, control in enumerate(controls[:max_controls]):
        logger.info(f"Analyzing control {i+1}/{max_controls}: {control.get('control', 'Unknown')}")
        
        try:
            result = reviewer.analyze_control_point(control, evidence_texts, framework)
            results.append(result)
        except Exception as e:
            logger.error(f"Error analyzing control {i+1}: {e}")
            continue
    
    # Calculate risk metrics
    risk_metrics = reviewer.calculate_risk_score(results)
    
    # Convert results to dict format
    results_dict = [
        {
            'control_id': r.control_id,
            'control_name': r.control_name,
            'requirement': r.requirement,
            'supplier_answer': r.supplier_answer,
            'evidence_validity': r.evidence_validity,
            'compliance_status': r.compliance_status,
            'risk_level': r.risk_level,
            'key_findings': r.key_findings,
            'risks': r.risks,
            'remediation': r.remediation,
            'evidence_references': r.evidence_references,
            'confidence_score': r.confidence_score
        }
        for r in results
    ]
    
    # Save results
    analysis_results = {
        'results': results_dict,
        'risk_metrics': risk_metrics,
        'framework': framework,
        'timestamp': datetime.now().isoformat(),
        'total_controls': len(controls),
        'analyzed_controls': len(results)
    }
    
    with open(session_path / 'results.json', 'w') as f:
        json.dump(analysis_results, f, indent=2)
    
    return jsonify({
        'success': True,
        'results': results_dict,
        'risk_metrics': risk_metrics,
        'summary': {
            'total_controls': len(controls),
            'analyzed_controls': len(results),
            'framework': framework
        }
    })

@app.route('/system_status')
def system_status():
    """Get system and model status"""
    status = {
        'backend': 'Ollama' if reviewer.use_ollama else 'Transformers',
        'model': reviewer.model_name,
        'frameworks': list(reviewer.frameworks.keys()),
        'ollama_available': reviewer.check_ollama_running() if hasattr(reviewer, 'check_ollama_running') else False
    }
    
    if reviewer.use_ollama:
        try:
            # Get Ollama models
            response = requests.get("http://localhost:11434/api/tags")
            if response.status_code == 200:
                models = [m['name'] for m in response.json().get('models', [])]
                status['available_models'] = models
        except:
            status['available_models'] = []
    
    return jsonify(status)

@app.route('/download_report/<session_id>')
def download_report(session_id):
    """Download the analysis report"""
    session_path = Path(f"sessions/{session_id}")
    results_path = session_path / 'results.json'
    
    if not results_path.exists():
        return jsonify({'error': 'Report not found'}), 404
    
    with open(results_path, 'r') as f:
        data = json.load(f)
    
    # Create Excel report with multiple sheets
    report_path = session_path / 'cyber_assessment_report.xlsx'
    
    with pd.ExcelWriter(report_path, engine='openpyxl') as writer:
        # Executive Summary sheet
        summary_data = {
            'Metric': [
                'Assessment Date',
                'Framework Used',
                'Total Controls',
                'Controls Analyzed',
                'Overall Risk Score',
                'Compliance Percentage',
                'Critical Risk Controls',
                'High Risk Controls'
            ],
            'Value': [
                data['timestamp'][:10],
                data['framework'],
                data['total_controls'],
                data['analyzed_controls'],
                f"{data['risk_metrics']['overall_risk_score']}%",
                f"{data['risk_metrics']['compliance_percentage']}%",
                data['risk_metrics']['risk_distribution'].get('Critical', 0),
                data['risk_metrics']['risk_distribution'].get('High', 0)
            ]
        }
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Executive Summary', index=False)
        
        # Risk Analysis sheet
        risk_df = pd.DataFrame([data['risk_metrics']['risk_distribution']])
        risk_df.to_excel(writer, sheet_name='Risk Analysis', index=False)
        
        # Detailed Results sheet
        results_df = pd.DataFrame(data['results'])
        
        # Convert lists to strings for Excel
        for col in ['key_findings', 'risks', 'remediation', 'evidence_references']:
            if col in results_df.columns:
                results_df[col] = results_df[col].apply(lambda x: '\n'.join(x) if isinstance(x, list) else x)
        
        results_df.to_excel(writer, sheet_name='Detailed Analysis', index=False)
        
        # High Priority Controls sheet
        if data['risk_metrics']['high_priority_controls']:
            priority_df = pd.DataFrame(data['risk_metrics']['high_priority_controls'])
            priority_df.to_excel(writer, sheet_name='High Priority Controls', index=False)
        
        # Format the Excel file
        workbook = writer.book
        for sheet_name in writer.sheets:
            worksheet = writer.sheets[sheet_name]
            
            # Auto-adjust column widths
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
    
    return send_file(str(report_path), as_attachment=True, 
                    download_name=f'cyber_assessment_report_{datetime.now().strftime("%Y%m%d")}.xlsx')

# Create enhanced HTML template
def get_html_template(backend="Ollama"):
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cyber Assessment Review System - Powered by Mistral 7B</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: #f5f5f5;
            color: #333;
            line-height: 1.6;
        }}
        
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            padding: 20px;
        }}
        
        header {{
            background: linear-gradient(135deg, #1a73e8 0%, #4285f4 100%);
            color: white;
            padding: 30px 0;
            text-align: center;
            margin-bottom: 40px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        }}
        
        h1 {{
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        
        .subtitle {{
            font-size: 1.2em;
            opacity: 0.9;
        }}
        
        .model-info {{
            font-size: 0.9em;
            opacity: 0.8;
            margin-top: 5px;
        }}
        
        .upload-section {{
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            margin-bottom: 30px;
        }}
        
        .upload-section h2 {{
            color: #1a73e8;
            margin-bottom: 20px;
            display: flex;
            align-items: center;
            gap: 10px;
        }}
        
        .step-number {{
            background: #1a73e8;
            color: white;
            width: 30px;
            height: 30px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
            font-size: 14px;
        }}
        
        .framework-selector {{
            margin-bottom: 20px;
        }}
        
        .framework-selector label {{
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }}
        
        .framework-selector select {{
            width: 100%;
            padding: 10px;
            border: 2px solid #e0e0e0;
            border-radius: 4px;
            font-size: 16px;
            background: white;
        }}
        
        .file-input-wrapper {{
            position: relative;
            margin-bottom: 20px;
        }}
        
        .file-input {{
            position: absolute;
            opacity: 0;
            width: 100%;
            height: 100%;
            cursor: pointer;
        }}
        
        .file-input-label {{
            display: block;
            padding: 20px;
            background: #f8f9fa;
            border: 2px dashed #1a73e8;
            border-radius: 8px;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s;
        }}
        
        .file-input-label:hover {{
            background: #e8f0fe;
            border-color: #1557b0;
        }}
        
        .file-input-label.drag-over {{
            background: #e8f0fe;
            border-color: #1557b0;
            transform: scale(1.02);
        }}
        
        .file-icon {{
            font-size: 48px;
            color: #1a73e8;
            margin-bottom: 10px;
        }}
        
        .file-list {{
            margin-top: 20px;
            display: none;
        }}
        
        .file-item {{
            padding: 12px 16px;
            margin: 8px 0;
            background: #f8f9fa;
            border-radius: 6px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            border: 1px solid #e0e0e0;
        }}
        
        .file-item.success {{
            border-color: #4caf50;
            background: #f1f8f4;
        }}
        
        .file-info {{
            display: flex;
            align-items: center;
            gap: 10px;
        }}
        
        .file-size {{
            color: #666;
            font-size: 14px;
        }}
        
        .btn {{
            padding: 12px 30px;
            background: #1a73e8;
            color: white;
            border: none;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            transition: all 0.3s;
            display: inline-flex;
            align-items: center;
            gap: 8px;
        }}
        
        .btn:hover {{
            background: #1557b0;
            transform: translateY(-1px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.2);
        }}
        
        .btn:disabled {{
            background: #ccc;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }}
        
        .btn-secondary {{
            background: #6c757d;
        }}
        
        .btn-secondary:hover {{
            background: #5a6268;
        }}
        
        .analysis-options {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            margin-bottom: 20px;
        }}
        
        .option-group {{
            margin-bottom: 15px;
        }}
        
        .option-group label {{
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }}
        
        .option-group input {{
            width: 200px;
            padding: 8px;
            border: 1px solid #ddd;
            border-radius: 4px;
        }}
        
        .results-section {{
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            display: none;
        }}
        
        .risk-score-display {{
            text-align: center;
            margin: 30px 0;
        }}
        
        .risk-score {{
            font-size: 72px;
            font-weight: bold;
            margin: 10px 0;
        }}
        
        .risk-score.critical {{ color: #d32f2f; }}
        .risk-score.high {{ color: #f57c00; }}
        .risk-score.medium {{ color: #fbc02d; }}
        .risk-score.low {{ color: #388e3c; }}
        
        .summary-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }}
        
        .summary-card {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            text-align: center;
            border: 1px solid #e0e0e0;
            transition: transform 0.2s;
        }}
        
        .summary-card:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }}
        
        .summary-card h3 {{
            color: #666;
            font-size: 14px;
            margin-bottom: 10px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}
        
        .summary-card .value {{
            font-size: 36px;
            font-weight: bold;
            color: #1a73e8;
        }}
        
        .summary-card .sub-value {{
            font-size: 14px;
            color: #666;
            margin-top: 5px;
        }}
        
        .tabs {{
            display: flex;
            gap: 10px;
            margin-bottom: 20px;
            border-bottom: 2px solid #e0e0e0;
        }}
        
        .tab {{
            padding: 10px 20px;
            background: none;
            border: none;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            color: #666;
            border-bottom: 3px solid transparent;
            transition: all 0.3s;
        }}
        
        .tab:hover {{
            color: #1a73e8;
        }}
        
        .tab.active {{
            color: #1a73e8;
            border-bottom-color: #1a73e8;
        }}
        
        .tab-content {{
            display: none;
        }}
        
        .tab-content.active {{
            display: block;
        }}
        
        .control-result {{
            background: #f8f9fa;
            padding: 20px;
            margin: 15px 0;
            border-radius: 8px;
            border-left: 4px solid #1a73e8;
        }}
        
        .control-result.critical {{
            border-left-color: #d32f2f;
        }}
        
        .control-result.high {{
            border-left-color: #f57c00;
        }}
        
        .control-result.medium {{
            border-left-color: #fbc02d;
        }}
        
        .control-result.low {{
            border-left-color: #388e3c;
        }}
        
        .control-header {{
            display: flex;
            justify-content: space-between;
            align-items: start;
            margin-bottom: 15px;
        }}
        
        .control-title {{
            flex: 1;
        }}
        
        .control-title h4 {{
            color: #333;
            margin-bottom: 5px;
        }}
        
        .control-id {{
            color: #666;
            font-size: 14px;
        }}
        
        .badges {{
            display: flex;
            gap: 8px;
            flex-wrap: wrap;
        }}
        
        .badge {{
            display: inline-block;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: bold;
            text-transform: uppercase;
        }}
        
        .badge.valid {{ background: #4caf50; color: white; }}
        .badge.partially-valid {{ background: #ff9800; color: white; }}
        .badge.invalid {{ background: #f44336; color: white; }}
        .badge.no-evidence {{ background: #9e9e9e; color: white; }}
        
        .badge.compliant {{ background: #4caf50; color: white; }}
        .badge.partially-compliant {{ background: #ff9800; color: white; }}
        .badge.non-compliant {{ background: #f44336; color: white; }}
        
        .badge.critical {{ background: #d32f2f; color: white; }}
        .badge.high {{ background: #f57c00; color: white; }}
        .badge.medium {{ background: #fbc02d; color: #333; }}
        .badge.low {{ background: #388e3c; color: white; }}
        
        .control-details {{
            margin-top: 15px;
        }}
        
        .detail-section {{
            margin: 15px 0;
        }}
        
        .detail-section h5 {{
            color: #666;
            margin-bottom: 8px;
            font-size: 14px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}
        
        .detail-list {{
            list-style: none;
            padding-left: 0;
        }}
        
        .detail-list li {{
            margin: 5px 0;
            padding-left: 20px;
            position: relative;
            color: #333;
        }}
        
        .detail-list li:before {{
            content: "▸";
            position: absolute;
            left: 0;
            color: #1a73e8;
        }}
        
        .confidence-meter {{
            margin-top: 10px;
            background: #e0e0e0;
            height: 8px;
            border-radius: 4px;
            overflow: hidden;
        }}
        
        .confidence-fill {{
            height: 100%;
            background: #1a73e8;
            transition: width 0.3s;
        }}
        
        .loading {{
            display: none;
            text-align: center;
            padding: 60px;
        }}
        
        .spinner {{
            border: 4px solid #f3f3f3;
            border-top: 4px solid #1a73e8;
            border-radius: 50%;
            width: 60px;
            height: 60px;
            animation: spin 1s linear infinite;
            margin: 0 auto 20px;
        }}
        
        @keyframes spin {{
            0% {{ transform: rotate(0deg); }}
            100% {{ transform: rotate(360deg); }}
        }}
        
        .progress-bar {{
            width: 100%;
            height: 20px;
            background: #e0e0e0;
            border-radius: 10px;
            overflow: hidden;
            margin: 20px 0;
        }}
        
        .progress-fill {{
            height: 100%;
            background: #1a73e8;
            width: 0%;
            transition: width 0.3s;
            display: flex;
            align-items: center;
            justify-content: center;
            color: white;
            font-size: 12px;
            font-weight: bold;
        }}
        
        .error-message {{
            background: #fee;
            color: #c33;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #c33;
        }}
        
        .success-message {{
            background: #e8f5e9;
            color: #2e7d32;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #2e7d32;
        }}
        
        .charts-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 20px;
            margin: 30px 0;
        }}
        
        .chart-container {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border: 1px solid #e0e0e0;
        }}
        
        .chart-title {{
            font-weight: 600;
            margin-bottom: 15px;
            color: #333;
        }}
    </style>
</head>
<body>
    <header>
        <h1>🛡️ Cyber Assessment Review System</h1>
        <p class="subtitle">AI-Powered Security Control Validation</p>
        <p class="model-info">Powered by Mistral 7B Instruct - Running Locally via {backend}</p>
    </header>
    
    <div class="container">
        <!-- Upload Assessment Section -->
        <div class="upload-section">
            <h2><span class="step-number">1</span> Upload Assessment File</h2>
            
            <div class="framework-selector">
                <label for="framework">Select Compliance Framework:</label>
                <select id="framework">
                    <option value="NIST">NIST Cybersecurity Framework</option>
                    <option value="ISO27001">ISO/IEC 27001:2022</option>
                    <option value="SOC2">SOC 2 Type II</option>
                    <option value="CIS">CIS Controls v8</option>
                    <option value="PCI-DSS">PCI-DSS v4.0</option>
                </select>
            </div>
            
            <div class="file-input-wrapper">
                <input type="file" id="assessmentFile" class="file-input" accept=".xlsx">
                <label for="assessmentFile" class="file-input-label" id="assessmentLabel">
                    <div class="file-icon">📊</div>
                    <div>Click to select or drag & drop assessment file (.xlsx)</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">Excel file with control assessments</div>
                </label>
            </div>
            <div id="assessmentFileList" class="file-list"></div>
        </div>
        
        <!-- Upload Evidence Section -->
        <div class="upload-section">
            <h2><span class="step-number">2</span> Upload Evidence Files</h2>
            <div class="file-input-wrapper">
                <input type="file" id="evidenceFiles" class="file-input" multiple accept=".pdf,.ppt,.pptx,.xlsx,.docx">
                <label for="evidenceFiles" class="file-input-label" id="evidenceLabel">
                    <div class="file-icon">📁</div>
                    <div>Click to select or drag & drop evidence files</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">PDF, PPT, XLSX, DOCX - Multiple files allowed</div>
                </label>
            </div>
            <div id="evidenceFileList" class="file-list"></div>
        </div>
        
        <!-- Analysis Options -->
        <div class="upload-section">
            <h2><span class="step-number">3</span> Analysis Options</h2>
            <div class="analysis-options">
                <div class="option-group">
                    <label for="maxControls">Maximum controls to analyze:</label>
                    <input type="number" id="maxControls" value="20" min="1" max="100">
                </div>
            </div>
            <div style="text-align: center;">
                <button id="analyzeBtn" class="btn" disabled>
                    <span>🔍</span> Analyze Assessment
                </button>
                <button id="statusBtn" class="btn btn-secondary" onclick="checkStatus()">
                    <span>ℹ️</span> System Status
                </button>
            </div>
        </div>
        
        <!-- Loading Section -->
        <div id="loadingSection" class="loading">
            <div class="spinner"></div>
            <h3>Analyzing Assessment with Mistral 7B ({backend})</h3>
            <p style="color: #666; margin-top: 10px;">This may take a few minutes depending on the number of controls...</p>
            <div class="progress-bar">
                <div class="progress-fill" id="progressFill">0%</div>
            </div>
            <p id="progressText" style="color: #666; margin-top: 10px;">Initializing analysis...</p>
        </div>
        
        <!-- Error/Success Messages -->
        <div id="errorMessage" class="error-message"></div>
        <div id="successMessage" class="success-message"></div>
        
        <!-- Results Section -->
        <div id="resultsSection" class="results-section">
            <h2>📊 Analysis Results</h2>
            
            <div class="risk-score-display">
                <h3>Overall Risk Score</h3>
                <div id="riskScore" class="risk-score">-</div>
                <p style="color: #666;">Lower is better (0-100 scale)</p>
            </div>
            
            <div id="summaryGrid" class="summary-grid"></div>
            
            <div class="charts-grid" id="chartsGrid"></div>
            
            <div class="tabs">
                <button class="tab active" onclick="showTab('detailed')">Detailed Analysis</button>
                <button class="tab" onclick="showTab('highrisk')">High Risk Controls</button>
                <button class="tab" onclick="showTab('summary')">Executive Summary</button>
            </div>
            
            <div id="detailed" class="tab-content active">
                <h3>Detailed Control Analysis</h3>
                <div id="detailedResults"></div>
            </div>
            
            <div id="highrisk" class="tab-content">
                <h3>High Priority Controls Requiring Immediate Attention</h3>
                <div id="highRiskResults"></div>
            </div>
            
            <div id="summary" class="tab-content">
                <h3>Executive Summary</h3>
                <div id="executiveSummary"></div>
            </div>
            
            <div style="text-align: center; margin-top: 30px;">
                <button id="downloadReportBtn" class="btn">
                    <span>📥</span> Download Full Report
                </button>
            </div>
        </div>
    </div>
    
    <script>
        let sessionId = null;
        let assessmentUploaded = false;
        let evidenceUploaded = false;
        
        // Drag and drop functionality
        function setupDragAndDrop(inputId, labelId) {{
            const label = document.getElementById(labelId);
            const input = document.getElementById(inputId);
            
            ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {{
                label.addEventListener(eventName, preventDefaults, false);
            }});
            
            function preventDefaults(e) {{
                e.preventDefault();
                e.stopPropagation();
            }}
            
            ['dragenter', 'dragover'].forEach(eventName => {{
                label.addEventListener(eventName, highlight, false);
            }});
            
            ['dragleave', 'drop'].forEach(eventName => {{
                label.addEventListener(eventName, unhighlight, false);
            }});
            
            function highlight(e) {{
                label.classList.add('drag-over');
            }}
            
            function unhighlight(e) {{
                label.classList.remove('drag-over');
            }}
            
            label.addEventListener('drop', handleDrop, false);
            
            function handleDrop(e) {{
                const dt = e.dataTransfer;
                const files = dt.files;
                input.files = files;
                input.dispatchEvent(new Event('change', {{ bubbles: true }}));
            }}
        }}
        
        setupDragAndDrop('assessmentFile', 'assessmentLabel');
        setupDragAndDrop('evidenceFiles', 'evidenceLabel');
        
        // File upload handlers
        document.getElementById('assessmentFile').addEventListener('change', async (e) => {{
            const file = e.target.files[0];
            if (!file) return;
            
            const formData = new FormData();
            formData.append('assessment', file);
            formData.append('framework', document.getElementById('framework').value);
            
            try {{
                const response = await fetch('/upload_assessment', {{
                    method: 'POST',
                    body: formData
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    sessionId = data.session_id;
                    assessmentUploaded = true;
                    
                    document.getElementById('assessmentFileList').style.display = 'block';
                    document.getElementById('assessmentFileList').innerHTML = `
                        <div class="file-item success">
                            <div class="file-info">
                                <span>📊 ${{file.name}}</span>
                                <span class="file-size">${{formatFileSize(file.size)}}</span>
                            </div>
                            <span style="color: #4caf50;">✓ ${{data.controls_count}} controls found</span>
                        </div>
                    `;
                    
                    if (data.sample_control) {{
                        console.log('Sample control:', data.sample_control);
                    }}
                    
                    checkAnalyzeButton();
                    showSuccess(`Assessment uploaded successfully! Found ${{data.controls_count}} controls.`);
                }} else {{
                    showError(data.error || 'Failed to upload assessment');
                }}
            }} catch (error) {{
                showError('Error uploading assessment: ' + error.message);
            }}
        }});
        
        document.getElementById('evidenceFiles').addEventListener('change', async (e) => {{
            const files = e.target.files;
            if (!files.length) return;
            
            if (!sessionId) {{
                showError('Please upload assessment file first');
                return;
            }}
            
            const formData = new FormData();
            for (let i = 0; i < files.length; i++) {{
                formData.append(`evidence_${{i}}`, files[i]);
            }}
            
            try {{
                const response = await fetch('/upload_evidence', {{
                    method: 'POST',
                    body: formData
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    evidenceUploaded = true;
                    
                    const fileListHtml = data.uploaded_files.map(file => {{
                        const icon = getFileIcon(file.filename);
                        return `
                            <div class="file-item success">
                                <div class="file-info">
                                    <span>${{icon}} ${{file.filename}}</span>
                                    <span class="file-size">${{formatFileSize(file.size)}}</span>
                                </div>
                                <span style="color: #4caf50;">✓</span>
                            </div>
                        `;
                    }}).join('');
                    
                    document.getElementById('evidenceFileList').style.display = 'block';
                    document.getElementById('evidenceFileList').innerHTML = fileListHtml;
                    
                    checkAnalyzeButton();
                    showSuccess(`${{data.total_files}} evidence files uploaded successfully!`);
                }} else {{
                    showError(data.error || 'Failed to upload evidence files');
                }}
            }} catch (error) {{
                showError('Error uploading evidence: ' + error.message);
            }}
        }});
        
        // Analyze button handler
        document.getElementById('analyzeBtn').addEventListener('click', async () => {{
            document.getElementById('loadingSection').style.display = 'block';
            document.getElementById('resultsSection').style.display = 'none';
            hideMessages();
            
            const maxControls = document.getElementById('maxControls').value;
            
            try {{
                const response = await fetch('/analyze', {{
                    method: 'POST',
                    headers: {{
                        'Content-Type': 'application/json'
                    }},
                    body: JSON.stringify({{ max_controls: maxControls }})
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    displayResults(data);
                    showSuccess('Analysis completed successfully!');
                }} else {{
                    showError(data.error || 'Analysis failed');
                }}
            }} catch (error) {{
                showError('Error during analysis: ' + error.message);
            }} finally {{
                document.getElementById('loadingSection').style.display = 'none';
            }}
        }});
        
        // Download report handler
        document.getElementById('downloadReportBtn').addEventListener('click', () => {{
            if (sessionId) {{
                window.location.href = `/download_report/${{sessionId}}`;
            }}
        }});
        
        // Helper functions
        function checkAnalyzeButton() {{
            const btn = document.getElementById('analyzeBtn');
            btn.disabled = !(assessmentUploaded && evidenceUploaded);
        }}
        
        function showError(message) {{
            const errorDiv = document.getElementById('errorMessage');
            errorDiv.textContent = '⚠️ ' + message;
            errorDiv.style.display = 'block';
            setTimeout(() => errorDiv.style.display = 'none', 7000);
        }}
        
        function showSuccess(message) {{
            const successDiv = document.getElementById('successMessage');
            successDiv.textContent = '✅ ' + message;
            successDiv.style.display = 'block';
            setTimeout(() => successDiv.style.display = 'none', 5000);
        }}
        
        function hideMessages() {{
            document.getElementById('errorMessage').style.display = 'none';
            document.getElementById('successMessage').style.display = 'none';
        }}
        
        function formatFileSize(bytes) {{
            if (bytes === 0) return '0 Bytes';
            const k = 1024;
            const sizes = ['Bytes', 'KB', 'MB', 'GB'];
            const i = Math.floor(Math.log(bytes) / Math.log(k));
            return parseFloat((bytes / Math.pow(k, i)).toFixed(2)) + ' ' + sizes[i];
        }}
        
        function getFileIcon(filename) {{
            const ext = filename.split('.').pop().toLowerCase();
            const icons = {{
                'pdf': '📄',
                'xlsx': '📊',
                'docx': '📝',
                'ppt': '📊',
                'pptx': '📊'
            }};
            return icons[ext] || '📎';
        }}
        
        function showTab(tabName) {{
            // Hide all tabs
            document.querySelectorAll('.tab-content').forEach(content => {{
                content.classList.remove('active');
            }});
            document.querySelectorAll('.tab').forEach(tab => {{
                tab.classList.remove('active');
            }});
            
            // Show selected tab
            document.getElementById(tabName).classList.add('active');
            event.target.classList.add('active');
        }}
        
        function getRiskScoreClass(score) {{
            if (score >= 70) return 'critical';
            if (score >= 50) return 'high';
            if (score >= 30) return 'medium';
            return 'low';
        }}
        
        function displayResults(data) {{
            // Display risk score
            const riskScore = data.risk_metrics.overall_risk_score;
            const riskScoreElement = document.getElementById('riskScore');
            riskScoreElement.textContent = riskScore + '%';
            riskScoreElement.className = 'risk-score ' + getRiskScoreClass(riskScore);
            
            // Display summary
            const summaryHtml = `
                <div class="summary-card">
                    <h3>Total Controls</h3>
                    <div class="value">${{data.summary.total_controls}}</div>
                    <div class="sub-value">In Assessment</div>
                </div>
                <div class="summary-card">
                    <h3>Analyzed</h3>
                    <div class="value">${{data.summary.analyzed_controls}}</div>
                    <div class="sub-value">By Mistral 7B</div>
                </div>
                <div class="summary-card">
                    <h3>Compliance Rate</h3>
                    <div class="value" style="color: #4caf50;">${{data.risk_metrics.compliance_percentage}}%</div>
                    <div class="sub-value">Controls Compliant</div>
                </div>
                <div class="summary-card">
                    <h3>Critical Risks</h3>
                    <div class="value" style="color: #d32f2f;">${{data.risk_metrics.risk_distribution.Critical || 0}}</div>
                    <div class="sub-value">Immediate Action</div>
                </div>
            `;
            document.getElementById('summaryGrid').innerHTML = summaryHtml;
            
            // Display charts placeholder
            const chartsHtml = `
                <div class="chart-container">
                    <div class="chart-title">Risk Distribution</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Critical: ${{data.risk_metrics.risk_distribution.Critical || 0}} | 
                        High: ${{data.risk_metrics.risk_distribution.High || 0}} | 
                        Medium: ${{data.risk_metrics.risk_distribution.Medium || 0}} | 
                        Low: ${{data.risk_metrics.risk_distribution.Low || 0}}
                    </div>
                </div>
                <div class="chart-container">
                    <div class="chart-title">Compliance Status</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Compliant: ${{data.risk_metrics.compliance_distribution.Compliant || 0}} | 
                        Partial: ${{data.risk_metrics.compliance_distribution['Partially Compliant'] || 0}} | 
                        Non-Compliant: ${{data.risk_metrics.compliance_distribution['Non-Compliant'] || 0}}
                    </div>
                </div>
            `;
            document.getElementById('chartsGrid').innerHTML = chartsHtml;
            
            // Display detailed results
            const resultsHtml = data.results.map(result => {{
                const riskClass = result.risk_level.toLowerCase().replace(' ', '-');
                const validityClass = result.evidence_validity.toLowerCase().replace(' ', '-');
                const complianceClass = result.compliance_status.toLowerCase().replace(' ', '-');
                
                return `
                    <div class="control-result ${{riskClass}}">
                        <div class="control-header">
                            <div class="control-title">
                                <h4>${{result.control_name || 'Unnamed Control'}}</h4>
                                <div class="control-id">ID: ${{result.control_id}}</div>
                            </div>
                            <div class="badges">
                                <span class="badge ${{validityClass}}">${{result.evidence_validity}}</span>
                                <span class="badge ${{complianceClass}}">${{result.compliance_status}}</span>
                                <span class="badge ${{riskClass}}">${{result.risk_level}}</span>
                            </div>
                        </div>
                        
                        <div class="control-details">
                            <div class="detail-section">
                                <h5>Requirement</h5>
                                <p>${{result.requirement}}</p>
                            </div>
                            
                            <div class="detail-section">
                                <h5>Supplier Response</h5>
                                <p>${{result.supplier_answer}}</p>
                            </div>
                            
                            ${{result.key_findings.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Key Findings</h5>
                                    <ul class="detail-list">
                                        ${{result.key_findings.map(f => `<li>${{f}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.risks.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Identified Risks</h5>
                                    <ul class="detail-list">
                                        ${{result.risks.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.remediation.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Remediation Recommendations</h5>
                                    <ul class="detail-list">
                                        ${{result.remediation.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.evidence_references.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Evidence Sources</h5>
                                    <ul class="detail-list">
                                        ${{result.evidence_references.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            <div class="detail-section">
                                <h5>Analysis Confidence</h5>
                                <div class="confidence-meter">
                                    <div class="confidence-fill" style="width: ${{result.confidence_score * 100}}%"></div>
                                </div>
                                <p style="font-size: 12px; color: #666; margin-top: 5px;">
                                    ${{(result.confidence_score * 100).toFixed(0)}}% confidence in this assessment
                                </p>
                            </div>
                        </div>
                    </div>
                `;
            }}).join('');
            
            document.getElementById('detailedResults').innerHTML = resultsHtml;
            
            // Display high risk controls
            const highRiskControls = data.results.filter(r => 
                r.risk_level === 'Critical' || r.risk_level === 'High'
            );
            
            if (highRiskControls.length > 0) {{
                const highRiskHtml = highRiskControls.map(result => {{
                    const riskClass = result.risk_level.toLowerCase();
                    return `
                        <div class="control-result ${{riskClass}}">
                            <h4>${{result.control_name}} (${{result.control_id}})</h4>
                            <p><strong>Risk Level:</strong> ${{result.risk_level}}</p>
                            <p><strong>Compliance Status:</strong> ${{result.compliance_status}}</p>
                            <p><strong>Primary Risk:</strong> ${{result.risks[0] || 'Not specified'}}</p>
                            <p><strong>Immediate Action:</strong> ${{result.remediation[0] || 'Review required'}}</p>
                        </div>
                    `;
                }}).join('');
                document.getElementById('highRiskResults').innerHTML = highRiskHtml;
            }} else {{
                document.getElementById('highRiskResults').innerHTML = 
                    '<p style="color: #4caf50; text-align: center; padding: 40px;">✅ No critical or high risk controls identified!</p>';
            }}
            
            // Display executive summary
            const executiveSummaryHtml = `
                <div style="padding: 20px;">
                    <h4>Assessment Overview</h4>
                    <p>This cyber security assessment was conducted using the <strong>${{data.summary.framework}}</strong> framework. 
                    A total of <strong>${{data.summary.analyzed_controls}}</strong> controls were analyzed out of 
                    <strong>${{data.summary.total_controls}}</strong> controls in the assessment.</p>
                    
                    <h4 style="margin-top: 20px;">Key Metrics</h4>
                    <ul>
                        <li>Overall Risk Score: <strong>${{data.risk_metrics.overall_risk_score}}%</strong> 
                            (${{getRiskScoreClass(data.risk_metrics.overall_risk_score)}} risk)</li>
                        <li>Compliance Rate: <strong>${{data.risk_metrics.compliance_percentage}}%</strong></li>
                        <li>Controls Requiring Immediate Attention: <strong>${{data.risk_metrics.high_priority_controls.length}}</strong></li>
                    </ul>
                    
                    <h4 style="margin-top: 20px;">Risk Summary</h4>
                    <p>The assessment identified ${{data.risk_metrics.risk_distribution.Critical || 0}} critical risks, 
                    ${{data.risk_metrics.risk_distribution.High || 0}} high risks, 
                    ${{data.risk_metrics.risk_distribution.Medium || 0}} medium risks, and 
                    ${{data.risk_metrics.risk_distribution.Low || 0}} low risks.</p>
                    
                    <h4 style="margin-top: 20px;">Recommendations</h4>
                    <p>Priority should be given to addressing the ${{data.risk_metrics.high_priority_controls.length}} 
                    high-priority controls that are both high-risk and non-compliant. These controls pose the greatest 
                    risk to the organization's security posture.</p>
                </div>
            `;
            document.getElementById('executiveSummary').innerHTML = executiveSummaryHtml;
            
            document.getElementById('resultsSection').style.display = 'block';
        }}
        
        // Check system status
        async function checkStatus() {{
            try {{
                const response = await fetch('/system_status');
                const status = await response.json();
                
                let message = `System Status:\\n\\n`;
                message += `Backend: ${{status.backend}}\\n`;
                message += `Model: ${{status.model}}\\n`;
                message += `Frameworks: ${{status.frameworks.join(', ')}}\\n`;
                
                if (status.backend === 'Ollama' && status.available_models) {{
                    message += `\\nAvailable Ollama Models:\\n`;
                    status.available_models.forEach(model => {{
                        message += `- ${{model}}\\n`;
                    }});
                }}
                
                alert(message);
            }} catch (error) {{
                alert('Error checking system status: ' + error.message);
            }}
        }}
        
        // Simulate progress during analysis
        let progressInterval;
        document.getElementById('analyzeBtn').addEventListener('click', () => {{
            let progress = 0;
            const progressFill = document.getElementById('progressFill');
            const progressText = document.getElementById('progressText');
            
            progressInterval = setInterval(() => {{
                progress += Math.random() * 15;
                if (progress > 90) progress = 90;
                
                progressFill.style.width = progress + '%';
                progressFill.textContent = Math.round(progress) + '%';
                
                if (progress < 30) {{
                    progressText.textContent = 'Loading Mistral 7B model...';
                }} else if (progress < 60) {{
                    progressText.textContent = 'Analyzing controls and evidence...';
                }} else {{
                    progressText.textContent = 'Generating risk assessments...';
                }}
            }}, 1000);
        }});
        
        // Clear progress interval when analysis completes
        const originalShowSuccess = showSuccess;
        showSuccess = function(message) {{
            if (progressInterval) {{
                clearInterval(progressInterval);
                progressInterval = null;
            }}
            originalShowSuccess(message);
        }};
    </script>
</body>
</html>
"""

# Save the template function
@app.route('/')
def index():
    """Render the main page"""
    backend = "Ollama" if reviewer.use_ollama else "Transformers"
    html_content = get_html_template(backend)
    return html_content
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: #f5f5f5;
            color: #333;
            line-height: 1.6;
        }}
        
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            padding: 20px;
        }}
        
        header {{
            background: linear-gradient(135deg, #1a73e8 0%, #4285f4 100%);
            color: white;
            padding: 30px 0;
            text-align: center;
            margin-bottom: 40px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        }}
        
        h1 {{
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        
        .subtitle {{
            font-size: 1.2em;
            opacity: 0.9;
        }}
        
        .model-info {{
            font-size: 0.9em;
            opacity: 0.8;
            margin-top: 5px;
        }}
        
        .upload-section {{
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            margin-bottom: 30px;
        }}
        
        .upload-section h2 {{
            color: #1a73e8;
            margin-bottom: 20px;
            display: flex;
            align-items: center;
            gap: 10px;
        }}
        
        .step-number {{
            background: #1a73e8;
            color: white;
            width: 30px;
            height: 30px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
            font-size: 14px;
        }}
        
        .framework-selector {{
            margin-bottom: 20px;
        }}
        
        .framework-selector label {{
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }}
        
        .framework-selector select {{
            width: 100%;
            padding: 10px;
            border: 2px solid #e0e0e0;
            border-radius: 4px;
            font-size: 16px;
            background: white;
        }}
        
        .file-input-wrapper {{
            position: relative;
            margin-bottom: 20px;
        }}
        
        .file-input {{
            position: absolute;
            opacity: 0;
            width: 100%;
            height: 100%;
            cursor: pointer;
        }}
        
        .file-input-label {{
            display: block;
            padding: 20px;
            background: #f8f9fa;
            border: 2px dashed #1a73e8;
            border-radius: 8px;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s;
        }}
        
        .file-input-label:hover {{
            background: #e8f0fe;
            border-color: #1557b0;
        }}
        
        .file-input-label.drag-over {{
            background: #e8f0fe;
            border-color: #1557b0;
            transform: scale(1.02);
        }}
        
        .file-icon {{
            font-size: 48px;
            color: #1a73e8;
            margin-bottom: 10px;
        }}
        
        .file-list {{
            margin-top: 20px;
            display: none;
        }}
        
        .file-item {{
            padding: 12px 16px;
            margin: 8px 0;
            background: #f8f9fa;
            border-radius: 6px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            border: 1px solid #e0e0e0;
        }}
        
        .file-item.success {{
            border-color: #4caf50;
            background: #f1f8f4;
        }}
        
        .file-info {{
            display: flex;
            align-items: center;
            gap: 10px;
        }}
        
        .file-size {{
            color: #666;
            font-size: 14px;
        }}
        
        .btn {{
            padding: 12px 30px;
            background: #1a73e8;
            color: white;
            border: none;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            transition: all 0.3s;
            display: inline-flex;
            align-items: center;
            gap: 8px;
        }}
        
        .btn:hover {{
            background: #1557b0;
            transform: translateY(-1px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.2);
        }}
        
        .btn:disabled {{
            background: #ccc;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }}
        
        .btn-secondary {{
            background: #6c757d;
        }}
        
        .btn-secondary:hover {{
            background: #5a6268;
        }}
        
        .analysis-options {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            margin-bottom: 20px;
        }}
        
        .option-group {{
            margin-bottom: 15px;
        }}
        
        .option-group label {{
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }}
        
        .option-group input {{
            width: 200px;
            padding: 8px;
            border: 1px solid #ddd;
            border-radius: 4px;
        }}
        
        .results-section {{
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            display: none;
        }}
        
        .risk-score-display {{
            text-align: center;
            margin: 30px 0;
        }}
        
        .risk-score {{
            font-size: 72px;
            font-weight: bold;
            margin: 10px 0;
        }}
        
        .risk-score.critical {{ color: #d32f2f; }}
        .risk-score.high {{ color: #f57c00; }}
        .risk-score.medium {{ color: #fbc02d; }}
        .risk-score.low {{ color: #388e3c; }}
        
        .summary-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }}
        
        .summary-card {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            text-align: center;
            border: 1px solid #e0e0e0;
            transition: transform 0.2s;
        }}
        
        .summary-card:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }}
        
        .summary-card h3 {{
            color: #666;
            font-size: 14px;
            margin-bottom: 10px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}
        
        .summary-card .value {{
            font-size: 36px;
            font-weight: bold;
            color: #1a73e8;
        }}
        
        .summary-card .sub-value {{
            font-size: 14px;
            color: #666;
            margin-top: 5px;
        }}
        
        .tabs {{
            display: flex;
            gap: 10px;
            margin-bottom: 20px;
            border-bottom: 2px solid #e0e0e0;
        }}
        
        .tab {{
            padding: 10px 20px;
            background: none;
            border: none;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            color: #666;
            border-bottom: 3px solid transparent;
            transition: all 0.3s;
        }}
        
        .tab:hover {{
            color: #1a73e8;
        }}
        
        .tab.active {{
            color: #1a73e8;
            border-bottom-color: #1a73e8;
        }}
        
        .tab-content {{
            display: none;
        }}
        
        .tab-content.active {{
            display: block;
        }}
        
        .control-result {{
            background: #f8f9fa;
            padding: 20px;
            margin: 15px 0;
            border-radius: 8px;
            border-left: 4px solid #1a73e8;
        }}
        
        .control-result.critical {{
            border-left-color: #d32f2f;
        }}
        
        .control-result.high {{
            border-left-color: #f57c00;
        }}
        
        .control-result.medium {{
            border-left-color: #fbc02d;
        }}
        
        .control-result.low {{
            border-left-color: #388e3c;
        }}
        
        .control-header {{
            display: flex;
            justify-content: space-between;
            align-items: start;
            margin-bottom: 15px;
        }}
        
        .control-title {{
            flex: 1;
        }}
        
        .control-title h4 {{
            color: #333;
            margin-bottom: 5px;
        }}
        
        .control-id {{
            color: #666;
            font-size: 14px;
        }}
        
        .badges {{
            display: flex;
            gap: 8px;
            flex-wrap: wrap;
        }}
        
        .badge {{
            display: inline-block;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: bold;
            text-transform: uppercase;
        }}
        
        .badge.valid {{ background: #4caf50; color: white; }}
        .badge.partially-valid {{ background: #ff9800; color: white; }}
        .badge.invalid {{ background: #f44336; color: white; }}
        .badge.no-evidence {{ background: #9e9e9e; color: white; }}
        
        .badge.compliant {{ background: #4caf50; color: white; }}
        .badge.partially-compliant {{ background: #ff9800; color: white; }}
        .badge.non-compliant {{ background: #f44336; color: white; }}
        
        .badge.critical {{ background: #d32f2f; color: white; }}
        .badge.high {{ background: #f57c00; color: white; }}
        .badge.medium {{ background: #fbc02d; color: #333; }}
        .badge.low {{ background: #388e3c; color: white; }}
        
        .control-details {{
            margin-top: 15px;
        }}
        
        .detail-section {{
            margin: 15px 0;
        }}
        
        .detail-section h5 {{
            color: #666;
            margin-bottom: 8px;
            font-size: 14px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}
        
        .detail-list {{
            list-style: none;
            padding-left: 0;
        }}
        
        .detail-list li {{
            margin: 5px 0;
            padding-left: 20px;
            position: relative;
            color: #333;
        }}
        
        .detail-list li:before {{
            content: "▸";
            position: absolute;
            left: 0;
            color: #1a73e8;
        }}
        
        .confidence-meter {{
            margin-top: 10px;
            background: #e0e0e0;
            height: 8px;
            border-radius: 4px;
            overflow: hidden;
        }}
        
        .confidence-fill {{
            height: 100%;
            background: #1a73e8;
            transition: width 0.3s;
        }}
        
        .loading {{
            display: none;
            text-align: center;
            padding: 60px;
        }}
        
        .spinner {{
            border: 4px solid #f3f3f3;
            border-top: 4px solid #1a73e8;
            border-radius: 50%;
            width: 60px;
            height: 60px;
            animation: spin 1s linear infinite;
            margin: 0 auto 20px;
        }}
        
        @keyframes spin {{
            0% {{ transform: rotate(0deg); }}
            100% {{ transform: rotate(360deg); }}
        }}
        
        .progress-bar {{
            width: 100%;
            height: 20px;
            background: #e0e0e0;
            border-radius: 10px;
            overflow: hidden;
            margin: 20px 0;
        }}
        
        .progress-fill {{
            height: 100%;
            background: #1a73e8;
            width: 0%;
            transition: width 0.3s;
            display: flex;
            align-items: center;
            justify-content: center;
            color: white;
            font-size: 12px;
            font-weight: bold;
        }}
        
        .error-message {{
            background: #fee;
            color: #c33;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #c33;
        }}
        
        .success-message {{
            background: #e8f5e9;
            color: #2e7d32;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #2e7d32;
        }}
        
        .charts-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 20px;
            margin: 30px 0;
        }}
        
        .chart-container {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border: 1px solid #e0e0e0;
        }}
        
        .chart-title {{
            font-weight: 600;
            margin-bottom: 15px;
            color: #333;
        }}
    </style>
</head>
<body>
    <header>
        <h1>🛡️ Cyber Assessment Review System</h1>
        <p class="subtitle">AI-Powered Security Control Validation</p>
        <p class="model-info">Powered by Mistral 7B Instruct - Running Locally via {backend}</p>
    </header>
    
    <div class="container">
        <!-- Upload Assessment Section -->
        <div class="upload-section">
            <h2><span class="step-number">1</span> Upload Assessment File</h2>
            
            <div class="framework-selector">
                <label for="framework">Select Compliance Framework:</label>
                <select id="framework">
                    <option value="NIST">NIST Cybersecurity Framework</option>
                    <option value="ISO27001">ISO/IEC 27001:2022</option>
                    <option value="SOC2">SOC 2 Type II</option>
                    <option value="CIS">CIS Controls v8</option>
                    <option value="PCI-DSS">PCI-DSS v4.0</option>
                </select>
            </div>
            
            <div class="file-input-wrapper">
                <input type="file" id="assessmentFile" class="file-input" accept=".xlsx">
                <label for="assessmentFile" class="file-input-label" id="assessmentLabel">
                    <div class="file-icon">📊</div>
                    <div>Click to select or drag & drop assessment file (.xlsx)</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">Excel file with control assessments</div>
                </label>
            </div>
            <div id="assessmentFileList" class="file-list"></div>
        </div>
        
        <!-- Upload Evidence Section -->
        <div class="upload-section">
            <h2><span class="step-number">2</span> Upload Evidence Files</h2>
            <div class="file-input-wrapper">
                <input type="file" id="evidenceFiles" class="file-input" multiple accept=".pdf,.ppt,.pptx,.xlsx,.docx">
                <label for="evidenceFiles" class="file-input-label" id="evidenceLabel">
                    <div class="file-icon">📁</div>
                    <div>Click to select or drag & drop evidence files</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">PDF, PPT, XLSX, DOCX - Multiple files allowed</div>
                </label>
            </div>
            <div id="evidenceFileList" class="file-list"></div>
        </div>
        
        <!-- Analysis Options -->
        <div class="upload-section">
            <h2><span class="step-number">3</span> Analysis Options</h2>
            <div class="analysis-options">
                <div class="option-group">
                    <label for="maxControls">Maximum controls to analyze:</label>
                    <input type="number" id="maxControls" value="20" min="1" max="100">
                </div>
            </div>
            <div style="text-align: center;">
                <button id="analyzeBtn" class="btn" disabled>
                    <span>🔍</span> Analyze Assessment
                </button>
                <button id="statusBtn" class="btn btn-secondary" onclick="checkStatus()">
                    <span>ℹ️</span> System Status
                </button>
            </div>
        </div>
        
        <!-- Loading Section -->
        <div id="loadingSection" class="loading">
            <div class="spinner"></div>
            <h3>Analyzing Assessment with Mistral 7B ({backend})</h3>
            <p style="color: #666; margin-top: 10px;">This may take a few minutes depending on the number of controls...</p>
            <div class="progress-bar">
                <div class="progress-fill" id="progressFill">0%</div>
            </div>
            <p id="progressText" style="color: #666; margin-top: 10px;">Initializing analysis...</p>
        </div>
        
        <!-- Error/Success Messages -->
        <div id="errorMessage" class="error-message"></div>
        <div id="successMessage" class="success-message"></div>
        
        <!-- Results Section -->
        <div id="resultsSection" class="results-section">
            <h2>📊 Analysis Results</h2>
            
            <div class="risk-score-display">
                <h3>Overall Risk Score</h3>
                <div id="riskScore" class="risk-score">-</div>
                <p style="color: #666;">Lower is better (0-100 scale)</p>
            </div>
            
            <div id="summaryGrid" class="summary-grid"></div>
            
            <div class="charts-grid" id="chartsGrid"></div>
            
            <div class="tabs">
                <button class="tab active" onclick="showTab('detailed')">Detailed Analysis</button>
                <button class="tab" onclick="showTab('highrisk')">High Risk Controls</button>
                <button class="tab" onclick="showTab('summary')">Executive Summary</button>
            </div>
            
            <div id="detailed" class="tab-content active">
                <h3>Detailed Control Analysis</h3>
                <div id="detailedResults"></div>
            </div>
            
            <div id="highrisk" class="tab-content">
                <h3>High Priority Controls Requiring Immediate Attention</h3>
                <div id="highRiskResults"></div>
            </div>
            
            <div id="summary" class="tab-content">
                <h3>Executive Summary</h3>
                <div id="executiveSummary"></div>
            </div>
            
            <div style="text-align: center; margin-top: 30px;">
                <button id="downloadReportBtn" class="btn">
                    <span>📥</span> Download Full Report
                </button>
            </div>
        </div>
    </div>
    
    <script>
        let sessionId = null;
        let assessmentUploaded = false;
        let evidenceUploaded = false;
        
        // Drag and drop functionality
        function setupDragAndDrop(inputId, labelId) {{
            const label = document.getElementById(labelId);
            const input = document.getElementById(inputId);
            
            ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {{
                label.addEventListener(eventName, preventDefaults, false);
            }});
            
            function preventDefaults(e) {{
                e.preventDefault();
                e.stopPropagation();
            }}
            
            ['dragenter', 'dragover'].forEach(eventName => {{
                label.addEventListener(eventName, highlight, false);
            }});
            
            ['dragleave', 'drop'].forEach(eventName => {{
                label.addEventListener(eventName, unhighlight, false);
            }});
            
            function highlight(e) {{
                label.classList.add('drag-over');
            }}
            
            function unhighlight(e) {{
                label.classList.remove('drag-over');
            }}
            
            label.addEventListener('drop', handleDrop, false);
            
            function handleDrop(e) {{
                const dt = e.dataTransfer;
                const files = dt.files;
                input.files = files;
                input.dispatchEvent(new Event('change', {{ bubbles: true }}));
            }}
        }}
        
        setupDragAndDrop('assessmentFile', 'assessmentLabel');
        setupDragAndDrop('evidenceFiles', 'evidenceLabel');
        
        // File upload handlers
        document.getElementById('assessmentFile').addEventListener('change', async (e) => {{
            const file = e.target.files[0];
            if (!file) return;
            
            const formData = new FormData();
            formData.append('assessment', file);
            formData.append('framework', document.getElementById('framework').value);
            
            try {{
                const response = await fetch('/upload_assessment', {{
                    method: 'POST',
                    body: formData
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    sessionId = data.session_id;
                    assessmentUploaded = true;
                    
                    document.getElementById('assessmentFileList').style.display = 'block';
                    document.getElementById('assessmentFileList').innerHTML = `
                        <div class="file-item success">
                            <div class="file-info">
                                <span>📊 ${{file.name}}</span>
                                <span class="file-size">${{formatFileSize(file.size)}}</span>
                            </div>
                            <span style="color: #4caf50;">✓ ${{data.controls_count}} controls found</span>
                        </div>
                    `;
                    
                    if (data.sample_control) {{
                        console.log('Sample control:', data.sample_control);
                    }}
                    
                    checkAnalyzeButton();
                    showSuccess(`Assessment uploaded successfully! Found ${{data.controls_count}} controls.`);
                }} else {{
                    showError(data.error || 'Failed to upload assessment');
                }}
            }} catch (error) {{
                showError('Error uploading assessment: ' + error.message);
            }}
        }});
        
        document.getElementById('evidenceFiles').addEventListener('change', async (e) => {{
            const files = e.target.files;
            if (!files.length) return;
            
            if (!sessionId) {{
                showError('Please upload assessment file first');
                return;
            }}
            
            const formData = new FormData();
            for (let i = 0; i < files.length; i++) {{
                formData.append(`evidence_${{i}}`, files[i]);
            }}
            
            try {{
                const response = await fetch('/upload_evidence', {{
                    method: 'POST',
                    body: formData
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    evidenceUploaded = true;
                    
                    const fileListHtml = data.uploaded_files.map(file => {{
                        const icon = getFileIcon(file.filename);
                        return `
                            <div class="file-item success">
                                <div class="file-info">
                                    <span>${{icon}} ${{file.filename}}</span>
                                    <span class="file-size">${{formatFileSize(file.size)}}</span>
                                </div>
                                <span style="color: #4caf50;">✓</span>
                            </div>
                        `;
                    }}).join('');
                    
                    document.getElementById('evidenceFileList').style.display = 'block';
                    document.getElementById('evidenceFileList').innerHTML = fileListHtml;
                    
                    checkAnalyzeButton();
                    showSuccess(`${{data.total_files}} evidence files uploaded successfully!`);
                }} else {{
                    showError(data.error || 'Failed to upload evidence files');
                }}
            }} catch (error) {{
                showError('Error uploading evidence: ' + error.message);
            }}
        }});
        
        // Analyze button handler
        document.getElementById('analyzeBtn').addEventListener('click', async () => {{
            document.getElementById('loadingSection').style.display = 'block';
            document.getElementById('resultsSection').style.display = 'none';
            hideMessages();
            
            const maxControls = document.getElementById('maxControls').value;
            
            try {{
                const response = await fetch('/analyze', {{
                    method: 'POST',
                    headers: {{
                        'Content-Type': 'application/json'
                    }},
                    body: JSON.stringify({{ max_controls: maxControls }})
                }});
                
                const data = await response.json();
                
                if (data.success) {{
                    displayResults(data);
                    showSuccess('Analysis completed successfully!');
                }} else {{
                    showError(data.error || 'Analysis failed');
                }}
            }} catch (error) {{
                showError('Error during analysis: ' + error.message);
            }} finally {{
                document.getElementById('loadingSection').style.display = 'none';
            }}
        }});
        
        // Download report handler
        document.getElementById('downloadReportBtn').addEventListener('click', () => {{
            if (sessionId) {{
                window.location.href = `/download_report/${{sessionId}}`;
            }}
        }});
        
        // Helper functions
        function checkAnalyzeButton() {{
            const btn = document.getElementById('analyzeBtn');
            btn.disabled = !(assessmentUploaded && evidenceUploaded);
        }}
        
        function showError(message) {{
            const errorDiv = document.getElementById('errorMessage');
            errorDiv.textContent = '⚠️ ' + message;
            errorDiv.style.display = 'block';
            setTimeout(() => errorDiv.style.display = 'none', 7000);
        }}
        
        function showSuccess(message) {{
            const successDiv = document.getElementById('successMessage');
            successDiv.textContent = '✅ ' + message;
            successDiv.style.display = 'block';
            setTimeout(() => successDiv.style.display = 'none', 5000);
        }}
        
        function hideMessages() {{
            document.getElementById('errorMessage').style.display = 'none';
            document.getElementById('successMessage').style.display = 'none';
        }}
        
        function formatFileSize(bytes) {{
            if (bytes === 0) return '0 Bytes';
            const k = 1024;
            const sizes = ['Bytes', 'KB', 'MB', 'GB'];
            const i = Math.floor(Math.log(bytes) / Math.log(k));
            return parseFloat((bytes / Math.pow(k, i)).toFixed(2)) + ' ' + sizes[i];
        }}
        
        function getFileIcon(filename) {{
            const ext = filename.split('.').pop().toLowerCase();
            const icons = {{
                'pdf': '📄',
                'xlsx': '📊',
                'docx': '📝',
                'ppt': '📊',
                'pptx': '📊'
            }};
            return icons[ext] || '📎';
        }}
        
        function showTab(tabName) {{
            // Hide all tabs
            document.querySelectorAll('.tab-content').forEach(content => {{
                content.classList.remove('active');
            }});
            document.querySelectorAll('.tab').forEach(tab => {{
                tab.classList.remove('active');
            }});
            
            // Show selected tab
            document.getElementById(tabName).classList.add('active');
            event.target.classList.add('active');
        }}
        
        function getRiskScoreClass(score) {{
            if (score >= 70) return 'critical';
            if (score >= 50) return 'high';
            if (score >= 30) return 'medium';
            return 'low';
        }}
        
        function displayResults(data) {{
            // Display risk score
            const riskScore = data.risk_metrics.overall_risk_score;
            const riskScoreElement = document.getElementById('riskScore');
            riskScoreElement.textContent = riskScore + '%';
            riskScoreElement.className = 'risk-score ' + getRiskScoreClass(riskScore);
            
            // Display summary
            const summaryHtml = `
                <div class="summary-card">
                    <h3>Total Controls</h3>
                    <div class="value">${{data.summary.total_controls}}</div>
                    <div class="sub-value">In Assessment</div>
                </div>
                <div class="summary-card">
                    <h3>Analyzed</h3>
                    <div class="value">${{data.summary.analyzed_controls}}</div>
                    <div class="sub-value">By Mistral 7B</div>
                </div>
                <div class="summary-card">
                    <h3>Compliance Rate</h3>
                    <div class="value" style="color: #4caf50;">${{data.risk_metrics.compliance_percentage}}%</div>
                    <div class="sub-value">Controls Compliant</div>
                </div>
                <div class="summary-card">
                    <h3>Critical Risks</h3>
                    <div class="value" style="color: #d32f2f;">${{data.risk_metrics.risk_distribution.Critical || 0}}</div>
                    <div class="sub-value">Immediate Action</div>
                </div>
            `;
            document.getElementById('summaryGrid').innerHTML = summaryHtml;
            
            // Display charts placeholder
            const chartsHtml = `
                <div class="chart-container">
                    <div class="chart-title">Risk Distribution</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Critical: ${{data.risk_metrics.risk_distribution.Critical || 0}} | 
                        High: ${{data.risk_metrics.risk_distribution.High || 0}} | 
                        Medium: ${{data.risk_metrics.risk_distribution.Medium || 0}} | 
                        Low: ${{data.risk_metrics.risk_distribution.Low || 0}}
                    </div>
                </div>
                <div class="chart-container">
                    <div class="chart-title">Compliance Status</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Compliant: ${{data.risk_metrics.compliance_distribution.Compliant || 0}} | 
                        Partial: ${{data.risk_metrics.compliance_distribution['Partially Compliant'] || 0}} | 
                        Non-Compliant: ${{data.risk_metrics.compliance_distribution['Non-Compliant'] || 0}}
                    </div>
                </div>
            `;
            document.getElementById('chartsGrid').innerHTML = chartsHtml;
            
            // Display detailed results
            const resultsHtml = data.results.map(result => {{
                const riskClass = result.risk_level.toLowerCase().replace(' ', '-');
                const validityClass = result.evidence_validity.toLowerCase().replace(' ', '-');
                const complianceClass = result.compliance_status.toLowerCase().replace(' ', '-');
                
                return `
                    <div class="control-result ${{riskClass}}">
                        <div class="control-header">
                            <div class="control-title">
                                <h4>${{result.control_name || 'Unnamed Control'}}</h4>
                                <div class="control-id">ID: ${{result.control_id}}</div>
                            </div>
                            <div class="badges">
                                <span class="badge ${{validityClass}}">${{result.evidence_validity}}</span>
                                <span class="badge ${{complianceClass}}">${{result.compliance_status}}</span>
                                <span class="badge ${{riskClass}}">${{result.risk_level}}</span>
                            </div>
                        </div>
                        
                        <div class="control-details">
                            <div class="detail-section">
                                <h5>Requirement</h5>
                                <p>${{result.requirement}}</p>
                            </div>
                            
                            <div class="detail-section">
                                <h5>Supplier Response</h5>
                                <p>${{result.supplier_answer}}</p>
                            </div>
                            
                            ${{result.key_findings.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Key Findings</h5>
                                    <ul class="detail-list">
                                        ${{result.key_findings.map(f => `<li>${{f}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.risks.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Identified Risks</h5>
                                    <ul class="detail-list">
                                        ${{result.risks.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.remediation.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Remediation Recommendations</h5>
                                    <ul class="detail-list">
                                        ${{result.remediation.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            ${{result.evidence_references.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Evidence Sources</h5>
                                    <ul class="detail-list">
                                        ${{result.evidence_references.map(r => `<li>${{r}}</li>`).join('')}}
                                    </ul>
                                </div>
                            ` : ''}}
                            
                            <div class="detail-section">
                                <h5>Analysis Confidence</h5>
                                <div class="confidence-meter">
                                    <div class="confidence-fill" style="width: ${{result.confidence_score * 100}}%"></div>
                                </div>
                                <p style="font-size: 12px; color: #666; margin-top: 5px;">
                                    ${{(result.confidence_score * 100).toFixed(0)}}% confidence in this assessment
                                </p>
                            </div>
                        </div>
                    </div>
                `;
            }}).join('');
            
            document.getElementById('detailedResults').innerHTML = resultsHtml;
            
            // Display high risk controls
            const highRiskControls = data.results.filter(r => 
                r.risk_level === 'Critical' || r.risk_level === 'High'
            );
            
            if (highRiskControls.length > 0) {{
                const highRiskHtml = highRiskControls.map(result => {{
                    const riskClass = result.risk_level.toLowerCase();
                    return `
                        <div class="control-result ${{riskClass}}">
                            <h4>${{result.control_name}} (${{result.control_id}})</h4>
                            <p><strong>Risk Level:</strong> ${{result.risk_level}}</p>
                            <p><strong>Compliance Status:</strong> ${{result.compliance_status}}</p>
                            <p><strong>Primary Risk:</strong> ${{result.risks[0] || 'Not specified'}}</p>
                            <p><strong>Immediate Action:</strong> ${{result.remediation[0] || 'Review required'}}</p>
                        </div>
                    `;
                }}).join('');
                document.getElementById('highRiskResults').innerHTML = highRiskHtml;
            }} else {{
                document.getElementById('highRiskResults').innerHTML = 
                    '<p style="color: #4caf50; text-align: center; padding: 40px;">✅ No critical or high risk controls identified!</p>';
            }}
            
            // Display executive summary
            const executiveSummaryHtml = `
                <div style="padding: 20px;">
                    <h4>Assessment Overview</h4>
                    <p>This cyber security assessment was conducted using the <strong>${{data.summary.framework}}</strong> framework. 
                    A total of <strong>${{data.summary.analyzed_controls}}</strong> controls were analyzed out of 
                    <strong>${{data.summary.total_controls}}</strong> controls in the assessment.</p>
                    
                    <h4 style="margin-top: 20px;">Key Metrics</h4>
                    <ul>
                        <li>Overall Risk Score: <strong>${{data.risk_metrics.overall_risk_score}}%</strong> 
                            (${{getRiskScoreClass(data.risk_metrics.overall_risk_score)}} risk)</li>
                        <li>Compliance Rate: <strong>${{data.risk_metrics.compliance_percentage}}%</strong></li>
                        <li>Controls Requiring Immediate Attention: <strong>${{data.risk_metrics.high_priority_controls.length}}</strong></li>
                    </ul>
                    
                    <h4 style="margin-top: 20px;">Risk Summary</h4>
                    <p>The assessment identified ${{data.risk_metrics.risk_distribution.Critical || 0}} critical risks, 
                    ${{data.risk_metrics.risk_distribution.High || 0}} high risks, 
                    ${{data.risk_metrics.risk_distribution.Medium || 0}} medium risks, and 
                    ${{data.risk_metrics.risk_distribution.Low || 0}} low risks.</p>
                    
                    <h4 style="margin-top: 20px;">Recommendations</h4>
                    <p>Priority should be given to addressing the ${{data.risk_metrics.high_priority_controls.length}} 
                    high-priority controls that are both high-risk and non-compliant. These controls pose the greatest 
                    risk to the organization's security posture.</p>
                </div>
            `;
            document.getElementById('executiveSummary').innerHTML = executiveSummaryHtml;
            
            document.getElementById('resultsSection').style.display = 'block';
        }}
        
        // Check system status
        async function checkStatus() {{
            try {{
                const response = await fetch('/system_status');
                const status = await response.json();
                
                let message = `System Status:\\n\\n`;
                message += `Backend: ${{status.backend}}\\n`;
                message += `Model: ${{status.model}}\\n`;
                message += `Frameworks: ${{status.frameworks.join(', ')}}\\n`;
                
                if (status.backend === 'Ollama' && status.available_models) {{
                    message += `\\nAvailable Ollama Models:\\n`;
                    status.available_models.forEach(model => {{
                        message += `- ${{model}}\\n`;
                    }});
                }}
                
                alert(message);
            }} catch (error) {{
                alert('Error checking system status: ' + error.message);
            }}
        }}
        
        // Simulate progress during analysis
        let progressInterval;
        document.getElementById('analyzeBtn').addEventListener('click', () => {{
            let progress = 0;
            const progressFill = document.getElementById('progressFill');
            const progressText = document.getElementById('progressText');
            
            progressInterval = setInterval(() => {{
                progress += Math.random() * 15;
                if (progress > 90) progress = 90;
                
                progressFill.style.width = progress + '%';
                progressFill.textContent = Math.round(progress) + '%';
                
                if (progress < 30) {{
                    progressText.textContent = 'Loading Mistral 7B model...';
                }} else if (progress < 60) {{
                    progressText.textContent = 'Analyzing controls and evidence...';
                }} else {{
                    progressText.textContent = 'Generating risk assessments...';
                }}
            }}, 1000);
        }});
        
        // Clear progress interval when analysis completes
        const originalShowSuccess = showSuccess;
        showSuccess = function(message) {{
            if (progressInterval) {{
                clearInterval(progressInterval);
                progressInterval = null;
            }}
            originalShowSuccess(message);
        }};
    </script>
</body>
</html>
"""

# Save the template function
@app.route('/')
def index():
    """Render the main page"""
    backend = "Ollama" if reviewer.use_ollama else "Transformers"
    html_content = get_html_template(backend)
    return html_content
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cyber Assessment Review System - Powered by Mistral 7B</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: #f5f5f5;
            color: #333;
            line-height: 1.6;
        }
        
        .container {
            max-width: 1400px;
            margin: 0 auto;
            padding: 20px;
        }
        
        header {
            background: linear-gradient(135deg, #1a73e8 0%, #4285f4 100%);
            color: white;
            padding: 30px 0;
            text-align: center;
            margin-bottom: 40px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        }
        
        h1 {
            font-size: 2.5em;
            margin-bottom: 10px;
        }
        
        .subtitle {
            font-size: 1.2em;
            opacity: 0.9;
        }
        
        .model-info {
            font-size: 0.9em;
            opacity: 0.8;
            margin-top: 5px;
        }
        
        .upload-section {
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            margin-bottom: 30px;
        }
        
        .upload-section h2 {
            color: #1a73e8;
            margin-bottom: 20px;
            display: flex;
            align-items: center;
            gap: 10px;
        }
        
        .step-number {
            background: #1a73e8;
            color: white;
            width: 30px;
            height: 30px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
            font-size: 14px;
        }
        
        .framework-selector {
            margin-bottom: 20px;
        }
        
        .framework-selector label {
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }
        
        .framework-selector select {
            width: 100%;
            padding: 10px;
            border: 2px solid #e0e0e0;
            border-radius: 4px;
            font-size: 16px;
            background: white;
        }
        
        .file-input-wrapper {
            position: relative;
            margin-bottom: 20px;
        }
        
        .file-input {
            position: absolute;
            opacity: 0;
            width: 100%;
            height: 100%;
            cursor: pointer;
        }
        
        .file-input-label {
            display: block;
            padding: 20px;
            background: #f8f9fa;
            border: 2px dashed #1a73e8;
            border-radius: 8px;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s;
        }
        
        .file-input-label:hover {
            background: #e8f0fe;
            border-color: #1557b0;
        }
        
        .file-input-label.drag-over {
            background: #e8f0fe;
            border-color: #1557b0;
            transform: scale(1.02);
        }
        
        .file-icon {
            font-size: 48px;
            color: #1a73e8;
            margin-bottom: 10px;
        }
        
        .file-list {
            margin-top: 20px;
            display: none;
        }
        
        .file-item {
            padding: 12px 16px;
            margin: 8px 0;
            background: #f8f9fa;
            border-radius: 6px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            border: 1px solid #e0e0e0;
        }
        
        .file-item.success {
            border-color: #4caf50;
            background: #f1f8f4;
        }
        
        .file-info {
            display: flex;
            align-items: center;
            gap: 10px;
        }
        
        .file-size {
            color: #666;
            font-size: 14px;
        }
        
        .btn {
            padding: 12px 30px;
            background: #1a73e8;
            color: white;
            border: none;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            transition: all 0.3s;
            display: inline-flex;
            align-items: center;
            gap: 8px;
        }
        
        .btn:hover {
            background: #1557b0;
            transform: translateY(-1px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.2);
        }
        
        .btn:disabled {
            background: #ccc;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }
        
        .btn-secondary {
            background: #6c757d;
        }
        
        .btn-secondary:hover {
            background: #5a6268;
        }
        
        .analysis-options {
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            margin-bottom: 20px;
        }
        
        .option-group {
            margin-bottom: 15px;
        }
        
        .option-group label {
            display: block;
            margin-bottom: 5px;
            font-weight: 600;
            color: #666;
        }
        
        .option-group input {
            width: 200px;
            padding: 8px;
            border: 1px solid #ddd;
            border-radius: 4px;
        }
        
        .results-section {
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            display: none;
        }
        
        .risk-score-display {
            text-align: center;
            margin: 30px 0;
        }
        
        .risk-score {
            font-size: 72px;
            font-weight: bold;
            margin: 10px 0;
        }
        
        .risk-score.critical { color: #d32f2f; }
        .risk-score.high { color: #f57c00; }
        .risk-score.medium { color: #fbc02d; }
        .risk-score.low { color: #388e3c; }
        
        .summary-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }
        
        .summary-card {
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            text-align: center;
            border: 1px solid #e0e0e0;
            transition: transform 0.2s;
        }
        
        .summary-card:hover {
            transform: translateY(-2px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }
        
        .summary-card h3 {
            color: #666;
            font-size: 14px;
            margin-bottom: 10px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .summary-card .value {
            font-size: 36px;
            font-weight: bold;
            color: #1a73e8;
        }
        
        .summary-card .sub-value {
            font-size: 14px;
            color: #666;
            margin-top: 5px;
        }
        
        .tabs {
            display: flex;
            gap: 10px;
            margin-bottom: 20px;
            border-bottom: 2px solid #e0e0e0;
        }
        
        .tab {
            padding: 10px 20px;
            background: none;
            border: none;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            color: #666;
            border-bottom: 3px solid transparent;
            transition: all 0.3s;
        }
        
        .tab:hover {
            color: #1a73e8;
        }
        
        .tab.active {
            color: #1a73e8;
            border-bottom-color: #1a73e8;
        }
        
        .tab-content {
            display: none;
        }
        
        .tab-content.active {
            display: block;
        }
        
        .control-result {
            background: #f8f9fa;
            padding: 20px;
            margin: 15px 0;
            border-radius: 8px;
            border-left: 4px solid #1a73e8;
        }
        
        .control-result.critical {
            border-left-color: #d32f2f;
        }
        
        .control-result.high {
            border-left-color: #f57c00;
        }
        
        .control-result.medium {
            border-left-color: #fbc02d;
        }
        
        .control-result.low {
            border-left-color: #388e3c;
        }
        
        .control-header {
            display: flex;
            justify-content: space-between;
            align-items: start;
            margin-bottom: 15px;
        }
        
        .control-title {
            flex: 1;
        }
        
        .control-title h4 {
            color: #333;
            margin-bottom: 5px;
        }
        
        .control-id {
            color: #666;
            font-size: 14px;
        }
        
        .badges {
            display: flex;
            gap: 8px;
            flex-wrap: wrap;
        }
        
        .badge {
            display: inline-block;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: bold;
            text-transform: uppercase;
        }
        
        .badge.valid { background: #4caf50; color: white; }
        .badge.partially-valid { background: #ff9800; color: white; }
        .badge.invalid { background: #f44336; color: white; }
        .badge.no-evidence { background: #9e9e9e; color: white; }
        
        .badge.compliant { background: #4caf50; color: white; }
        .badge.partially-compliant { background: #ff9800; color: white; }
        .badge.non-compliant { background: #f44336; color: white; }
        
        .badge.critical { background: #d32f2f; color: white; }
        .badge.high { background: #f57c00; color: white; }
        .badge.medium { background: #fbc02d; color: #333; }
        .badge.low { background: #388e3c; color: white; }
        
        .control-details {
            margin-top: 15px;
        }
        
        .detail-section {
            margin: 15px 0;
        }
        
        .detail-section h5 {
            color: #666;
            margin-bottom: 8px;
            font-size: 14px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .detail-list {
            list-style: none;
            padding-left: 0;
        }
        
        .detail-list li {
            margin: 5px 0;
            padding-left: 20px;
            position: relative;
            color: #333;
        }
        
        .detail-list li:before {
            content: "▸";
            position: absolute;
            left: 0;
            color: #1a73e8;
        }
        
        .confidence-meter {
            margin-top: 10px;
            background: #e0e0e0;
            height: 8px;
            border-radius: 4px;
            overflow: hidden;
        }
        
        .confidence-fill {
            height: 100%;
            background: #1a73e8;
            transition: width 0.3s;
        }
        
        .loading {
            display: none;
            text-align: center;
            padding: 60px;
        }
        
        .spinner {
            border: 4px solid #f3f3f3;
            border-top: 4px solid #1a73e8;
            border-radius: 50%;
            width: 60px;
            height: 60px;
            animation: spin 1s linear infinite;
            margin: 0 auto 20px;
        }
        
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        
        .progress-bar {
            width: 100%;
            height: 20px;
            background: #e0e0e0;
            border-radius: 10px;
            overflow: hidden;
            margin: 20px 0;
        }
        
        .progress-fill {
            height: 100%;
            background: #1a73e8;
            width: 0%;
            transition: width 0.3s;
            display: flex;
            align-items: center;
            justify-content: center;
            color: white;
            font-size: 12px;
            font-weight: bold;
        }
        
        .error-message {
            background: #fee;
            color: #c33;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #c33;
        }
        
        .success-message {
            background: #e8f5e9;
            color: #2e7d32;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            display: none;
            border-left: 4px solid #2e7d32;
        }
        
        .charts-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 20px;
            margin: 30px 0;
        }
        
        .chart-container {
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border: 1px solid #e0e0e0;
        }
        
        .chart-title {
            font-weight: 600;
            margin-bottom: 15px;
            color: #333;
        }
    </style>
</head>
<body>
    <header>
        <h1>🛡️ Cyber Assessment Review System</h1>
        <p class="subtitle">AI-Powered Security Control Validation</p>
        <p class="model-info">Powered by Mistral 7B Instruct - Running Locally via {{ backend }}</p>
    </header>
    
    <div class="container">
        <!-- Upload Assessment Section -->
        <div class="upload-section">
            <h2><span class="step-number">1</span> Upload Assessment File</h2>
            
            <div class="framework-selector">
                <label for="framework">Select Compliance Framework:</label>
                <select id="framework">
                    <option value="NIST">NIST Cybersecurity Framework</option>
                    <option value="ISO27001">ISO/IEC 27001:2022</option>
                    <option value="SOC2">SOC 2 Type II</option>
                    <option value="CIS">CIS Controls v8</option>
                    <option value="PCI-DSS">PCI-DSS v4.0</option>
                </select>
            </div>
            
            <div class="file-input-wrapper">
                <input type="file" id="assessmentFile" class="file-input" accept=".xlsx">
                <label for="assessmentFile" class="file-input-label" id="assessmentLabel">
                    <div class="file-icon">📊</div>
                    <div>Click to select or drag & drop assessment file (.xlsx)</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">Excel file with control assessments</div>
                </label>
            </div>
            <div id="assessmentFileList" class="file-list"></div>
        </div>
        
        <!-- Upload Evidence Section -->
        <div class="upload-section">
            <h2><span class="step-number">2</span> Upload Evidence Files</h2>
            <div class="file-input-wrapper">
                <input type="file" id="evidenceFiles" class="file-input" multiple accept=".pdf,.ppt,.pptx,.xlsx,.docx">
                <label for="evidenceFiles" class="file-input-label" id="evidenceLabel">
                    <div class="file-icon">📁</div>
                    <div>Click to select or drag & drop evidence files</div>
                    <div style="color: #666; font-size: 14px; margin-top: 5px;">PDF, PPT, XLSX, DOCX - Multiple files allowed</div>
                </label>
            </div>
            <div id="evidenceFileList" class="file-list"></div>
        </div>
        
        <!-- Analysis Options -->
        <div class="upload-section">
            <h2><span class="step-number">3</span> Analysis Options</h2>
            <div class="analysis-options">
                <div class="option-group">
                    <label for="maxControls">Maximum controls to analyze:</label>
                    <input type="number" id="maxControls" value="20" min="1" max="100">
                </div>
            </div>
            <div style="text-align: center;">
                <button id="analyzeBtn" class="btn" disabled>
                    <span>🔍</span> Analyze Assessment
                </button>
                <button id="statusBtn" class="btn btn-secondary" onclick="checkStatus()">
                    <span>ℹ️</span> System Status
                </button>
            </div>
        </div>
        
        <!-- Loading Section -->
        <div id="loadingSection" class="loading">
            <div class="spinner"></div>
            <h3>Analyzing Assessment with Mistral 7B{% if backend %} ({{ backend }}){% endif %}</h3>
            <p style="color: #666; margin-top: 10px;">This may take a few minutes depending on the number of controls...</p>
            <div class="progress-bar">
                <div class="progress-fill" id="progressFill">0%</div>
            </div>
            <p id="progressText" style="color: #666; margin-top: 10px;">Initializing analysis...</p>
        </div>
        
        <!-- Error/Success Messages -->
        <div id="errorMessage" class="error-message"></div>
        <div id="successMessage" class="success-message"></div>
        
        <!-- Results Section -->
        <div id="resultsSection" class="results-section">
            <h2>📊 Analysis Results</h2>
            
            <div class="risk-score-display">
                <h3>Overall Risk Score</h3>
                <div id="riskScore" class="risk-score">-</div>
                <p style="color: #666;">Lower is better (0-100 scale)</p>
            </div>
            
            <div id="summaryGrid" class="summary-grid"></div>
            
            <div class="charts-grid" id="chartsGrid"></div>
            
            <div class="tabs">
                <button class="tab active" onclick="showTab('detailed')">Detailed Analysis</button>
                <button class="tab" onclick="showTab('highrisk')">High Risk Controls</button>
                <button class="tab" onclick="showTab('summary')">Executive Summary</button>
            </div>
            
            <div id="detailed" class="tab-content active">
                <h3>Detailed Control Analysis</h3>
                <div id="detailedResults"></div>
            </div>
            
            <div id="highrisk" class="tab-content">
                <h3>High Priority Controls Requiring Immediate Attention</h3>
                <div id="highRiskResults"></div>
            </div>
            
            <div id="summary" class="tab-content">
                <h3>Executive Summary</h3>
                <div id="executiveSummary"></div>
            </div>
            
            <div style="text-align: center; margin-top: 30px;">
                <button id="downloadReportBtn" class="btn">
                    <span>📥</span> Download Full Report
                </button>
            </div>
        </div>
    </div>
    
    <script>
        let sessionId = null;
        let assessmentUploaded = false;
        let evidenceUploaded = false;
        
        // Drag and drop functionality
        function setupDragAndDrop(inputId, labelId) {
            const label = document.getElementById(labelId);
            const input = document.getElementById(inputId);
            
            ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {
                label.addEventListener(eventName, preventDefaults, false);
            });
            
            function preventDefaults(e) {
                e.preventDefault();
                e.stopPropagation();
            }
            
            ['dragenter', 'dragover'].forEach(eventName => {
                label.addEventListener(eventName, highlight, false);
            });
            
            ['dragleave', 'drop'].forEach(eventName => {
                label.addEventListener(eventName, unhighlight, false);
            });
            
            function highlight(e) {
                label.classList.add('drag-over');
            }
            
            function unhighlight(e) {
                label.classList.remove('drag-over');
            }
            
            label.addEventListener('drop', handleDrop, false);
            
            function handleDrop(e) {
                const dt = e.dataTransfer;
                const files = dt.files;
                input.files = files;
                input.dispatchEvent(new Event('change', { bubbles: true }));
            }
        }
        
        setupDragAndDrop('assessmentFile', 'assessmentLabel');
        setupDragAndDrop('evidenceFiles', 'evidenceLabel');
        
        // File upload handlers
        document.getElementById('assessmentFile').addEventListener('change', async (e) => {
            const file = e.target.files[0];
            if (!file) return;
            
            const formData = new FormData();
            formData.append('assessment', file);
            formData.append('framework', document.getElementById('framework').value);
            
            try {
                const response = await fetch('/upload_assessment', {
                    method: 'POST',
                    body: formData
                });
                
                const data = await response.json();
                
                if (data.success) {
                    sessionId = data.session_id;
                    assessmentUploaded = true;
                    
                    document.getElementById('assessmentFileList').style.display = 'block';
                    document.getElementById('assessmentFileList').innerHTML = `
                        <div class="file-item success">
                            <div class="file-info">
                                <span>📊 ${file.name}</span>
                                <span class="file-size">${formatFileSize(file.size)}</span>
                            </div>
                            <span style="color: #4caf50;">✓ ${data.controls_count} controls found</span>
                        </div>
                    `;
                    
                    if (data.sample_control) {
                        console.log('Sample control:', data.sample_control);
                    }
                    
                    checkAnalyzeButton();
                    showSuccess(`Assessment uploaded successfully! Found ${data.controls_count} controls.`);
                } else {
                    showError(data.error || 'Failed to upload assessment');
                }
            } catch (error) {
                showError('Error uploading assessment: ' + error.message);
            }
        });
        
        document.getElementById('evidenceFiles').addEventListener('change', async (e) => {
            const files = e.target.files;
            if (!files.length) return;
            
            if (!sessionId) {
                showError('Please upload assessment file first');
                return;
            }
            
            const formData = new FormData();
            for (let i = 0; i < files.length; i++) {
                formData.append(`evidence_${i}`, files[i]);
            }
            
            try {
                const response = await fetch('/upload_evidence', {
                    method: 'POST',
                    body: formData
                });
                
                const data = await response.json();
                
                if (data.success) {
                    evidenceUploaded = true;
                    
                    const fileListHtml = data.uploaded_files.map(file => {
                        const icon = getFileIcon(file.filename);
                        return `
                            <div class="file-item success">
                                <div class="file-info">
                                    <span>${icon} ${file.filename}</span>
                                    <span class="file-size">${formatFileSize(file.size)}</span>
                                </div>
                                <span style="color: #4caf50;">✓</span>
                            </div>
                        `;
                    }).join('');
                    
                    document.getElementById('evidenceFileList').style.display = 'block';
                    document.getElementById('evidenceFileList').innerHTML = fileListHtml;
                    
                    checkAnalyzeButton();
                    showSuccess(`${data.total_files} evidence files uploaded successfully!`);
                } else {
                    showError(data.error || 'Failed to upload evidence files');
                }
            } catch (error) {
                showError('Error uploading evidence: ' + error.message);
            }
        });
        
        // Analyze button handler
        document.getElementById('analyzeBtn').addEventListener('click', async () => {
            document.getElementById('loadingSection').style.display = 'block';
            document.getElementById('resultsSection').style.display = 'none';
            hideMessages();
            
            const maxControls = document.getElementById('maxControls').value;
            
            try {
                const response = await fetch('/analyze', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify({ max_controls: maxControls })
                });
                
                const data = await response.json();
                
                if (data.success) {
                    displayResults(data);
                    showSuccess('Analysis completed successfully!');
                } else {
                    showError(data.error || 'Analysis failed');
                }
            } catch (error) {
                showError('Error during analysis: ' + error.message);
            } finally {
                document.getElementById('loadingSection').style.display = 'none';
            }
        });
        
        // Download report handler
        document.getElementById('downloadReportBtn').addEventListener('click', () => {
            if (sessionId) {
                window.location.href = `/download_report/${sessionId}`;
            }
        });
        
        // Helper functions
        function checkAnalyzeButton() {
            const btn = document.getElementById('analyzeBtn');
            btn.disabled = !(assessmentUploaded && evidenceUploaded);
        }
        
        function showError(message) {
            const errorDiv = document.getElementById('errorMessage');
            errorDiv.textContent = '⚠️ ' + message;
            errorDiv.style.display = 'block';
            setTimeout(() => errorDiv.style.display = 'none', 7000);
        }
        
        function showSuccess(message) {
            const successDiv = document.getElementById('successMessage');
            successDiv.textContent = '✅ ' + message;
            successDiv.style.display = 'block';
            setTimeout(() => successDiv.style.display = 'none', 5000);
        }
        
        function hideMessages() {
            document.getElementById('errorMessage').style.display = 'none';
            document.getElementById('successMessage').style.display = 'none';
        }
        
        function formatFileSize(bytes) {
            if (bytes === 0) return '0 Bytes';
            const k = 1024;
            const sizes = ['Bytes', 'KB', 'MB', 'GB'];
            const i = Math.floor(Math.log(bytes) / Math.log(k));
            return parseFloat((bytes / Math.pow(k, i)).toFixed(2)) + ' ' + sizes[i];
        }
        
        function getFileIcon(filename) {
            const ext = filename.split('.').pop().toLowerCase();
            const icons = {
                'pdf': '📄',
                'xlsx': '📊',
                'docx': '📝',
                'ppt': '📊',
                'pptx': '📊'
            };
            return icons[ext] || '📎';
        }
        
        function showTab(tabName) {
            // Hide all tabs
            document.querySelectorAll('.tab-content').forEach(content => {
                content.classList.remove('active');
            });
            document.querySelectorAll('.tab').forEach(tab => {
                tab.classList.remove('active');
            });
            
            // Show selected tab
            document.getElementById(tabName).classList.add('active');
            event.target.classList.add('active');
        }
        
        function getRiskScoreClass(score) {
            if (score >= 70) return 'critical';
            if (score >= 50) return 'high';
            if (score >= 30) return 'medium';
            return 'low';
        }
        
        function displayResults(data) {
            // Display risk score
            const riskScore = data.risk_metrics.overall_risk_score;
            const riskScoreElement = document.getElementById('riskScore');
            riskScoreElement.textContent = riskScore + '%';
            riskScoreElement.className = 'risk-score ' + getRiskScoreClass(riskScore);
            
            // Display summary
            const summaryHtml = `
                <div class="summary-card">
                    <h3>Total Controls</h3>
                    <div class="value">${data.summary.total_controls}</div>
                    <div class="sub-value">In Assessment</div>
                </div>
                <div class="summary-card">
                    <h3>Analyzed</h3>
                    <div class="value">${data.summary.analyzed_controls}</div>
                    <div class="sub-value">By Mistral 7B</div>
                </div>
                <div class="summary-card">
                    <h3>Compliance Rate</h3>
                    <div class="value" style="color: #4caf50;">${data.risk_metrics.compliance_percentage}%</div>
                    <div class="sub-value">Controls Compliant</div>
                </div>
                <div class="summary-card">
                    <h3>Critical Risks</h3>
                    <div class="value" style="color: #d32f2f;">${data.risk_metrics.risk_distribution.Critical || 0}</div>
                    <div class="sub-value">Immediate Action</div>
                </div>
            `;
            document.getElementById('summaryGrid').innerHTML = summaryHtml;
            
            // Display charts placeholder
            const chartsHtml = `
                <div class="chart-container">
                    <div class="chart-title">Risk Distribution</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Critical: ${data.risk_metrics.risk_distribution.Critical || 0} | 
                        High: ${data.risk_metrics.risk_distribution.High || 0} | 
                        Medium: ${data.risk_metrics.risk_distribution.Medium || 0} | 
                        Low: ${data.risk_metrics.risk_distribution.Low || 0}
                    </div>
                </div>
                <div class="chart-container">
                    <div class="chart-title">Compliance Status</div>
                    <div style="text-align: center; padding: 20px; color: #666;">
                        Compliant: ${data.risk_metrics.compliance_distribution.Compliant || 0} | 
                        Partial: ${data.risk_metrics.compliance_distribution['Partially Compliant'] || 0} | 
                        Non-Compliant: ${data.risk_metrics.compliance_distribution['Non-Compliant'] || 0}
                    </div>
                </div>
            `;
            document.getElementById('chartsGrid').innerHTML = chartsHtml;
            
            // Display detailed results
            const resultsHtml = data.results.map(result => {
                const riskClass = result.risk_level.toLowerCase().replace(' ', '-');
                const validityClass = result.evidence_validity.toLowerCase().replace(' ', '-');
                const complianceClass = result.compliance_status.toLowerCase().replace(' ', '-');
                
                return `
                    <div class="control-result ${riskClass}">
                        <div class="control-header">
                            <div class="control-title">
                                <h4>${result.control_name || 'Unnamed Control'}</h4>
                                <div class="control-id">ID: ${result.control_id}</div>
                            </div>
                            <div class="badges">
                                <span class="badge ${validityClass}">${result.evidence_validity}</span>
                                <span class="badge ${complianceClass}">${result.compliance_status}</span>
                                <span class="badge ${riskClass}">${result.risk_level}</span>
                            </div>
                        </div>
                        
                        <div class="control-details">
                            <div class="detail-section">
                                <h5>Requirement</h5>
                                <p>${result.requirement}</p>
                            </div>
                            
                            <div class="detail-section">
                                <h5>Supplier Response</h5>
                                <p>${result.supplier_answer}</p>
                            </div>
                            
                            ${result.key_findings.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Key Findings</h5>
                                    <ul class="detail-list">
                                        ${result.key_findings.map(f => `<li>${f}</li>`).join('')}
                                    </ul>
                                </div>
                            ` : ''}
                            
                            ${result.risks.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Identified Risks</h5>
                                    <ul class="detail-list">
                                        ${result.risks.map(r => `<li>${r}</li>`).join('')}
                                    </ul>
                                </div>
                            ` : ''}
                            
                            ${result.remediation.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Remediation Recommendations</h5>
                                    <ul class="detail-list">
                                        ${result.remediation.map(r => `<li>${r}</li>`).join('')}
                                    </ul>
                                </div>
                            ` : ''}
                            
                            ${result.evidence_references.length > 0 ? `
                                <div class="detail-section">
                                    <h5>Evidence Sources</h5>
                                    <ul class="detail-list">
                                        ${result.evidence_references.map(r => `<li>${r}</li>`).join('')}
                                    </ul>
                                </div>
                            ` : ''}
                            
                            <div class="detail-section">
                                <h5>Analysis Confidence</h5>
                                <div class="confidence-meter">
                                    <div class="confidence-fill" style="width: ${result.confidence_score * 100}%"></div>
                                </div>
                                <p style="font-size: 12px; color: #666; margin-top: 5px;">
                                    ${(result.confidence_score * 100).toFixed(0)}% confidence in this assessment
                                </p>
                            </div>
                        </div>
                    </div>
                `;
            }).join('');
            
            document.getElementById('detailedResults').innerHTML = resultsHtml;
            
            // Display high risk controls
            const highRiskControls = data.results.filter(r => 
                r.risk_level === 'Critical' || r.risk_level === 'High'
            );
            
            if (highRiskControls.length > 0) {
                const highRiskHtml = highRiskControls.map(result => {
                    const riskClass = result.risk_level.toLowerCase();
                    return `
                        <div class="control-result ${riskClass}">
                            <h4>${result.control_name} (${result.control_id})</h4>
                            <p><strong>Risk Level:</strong> ${result.risk_level}</p>
                            <p><strong>Compliance Status:</strong> ${result.compliance_status}</p>
                            <p><strong>Primary Risk:</strong> ${result.risks[0] || 'Not specified'}</p>
                            <p><strong>Immediate Action:</strong> ${result.remediation[0] || 'Review required'}</p>
                        </div>
                    `;
                }).join('');
                document.getElementById('highRiskResults').innerHTML = highRiskHtml;
            } else {
                document.getElementById('highRiskResults').innerHTML = 
                    '<p style="color: #4caf50; text-align: center; padding: 40px;">✅ No critical or high risk controls identified!</p>';
            }
            
            // Display executive summary
            const executiveSummaryHtml = `
                <div style="padding: 20px;">
                    <h4>Assessment Overview</h4>
                    <p>This cyber security assessment was conducted using the <strong>${data.summary.framework}</strong> framework. 
                    A total of <strong>${data.summary.analyzed_controls}</strong> controls were analyzed out of 
                    <strong>${data.summary.total_controls}</strong> controls in the assessment.</p>
                    
                    <h4 style="margin-top: 20px;">Key Metrics</h4>
                    <ul>
                        <li>Overall Risk Score: <strong>${data.risk_metrics.overall_risk_score}%</strong> 
                            (${getRiskScoreClass(data.risk_metrics.overall_risk_score)} risk)</li>
                        <li>Compliance Rate: <strong>${data.risk_metrics.compliance_percentage}%</strong></li>
                        <li>Controls Requiring Immediate Attention: <strong>${data.risk_metrics.high_priority_controls.length}</strong></li>
                    </ul>
                    
                    <h4 style="margin-top: 20px;">Risk Summary</h4>
                    <p>The assessment identified ${data.risk_metrics.risk_distribution.Critical || 0} critical risks, 
                    ${data.risk_metrics.risk_distribution.High || 0} high risks, 
                    ${data.risk_metrics.risk_distribution.Medium || 0} medium risks, and 
                    ${data.risk_metrics.risk_distribution.Low || 0} low risks.</p>
                    
                    <h4 style="margin-top: 20px;">Recommendations</h4>
                    <p>Priority should be given to addressing the ${data.risk_metrics.high_priority_controls.length} 
                    high-priority controls that are both high-risk and non-compliant. These controls pose the greatest 
                    risk to the organization's security posture.</p>
                </div>
            `;
            document.getElementById('executiveSummary').innerHTML = executiveSummaryHtml;
            
            document.getElementById('resultsSection').style.display = 'block';
        }
        
        // Check system status
        async function checkStatus() {
            try {
                const response = await fetch('/system_status');
                const status = await response.json();
                
                let message = `System Status:\n\n`;
                message += `Backend: ${status.backend}\n`;
                message += `Model: ${status.model}\n`;
                message += `Frameworks: ${status.frameworks.join(', ')}\n`;
                
                if (status.backend === 'Ollama' && status.available_models) {
                    message += `\nAvailable Ollama Models:\n`;
                    status.available_models.forEach(model => {
                        message += `- ${model}\n`;
                    });
                }
                
                alert(message);
            } catch (error) {
                alert('Error checking system status: ' + error.message);
            }
        }
        
        // Simulate progress during analysis
        let progressInterval;
        document.getElementById('analyzeBtn').addEventListener('click', () => {
            let progress = 0;
            const progressFill = document.getElementById('progressFill');
            const progressText = document.getElementById('progressText');
            
            progressInterval = setInterval(() => {
                progress += Math.random() * 15;
                if (progress > 90) progress = 90;
                
                progressFill.style.width = progress + '%';
                progressFill.textContent = Math.round(progress) + '%';
                
                if (progress < 30) {
                    progressText.textContent = 'Loading Mistral 7B model...';
                } else if (progress < 60) {
                    progressText.textContent = 'Analyzing controls and evidence...';
                } else {
                    progressText.textContent = 'Generating risk assessments...';
                }
            }, 1000);
        });
        
        // Clear progress interval when analysis completes
        const originalShowSuccess = showSuccess;
        showSuccess = function(message) {
            if (progressInterval) {
                clearInterval(progressInterval);
                progressInterval = null;
            }
            originalShowSuccess(message);
        };
    </script>
</body>
</html>
"""

# Save the HTML template
with open('templates/index.html', 'w') as f:
    f.write(html_template)

# Updated requirements file with Ollama support
requirements_content = """# Core requirements (always needed)
flask==3.0.0
pandas==2.1.4
openpyxl==3.1.2
PyPDF2==3.0.1
python-docx==1.1.0
python-pptx==0.6.23
werkzeug==3.0.1
requests==2.31.0
numpy==1.24.3

# Optional: For Transformers mode (if not using Ollama)
# Uncomment these if you want to use Transformers backend:
# torch==2.1.0
# transformers==4.36.0
# accelerate==0.25.0
# sentencepiece==0.1.99
# bitsandbytes==0.41.3
# scipy==1.11.4
# safetensors==0.4.1
"""

# Save requirements.txt
with open('requirements.txt', 'w') as f:
    f.write(requirements_content)

# Enhanced installation and run instructions
instructions = """
# Cyber Assessment Review System with Mistral 7B - Setup Instructions

## 🚀 Automatic Setup (NEW!)

The application now automatically checks and installs missing dependencies, including pip itself!

### Quick Start:
```bash
# Just run this - everything else is automatic!
python app.py
```

The script will:
1. ✅ Check if pip is installed (installs if missing)
2. 📦 Install all required Python packages
3. 🤖 Detect Ollama or fall back to Transformers
4. 🚀 Start the application

### Alternative: Guided Setup
```bash
python setup.py
# Follow the interactive setup wizard
```

## 🔧 What Gets Installed Automatically

### 1. **pip** (if not installed)
- Linux/macOS: Uses `ensurepip` module
- Windows: Downloads and runs `get-pip.py`
- Fallback: Shows manual installation instructions

### 2. **Core Dependencies** (always installed)
- Flask, pandas, openpyxl, PyPDF2
- python-docx, python-pptx, requests, numpy

### 3. **Optional Dependencies** (only if Ollama not found)
- PyTorch (with CUDA support on Linux/Windows)
- Transformers, accelerate, bitsandbytes
- scipy, safetensors, sentencepiece

## 🎯 Manual Setup Options

### Option 1: Using Ollama (Recommended)

#### 1. Install pip (if needed):
**Ubuntu/Debian:**
```bash
sudo apt-get update
sudo apt-get install python3-pip
```

**Fedora:**
```bash
sudo dnf install python3-pip
```

**macOS:**
```bash
# pip usually comes with Python on macOS
python3 -m ensurepip --upgrade
```

**Windows:**
```bash
# Download get-pip.py from https://bootstrap.pypa.io/get-pip.py
python get-pip.py
```

#### 2. Install Ollama:
**Windows:**
Download from: https://ollama.com/download/windows

**macOS/Linux:**
```bash
curl -fsSL https://ollama.com/install.sh | sh
```

#### 3. Get Mistral model:
```bash
ollama pull mistral:7b-instruct
```

#### 4. Run the application:
```bash
python app.py
```

### Option 2: Using Transformers (Advanced)

```bash
# The app will detect missing packages and install automatically
python app.py
```

## 🛡️ Security & Permissions

### Linux/macOS Users:
If you get permission errors:
```bash
# Option 1: Use user install
python -m pip install --user package_name

# Option 2: Use virtual environment (recommended)
python -m venv venv
source venv/bin/activate
pip install -r requirements.txt
```

### Windows Users:
Run as Administrator if you encounter permission issues.

## 📊 System Requirements:

### Minimum:
- Python 3.8+ (with pip)
- 8GB RAM
- 10GB disk space
- Internet connection (first run only)

### Recommended:
- Python 3.10+
- 16-32GB RAM
- NVIDIA GPU (optional)
- 20GB disk space

## ⚡ Performance Tips:

### First Run:
- pip installation: ~1 minute
- Package installation: 5-10 minutes
- Model download: 5-30 minutes (depending on internet speed)

### Subsequent Runs:
- Instant start (all dependencies cached)
- No internet required

## 🛠️ Troubleshooting:

### pip Installation Issues:
```bash
# Linux - if apt/dnf fails:
wget https://bootstrap.pypa.io/get-pip.py
python3 get-pip.py --user

# macOS - if ensurepip fails:
brew install python3

# Windows - if automatic install fails:
1. Download https://bootstrap.pypa.io/get-pip.py
2. Open Command Prompt as Administrator
3. Run: python get-pip.py
```

### Permission Errors:
```bash
# Use --user flag
python -m pip install --user package_name

# Or use virtual environment
python -m venv myenv
source myenv/bin/activate  # Windows: myenv\\Scripts\\activate
```

### Dependency Conflicts:
```bash
# Clear pip cache
pip cache purge

# Force reinstall
pip install --force-reinstall package_name

# Use specific versions
pip install package_name==version
```

### Skip Automatic Installation:
```bash
# If you want to handle dependencies manually
export SKIP_DEP_CHECK=1
python app.py
```

## 📝 Configuration:

### Change Model:
```python
# In app.py, modify:
reviewer = CyberAssessmentReviewer(model_name="phi3:medium")    # Smaller
reviewer = CyberAssessmentReviewer(model_name="mixtral:8x7b")   # Larger
```

### Force Backend:
```python
# Force Ollama only
reviewer = CyberAssessmentReviewer(use_ollama=True)

# Force Transformers only  
reviewer = CyberAssessmentReviewer(use_ollama=False)
```

## 🔐 What Happens During Auto-Install:

1. **Check Python & pip**: Ensures basic requirements
2. **Scan Imports**: Detects missing packages
3. **Check Ollama**: Determines which packages needed
4. **Install Packages**: Uses pip to install missing deps
5. **Verify Installation**: Confirms all packages work
6. **Start Application**: Launches the web interface

## 📚 File Structure:
```
cyber-assessment-review/
├── app.py                # Main application (auto-installs deps)
├── setup.py             # Interactive setup wizard
├── requirements.txt     # Python dependencies list
├── README.md           # Project documentation
├── QUICKSTART.md       # Quick start guide
├── sessions/           # Temporary data (auto-created)
├── uploads/            # User uploads (auto-created)
└── models/             # Model cache (auto-created)
```

## 💡 Tips:

- **First Time?** Just run `python app.py` - it handles everything!
- **No Internet?** Use Ollama mode with pre-downloaded models
- **Limited RAM?** Use smaller models (phi3:medium)
- **Fast GPU?** Use larger models (mixtral:8x7b)

## 🎯 Best Practices:

1. **Use Virtual Environment**:
   ```bash
   python -m venv venv
   source venv/bin/activate  # Windows: venv\\Scripts\\activate
   python app.py
   ```

2. **Pre-download Models** (for offline use):
   ```bash
   ollama pull mistral:7b-instruct
   ollama pull phi3:medium
   ```

3. **Monitor First Run**:
   - Watch for any installation errors
   - Note which backend is selected
   - Check available disk space

## ❓ FAQ:

**Q: What if I don't have admin rights?**
A: Use `--user` flag or virtual environments

**Q: Can I use this offline after setup?**
A: Yes! After first run, no internet needed

**Q: How do I update packages later?**
A: Run `pip install --upgrade -r requirements.txt`

**Q: What if auto-install fails?**
A: Set `SKIP_DEP_CHECK=1` and install manually

Enjoy your AI-powered cyber assessment reviews!
"""

# Save instructions
with open('SETUP_INSTRUCTIONS.md', 'w') as f:
    f.write(instructions)

# Update the quickstart guide
quickstart = """
# 🚀 Cyber Assessment Review - Quick Start

## Option 1: Automatic Setup (Easiest)
```bash
python app.py
# Dependencies will be installed automatically!
```

## Option 2: Using Setup Script
```bash
python setup.py
# Follow the guided setup process
```

## Option 3: Manual Ollama Setup
```bash
# 1. Install Ollama
curl -fsSL https://ollama.com/install.sh | sh

# 2. Get Mistral model
ollama pull mistral:7b-instruct

# 3. Install Python deps & run
pip install -r requirements.txt
python app.py
```

## First Time?
The app will:
1. ✅ Check all dependencies
2. 📦 Install missing packages
3. 🤖 Detect Ollama or use Transformers
4. 🚀 Start the web interface

Open http://localhost:5000 and start analyzing!

## Need Help?
- Check system status with the "System Status" button
- See full docs in SETUP_INSTRUCTIONS.md
- Ollama issues? Run: `ollama serve`
"""

with open('QUICKSTART.md', 'w') as f:
    f.write(quickstart)

# Create a setup script for easier installation
setup_script = """#!/usr/bin/env python3
# setup.py - Easy setup script for Cyber Assessment Review System

import sys
import subprocess
import os
import platform

def print_banner():
    print("=" * 60)
    print("🛡️  Cyber Assessment Review System Setup")
    print("=" * 60)

def check_python_version():
    print("\\n🐍 Checking Python version...")
    version = sys.version_info
    if version.major == 3 and version.minor >= 8:
        print(f"✅ Python {version.major}.{version.minor}.{version.micro} - OK")
        return True
    else:
        print(f"❌ Python {version.major}.{version.minor}.{version.micro} - Too old")
        print("   Please install Python 3.8 or newer")
        return False

def check_pip():
    print("\\n📦 Checking for pip...")
    try:
        result = subprocess.run([sys.executable, '-m', 'pip', '--version'], 
                              capture_output=True, text=True)
        if result.returncode == 0:
            print("✅ pip is installed")
            return True
    except:
        pass
    
    print("❌ pip not found")
    print("\\nAttempting to install pip...")
    
    try:
        # Try ensurepip first
        subprocess.run([sys.executable, '-m', 'ensurepip', '--upgrade'], check=True)
        print("✅ pip installed successfully")
        return True
    except:
        print("\\n📥 To install pip manually:")
        
        if platform.system() == "Linux":
            print("   Ubuntu/Debian: sudo apt-get install python3-pip")
            print("   Fedora: sudo dnf install python3-pip")
            print("   Arch: sudo pacman -S python-pip")
        elif platform.system() == "Darwin":  # macOS
            print("   brew install python3")
            print("   or: python3 -m ensurepip")
        elif platform.system() == "Windows":
            print("   Download: https://bootstrap.pypa.io/get-pip.py")
            print("   Run: python get-pip.py")
        
        return False

def check_ollama():
    print("\\n🤖 Checking for Ollama...")
    try:
        result = subprocess.run(['ollama', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            print("✅ Ollama is installed")
            
            # Check if Ollama is running
            import requests
            try:
                response = requests.get("http://localhost:11434/api/tags", timeout=2)
                if response.status_code == 200:
                    print("✅ Ollama is running")
                    return True, True
                else:
                    print("⚠️  Ollama is installed but not running")
                    print("   Run: ollama serve")
                    return True, False
            except:
                print("⚠️  Ollama is installed but not running")
                print("   Run: ollama serve")
                return True, False
    except FileNotFoundError:
        print("❌ Ollama not found")
        print("\\n📥 To install Ollama:")
        
        if platform.system() == "Darwin":  # macOS
            print("   curl -fsSL https://ollama.com/install.sh | sh")
        elif platform.system() == "Linux":
            print("   curl -fsSL https://ollama.com/install.sh | sh")
        elif platform.system() == "Windows":
            print("   Download from: https://ollama.com/download/windows")
        
        return False, False

def check_mistral_model():
    print("\\n🎯 Checking for Mistral model...")
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        if response.status_code == 200:
            models = response.json().get("models", [])
            if any(model["name"] == "mistral:7b-instruct" for model in models):
                print("✅ Mistral 7B model is available")
                return True
            else:
                print("❌ Mistral 7B model not found")
                print("   Run: ollama pull mistral:7b-instruct")
                return False
    except:
        print("⚠️  Cannot check models (Ollama not running)")
        return False

def install_dependencies():
    print("\\n📦 Installing Python dependencies...")
    
    # Create virtual environment if it doesn't exist
    if not os.path.exists('venv'):
        print("Creating virtual environment...")
        subprocess.run([sys.executable, '-m', 'venv', 'venv'])
        print("✅ Virtual environment created")
        print("\\n⚠️  Please activate the virtual environment and run this script again:")
        if platform.system() == "Windows":
            print("   venv\\\\Scripts\\\\activate")
        else:
            print("   source venv/bin/activate")
        return False
    
    # Install requirements
    print("Installing requirements...")
    subprocess.run([sys.executable, '-m', 'pip', 'install', '--upgrade', 'pip'])
    subprocess.run([sys.executable, '-m', 'pip', 'install', '-r', 'requirements.txt'])
    print("✅ Dependencies installed")
    return True

def main():
    print_banner()
    
    # Check Python version
    if not check_python_version():
        sys.exit(1)
    
    # Check pip
    if not check_pip():
        sys.exit(1)
    
    # Check Ollama
    ollama_installed, ollama_running = check_ollama()
    
    # Check Mistral model if Ollama is running
    if ollama_running:
        check_mistral_model()
    
    # Install Python dependencies
    deps_installed = install_dependencies()
    
    print("\\n" + "=" * 60)
    print("📋 Setup Summary:")
    print("=" * 60)
    
    if ollama_installed and ollama_running and deps_installed:
        print("\\n✅ Everything is ready! You can now run:")
        print("   python app.py")
    else:
        print("\\n⚠️  Some steps need completion:")
        if not ollama_installed:
            print("   1. Install Ollama (see instructions above)")
        elif not ollama_running:
            print("   1. Start Ollama: ollama serve")
        if not deps_installed:
            print("   2. Activate venv and run setup again")
        else:
            print("   2. Pull Mistral model: ollama pull mistral:7b-instruct")
            print("   3. Run the app: python app.py")
    
    print("\\n💡 For more help, see SETUP_INSTRUCTIONS.md")
    print("=" * 60)

if __name__ == "__main__":
    main()
"""

# Save setup script
with open('setup.py', 'w') as f:
    f.write(setup_script)

# Make it executable on Unix-like systems
try:
    os.chmod('setup.py', 0o755)
except:
    pass

# Updated requirements file with clearer organization
requirements_content = """# Core Requirements - Always needed
# =====================================
flask==3.0.0
pandas==2.1.4
openpyxl==3.1.2
PyPDF2==3.0.1
python-docx==1.1.0
python-pptx==0.6.23
werkzeug==3.0.1
requests==2.31.0
numpy==1.24.3

# Optional: Transformers Backend
# =====================================
# Only needed if not using Ollama
# Uncomment the lines below if you want to use Transformers:

# torch==2.1.0
# transformers==4.36.0
# accelerate==0.25.0
# sentencepiece==0.1.99
# bitsandbytes==0.41.3
# scipy==1.11.4
# safetensors==0.4.1

# Note: For PyTorch with CUDA support, install manually:
# pip install torch torchvision torchaudio --index-url https://download.pytorch.org/whl/cu118
"""

# Save requirements.txt
with open('requirements.txt', 'w') as f:
    f.write(requirements_content)

# Create a simple run script
run_script = """#!/usr/bin/env python3
# run.py - Simple runner for the Cyber Assessment Review System

import subprocess
import sys
import os

print("🚀 Starting Cyber Assessment Review System...")

# Check if running in virtual environment
if sys.prefix == sys.base_prefix:
    print("⚠️  Warning: Not running in a virtual environment")
    print("   Consider using: source venv/bin/activate")

# Run the main application
try:
    subprocess.run([sys.executable, 'app.py'])
except KeyboardInterrupt:
    print("\\\\n👋 Shutting down gracefully...")
except Exception as e:
    print(f"❌ Error: {e}")
"""

with open('run.py', 'w') as f:
    f.write(run_script)

with open('run.py', 'w') as f:
    f.write(run_script)

try:
    os.chmod('run.py', 0o755)
except:
    pass

# Create a comprehensive README
readme_content = """# 🛡️ Cyber Assessment Review System

AI-powered security control validation using Mistral 7B, running entirely on your local machine.

## 🚀 Quick Start

```bash
# Option 1: Automatic setup (easiest)
python app.py
# Dependencies will be installed automatically!

# Option 2: Guided setup
python setup.py

# Option 3: Manual setup
pip install -r requirements.txt
python app.py
```

## 📋 Features

- **Local AI Processing**: All data stays on your machine
- **Multi-Framework Support**: NIST, ISO27001, SOC2, CIS, PCI-DSS
- **Smart Evidence Analysis**: Processes PDF, DOCX, PPTX, XLSX files
- **Risk Assessment**: Automated risk scoring and prioritization
- **Comprehensive Reports**: Excel reports with detailed findings
- **Easy Model Management**: Ollama integration for simple setup

## 🔧 Requirements

- Python 3.8+
- 8-16GB RAM
- 10GB disk space
- Optional: NVIDIA GPU for faster processing

## 📖 Documentation

- [Quick Start Guide](QUICKSTART.md)
- [Full Setup Instructions](SETUP_INSTRUCTIONS.md)
- [Troubleshooting Guide](SETUP_INSTRUCTIONS.md#-troubleshooting)

## 🤖 Supported Models

Default: Mistral 7B Instruct

Alternative models (via Ollama):
- `phi3:medium` - Smaller, faster
- `mixtral:8x7b` - Larger, more capable
- `qwen2.5:7b` - Good for technical analysis

## 🛠️ Usage

1. Upload your cyber assessment Excel file
2. Upload supporting evidence documents
3. Select compliance framework
4. Click Analyze
5. Download comprehensive report

## 📊 Output

The system generates:
- Overall risk scores
- Compliance percentages
- Control-by-control analysis
- Evidence validation
- Remediation recommendations
- Executive summary

## 🔐 Security

- 100% local processing
- No data leaves your machine
- No cloud dependencies
- Session data auto-cleanup

## 📝 License

MIT License - See LICENSE file for details

## 🙏 Acknowledgments

Built with:
- Mistral AI's language models
- Ollama for model management
- Flask for web interface
- Transformers by Hugging Face

---

For help, check the documentation or open an issue.
"""

with open('README.md', 'w') as f:
    f.write(readme_content)

if __name__ == '__main__':
    print("=" * 60)
    print("Cyber Assessment Review System - Powered by Mistral 7B")
    print("=" * 60)
    print("\nInitializing system...")
    
    if reviewer.use_ollama:
        print("✅ Using Ollama for model management")
        print(f"   Model: {reviewer.model_name}")
    else:
        print("📦 Using Transformers backend")
        print("   Note: First run will download Mistral 7B model (~13GB)")
    
    print("\n🌐 Once started, open http://localhost:5000 in your browser")
    print("\n💡 Tip: For easier setup, install Ollama from https://ollama.com")
    print("=" * 60)
    
    app.run(debug=True, host='0.0.0.0', port=5000)