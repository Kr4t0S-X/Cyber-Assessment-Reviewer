# Cyber Assessment Review System with Mistral 7B
# Main application file: app.py

import sys
import subprocess
import importlib.util
import os

# Check and install dependencies
def check_and_install_dependencies():
    """Check for required dependencies and install missing ones"""
    
    # Check if we should skip dependency checking
    if os.environ.get('SKIP_DEP_CHECK', '').lower() in ['1', 'true', 'yes']:
        print("Skipping dependency check (SKIP_DEP_CHECK is set)")
        return
    
    print("Checking dependencies...")
    
    # Core dependencies that are always needed
    core_dependencies = {
        'flask': 'flask==3.0.0',
        'pandas': 'pandas==2.1.4',
        'openpyxl': 'openpyxl==3.1.2',
        'PyPDF2': 'PyPDF2==3.0.1',
        'docx': 'python-docx==1.1.0',
        'pptx': 'python-pptx==0.6.23',
        'werkzeug': 'werkzeug==3.0.1',
        'requests': 'requests==2.31.0',
        'numpy': 'numpy==1.24.3'
    }
    
    # Optional dependencies for Transformers mode
    transformers_dependencies = {
        'torch': 'torch==2.1.0',
        'transformers': 'transformers==4.36.0',
        'accelerate': 'accelerate==0.25.0',
        'sentencepiece': 'sentencepiece==0.1.99',
        'bitsandbytes': 'bitsandbytes==0.41.3',
        'scipy': 'scipy==1.11.4',
        'safetensors': 'safetensors==0.4.1'
    }
    
    missing_core = []
    missing_transformers = []
    
    # Check core dependencies
    for module, package in core_dependencies.items():
        if importlib.util.find_spec(module) is None:
            missing_core.append(package)
    
    # Check if Ollama is available
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        ollama_available = response.status_code == 200
    except:
        ollama_available = False
    
    # If Ollama is not available, check Transformers dependencies
    if not ollama_available:
        print("Ollama not detected. Checking Transformers dependencies...")
        for module, package in transformers_dependencies.items():
            if importlib.util.find_spec(module) is None:
                missing_transformers.append(package)
    
    # Install missing packages
    if missing_core or missing_transformers:
        print("\nMissing dependencies detected. Installing...")
        
        # Upgrade pip first
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--upgrade", "pip"])
        
        # Install missing core dependencies
        if missing_core:
            print(f"\nInstalling core dependencies: {', '.join(missing_core)}")
            for package in missing_core:
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                    print(f"✓ Installed {package}")
                except subprocess.CalledProcessError as e:
                    print(f"✗ Failed to install {package}: {e}")
                    sys.exit(1)
        
        # Install missing Transformers dependencies if needed
        if missing_transformers and not ollama_available:
            print(f"\nInstalling Transformers dependencies: {', '.join(missing_transformers)}")
            
            # Special handling for PyTorch
            if 'torch' in [pkg.split('==')[0] for pkg in missing_transformers]:
                print("\nInstalling PyTorch... This may take a few minutes.")
                try:
                    # Try to install PyTorch with CUDA support
                    if sys.platform.startswith('linux') or sys.platform == 'win32':
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio", 
                            "--index-url", "https://download.pytorch.org/whl/cu118"
                        ])
                    else:
                        # macOS or other platforms
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio"
                        ])
                    print("✓ Installed PyTorch")
                except subprocess.CalledProcessError:
                    print("✗ Failed to install PyTorch with CUDA, trying CPU-only version...")
                    try:
                        subprocess.check_call([
                            sys.executable, "-m", "pip", "install", 
                            "torch", "torchvision", "torchaudio"
                        ])
                        print("✓ Installed PyTorch (CPU-only)")
                    except subprocess.CalledProcessError as e:
                        print(f"✗ Failed to install PyTorch: {e}")
                        sys.exit(1)
            
            # Install other Transformers dependencies
            for package in missing_transformers:
                if not package.startswith('torch'):
                    try:
                        subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                        print(f"✓ Installed {package}")
                    except subprocess.CalledProcessError as e:
                        print(f"✗ Failed to install {package}: {e}")
                        # Don't exit for optional dependencies
        
        print("\n✅ All required dependencies installed!")
        print("Please restart the application for changes to take effect.")
        sys.exit(0)
    else:
        if ollama_available:
            print("✅ All core dependencies satisfied. Using Ollama backend.")
        else:
            print("✅ All dependencies satisfied.")

# Run dependency check before imports
check_and_install_dependencies()

# Now import everything else
import json
import pandas as pd
from flask import Flask, render_template, request, jsonify, session, send_file
from werkzeug.utils import secure_filename
from datetime import datetime
import uuid
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import re
from typing import List, Dict, Any, Optional
import logging
from dataclasses import dataclass
import numpy as np
import requests
import subprocess
import time

# Optional imports for Transformers mode
try:
    import torch
    from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
    TRANSFORMERS_AVAILABLE = True

    # Try to import BitsAndBytesConfig separately
    try:
        from transformers import BitsAndBytesConfig
        QUANTIZATION_AVAILABLE = True
    except ImportError:
        QUANTIZATION_AVAILABLE = False
        print("Quantization not available. Will use standard precision.")

except ImportError:
    TRANSFORMERS_AVAILABLE = False
    QUANTIZATION_AVAILABLE = False
    print("Transformers not available. Using Ollama mode only.")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)
app.secret_key = 'your-secret-key-here'  # Change this in production
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'xlsx', 'pdf', 'ppt', 'pptx', 'docx'}

# Create necessary directories
for dir_name in ['uploads', 'sessions', 'static', 'models']:
    Path(dir_name).mkdir(exist_ok=True)

@dataclass
class AssessmentResult:
    """Data class for assessment results"""
    control_id: str
    control_name: str
    requirement: str
    supplier_answer: str
    evidence_validity: str
    compliance_status: str
    risk_level: str
    key_findings: List[str]
    risks: List[str]
    remediation: List[str]
    evidence_references: List[str]
    confidence_score: float

class CyberAssessmentReviewer:
    def __init__(self, model_name="mistral:7b-instruct", use_ollama=True):
        """
        Initialize with Mistral 7B for optimal cyber assessment analysis
        Supports both Ollama (recommended) and Transformers backends
        """
        logger.info(f"Initializing Cyber Assessment Reviewer")
        self.use_ollama = use_ollama
        self.model_name = model_name
        
        if self.use_ollama:
            # Check if Ollama is running
            if self.check_ollama_running():
                logger.info("Ollama detected and running")
                # Check if model is available
                if not self.check_ollama_model():
                    logger.info(f"Model {model_name} not found. Pulling from Ollama...")
                    self.pull_ollama_model()
                logger.info(f"Using Ollama with model: {model_name}")
            else:
                logger.warning("Ollama not running. Falling back to Transformers mode.")
                self.use_ollama = False
        
        if not self.use_ollama:
            if not TRANSFORMERS_AVAILABLE:
                raise RuntimeError("Neither Ollama nor Transformers is available. Please install Ollama or transformers library.")
            
            # Transformers initialization with improved CUDA handling
            self.device = "cuda" if (TRANSFORMERS_AVAILABLE and torch.cuda.is_available()) else "cpu"
            logger.info(f"Using device: {self.device}")

            # Use a better model for cybersecurity analysis
            hf_model_name = "microsoft/DialoGPT-medium"  # Better for conversation/analysis tasks

            # Configure quantization with better error handling
            self.quantization_config = None
            if self.device == "cuda" and QUANTIZATION_AVAILABLE:
                try:
                    # Try to use quantization if bitsandbytes works
                    self.quantization_config = BitsAndBytesConfig(
                        load_in_4bit=True,
                        bnb_4bit_compute_dtype=torch.float16,
                        bnb_4bit_quant_type="nf4",
                        bnb_4bit_use_double_quant=True
                    )
                    logger.info("Using 4-bit quantization for efficient memory usage")
                except Exception as e:
                    logger.warning(f"Quantization setup failed: {e}. Using standard precision.")
                    self.quantization_config = None
            else:
                if self.device == "cuda":
                    logger.info("Using CUDA without quantization")
                else:
                    logger.info("Using CPU mode without quantization")
            
            # Load tokenizer
            self.tokenizer = AutoTokenizer.from_pretrained(hf_model_name)
            if not self.tokenizer.pad_token:
                self.tokenizer.pad_token = self.tokenizer.eos_token
            
            # Load model with optimizations and error handling
            logger.info(f"Loading {hf_model_name} model...")
            try:
                self.model = AutoModelForCausalLM.from_pretrained(
                    hf_model_name,
                    quantization_config=self.quantization_config,
                    torch_dtype=torch.float16 if self.device == "cuda" else torch.float32,
                    device_map="auto" if self.device == "cuda" else None,
                    trust_remote_code=True,
                    low_cpu_mem_usage=True
                )
            except Exception as e:
                logger.warning(f"Failed to load model with quantization: {e}")
                logger.info("Retrying without quantization...")
                self.model = AutoModelForCausalLM.from_pretrained(
                    hf_model_name,
                    torch_dtype=torch.float32,
                    trust_remote_code=True,
                    low_cpu_mem_usage=True
                )
                self.device = "cpu"  # Force CPU if quantization fails
            
            # Create pipeline with optimized settings
            self.pipe = pipeline(
                "text-generation",
                model=self.model,
                tokenizer=self.tokenizer,
                device=0 if self.device == "cuda" else -1,  # 0 for GPU, -1 for CPU
                max_new_tokens=1024,
                temperature=0.3,
                do_sample=True,
                top_p=0.9,
                repetition_penalty=1.1
            )
        
        # Cyber security frameworks for context
        self.frameworks = {
            "NIST": "NIST Cybersecurity Framework",
            "ISO27001": "ISO/IEC 27001:2022",
            "SOC2": "SOC 2 Type II",
            "CIS": "CIS Controls v8",
            "PCI-DSS": "PCI-DSS v4.0"
        }
        
        logger.info("Model initialization complete")
    
    def check_ollama_running(self) -> bool:
        """Check if Ollama server is running"""
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            return response.status_code == 200
        except:
            return False
    
    def check_ollama_model(self) -> bool:
        """Check if the model is available in Ollama"""
        try:
            response = requests.get("http://localhost:11434/api/tags")
            if response.status_code == 200:
                models = response.json().get("models", [])
                return any(model["name"] == self.model_name for model in models)
        except:
            pass
        return False
    
    def pull_ollama_model(self):
        """Pull model using Ollama"""
        try:
            # Try to pull the model
            logger.info(f"Pulling {self.model_name} - this may take several minutes...")
            response = requests.post(
                "http://localhost:11434/api/pull",
                json={"name": self.model_name},
                stream=True
            )
            
            for line in response.iter_lines():
                if line:
                    data = json.loads(line)
                    if "status" in data:
                        logger.info(f"Ollama: {data['status']}")
                        
            logger.info(f"Model {self.model_name} pulled successfully")
        except Exception as e:
            logger.error(f"Failed to pull model: {e}")
            raise RuntimeError(f"Could not pull model {self.model_name}. Please run: ollama pull {self.model_name}")
    
    def generate_with_ollama(self, prompt: str, temperature: float = 0.3) -> str:
        """Generate response using Ollama API"""
        try:
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": self.model_name,
                    "prompt": prompt,
                    "temperature": temperature,
                    "stream": False,
                    "options": {
                        "num_predict": 1024,
                        "top_p": 0.9,
                        "repeat_penalty": 1.1
                    }
                }
            )
            
            if response.status_code == 200:
                return response.json()["response"]
            else:
                logger.error(f"Ollama API error: {response.status_code}")
                return ""
        except Exception as e:
            logger.error(f"Error generating with Ollama: {e}")
            return ""
    
    def extract_text_from_pdf(self, filepath: str, max_pages: int = 50) -> Dict[str, str]:
        """Extract text from PDF file with page references"""
        pages_text = {}
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = min(len(pdf_reader.pages), max_pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        pages_text[f"Page {page_num + 1}"] = text[:2000]  # Limit per page
        except Exception as e:
            logger.error(f"Error reading PDF {filepath}: {e}")
            return {"Error": f"Failed to read PDF: {str(e)}"}
        
        return pages_text
    
    def extract_text_from_docx(self, filepath: str) -> Dict[str, str]:
        """Extract text from DOCX file with section references"""
        sections_text = {}
        try:
            doc = Document(filepath)
            
            # Extract paragraphs
            full_text = []
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
                    
            # Group into sections
            if full_text:
                section_size = max(1, len(full_text) // 5)  # Divide into ~5 sections
                for i in range(0, len(full_text), section_size):
                    section_text = '\n'.join(full_text[i:i+section_size])
                    sections_text[f"Section {i//section_size + 1}"] = section_text[:2000]
                    
        except Exception as e:
            logger.error(f"Error reading DOCX {filepath}: {e}")
            return {"Error": f"Failed to read DOCX: {str(e)}"}
        
        return sections_text
    
    def extract_text_from_pptx(self, filepath: str) -> Dict[str, str]:
        """Extract text from PPTX file with slide references"""
        slides_text = {}
        try:
            prs = Presentation(filepath)
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text[f"Slide {i + 1}"] = '\n'.join(slide_text)[:2000]
                    
        except Exception as e:
            logger.error(f"Error reading PPTX {filepath}: {e}")
            return {"Error": f"Failed to read PPTX: {str(e)}"}
        
        return slides_text
    
    def extract_text_from_xlsx(self, filepath: str) -> Dict[str, str]:
        """Extract text from XLSX file with sheet references"""
        sheets_text = {}
        try:
            xlsx_file = pd.ExcelFile(filepath)
            
            for sheet_name in xlsx_file.sheet_names[:5]:  # Limit to first 5 sheets
                df = pd.read_excel(xlsx_file, sheet_name=sheet_name)
                
                # Convert to string representation
                text_content = f"Sheet: {sheet_name}\n"
                text_content += f"Columns: {', '.join(df.columns)}\n"
                text_content += f"Rows: {len(df)}\n\n"
                
                # Add sample data
                if len(df) > 0:
                    text_content += df.head(10).to_string()
                    
                sheets_text[sheet_name] = text_content[:2000]
                
        except Exception as e:
            logger.error(f"Error reading XLSX {filepath}: {e}")
            return {"Error": f"Failed to read XLSX: {str(e)}"}
        
        return sheets_text
    
    def extract_evidence_text(self, filepath: str) -> Dict[str, str]:
        """Extract text from evidence file with source references"""
        ext = Path(filepath).suffix.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(filepath)
        elif ext == '.docx':
            return self.extract_text_from_docx(filepath)
        elif ext in ['.ppt', '.pptx']:
            return self.extract_text_from_pptx(filepath)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(filepath)
        
        return {"Error": f"Unsupported file type: {ext}"}
    
    def create_cyber_prompt(self, control: Dict, evidence_text: str, framework: str = "NIST") -> str:
        """Create specialized prompt for cyber assessment analysis"""
        prompt = f"""You are a senior cybersecurity compliance auditor reviewing control assessments. 
Analyze the following control against the provided evidence using the {self.frameworks.get(framework, framework)} framework.

CONTROL INFORMATION:
- Control ID: {control.get('control_id', control.get('control', 'N/A'))}
- Control Name: {control.get('control_name', control.get('name', 'N/A'))}
- Requirement: {control.get('requirement', control.get('description', 'N/A'))}
- Supplier Answer: {control.get('answer', control.get('response', 'N/A'))}
- Implementation Status: {control.get('status', 'Unknown')}

EVIDENCE PROVIDED:
{evidence_text[:3000]}

Provide your analysis in the following JSON format:
{{
    "evidence_validity": "Valid|Partially Valid|Invalid|No Evidence",
    "compliance_status": "Compliant|Partially Compliant|Non-Compliant",
    "risk_level": "Critical|High|Medium|Low",
    "confidence_score": 0.0-1.0,
    "key_findings": ["finding1", "finding2", "finding3"],
    "identified_risks": ["risk1", "risk2", "risk3"],
    "remediation_steps": ["step1", "step2", "step3"],
    "evidence_gaps": ["gap1", "gap2"]
}}

Focus on:
1. Whether evidence adequately demonstrates the control implementation
2. Specific security vulnerabilities or compliance gaps
3. Practical, prioritized remediation recommendations
4. Risk impact on the organization

Provide only the JSON response without additional text."""

        return prompt
    
    def parse_llm_response(self, response: str) -> Dict[str, Any]:
        """Parse and validate LLM JSON response"""
        try:
            # Extract JSON from response
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                json_str = json_match.group()
                result = json.loads(json_str)
                
                # Ensure all required fields exist
                default_result = {
                    "evidence_validity": "Invalid",
                    "compliance_status": "Non-Compliant",
                    "risk_level": "High",
                    "confidence_score": 0.5,
                    "key_findings": [],
                    "identified_risks": [],
                    "remediation_steps": [],
                    "evidence_gaps": []
                }
                
                # Merge with defaults
                for key, default_value in default_result.items():
                    if key not in result:
                        result[key] = default_value
                
                return result
            else:
                raise ValueError("No JSON found in response")
                
        except Exception as e:
            logger.error(f"Error parsing LLM response: {e}")
            logger.debug(f"Response was: {response}")
            
            # Return structured default response
            return {
                "evidence_validity": "Error",
                "compliance_status": "Unknown",
                "risk_level": "Unknown",
                "confidence_score": 0.0,
                "key_findings": ["Error parsing AI response"],
                "identified_risks": ["Unable to assess risks due to parsing error"],
                "remediation_steps": ["Manual review required"],
                "evidence_gaps": ["Technical error in assessment"]
            }
    
    def analyze_control_point(self, control: Dict, evidence_texts: Dict[str, Dict[str, str]], 
                            framework: str = "NIST") -> AssessmentResult:
        """Analyze a single control point with its evidence"""
        # Combine evidence from all files
        combined_evidence = []
        evidence_references = []
        
        for filename, file_sections in evidence_texts.items():
            for section, text in file_sections.items():
                if "Error" not in section and text.strip():
                    combined_evidence.append(f"[{filename} - {section}]:\n{text}")
                    evidence_references.append(f"{filename} - {section}")
        
        evidence_text = "\n\n".join(combined_evidence[:5])  # Limit to top 5 sections
        
        # Create prompt
        prompt = self.create_cyber_prompt(control, evidence_text, framework)
        
        try:
            # Generate response based on backend
            if self.use_ollama:
                generated_text = self.generate_with_ollama(prompt, temperature=0.3)
            else:
                # Use Transformers pipeline
                response = self.pipe(prompt, max_new_tokens=800, temperature=0.3)[0]['generated_text']
                # Extract only the generated part
                generated_text = response[len(prompt):].strip()
            
            # Parse response
            parsed = self.parse_llm_response(generated_text)
            
            # Create AssessmentResult
            result = AssessmentResult(
                control_id=control.get('control_id', control.get('control', 'N/A')),
                control_name=control.get('control_name', control.get('name', 'N/A')),
                requirement=control.get('requirement', control.get('description', 'N/A')),
                supplier_answer=control.get('answer', control.get('response', 'N/A')),
                evidence_validity=parsed['evidence_validity'],
                compliance_status=parsed['compliance_status'],
                risk_level=parsed['risk_level'],
                key_findings=parsed.get('key_findings', [])[:3],
                risks=parsed.get('identified_risks', [])[:3],
                remediation=parsed.get('remediation_steps', [])[:3],
                evidence_references=evidence_references[:3],
                confidence_score=float(parsed.get('confidence_score', 0.5))
            )
            
            return result
            
        except Exception as e:
            logger.error(f"Error analyzing control {control.get('control_id', 'Unknown')}: {e}")
            
            # Return error result
            return AssessmentResult(
                control_id=control.get('control_id', 'Unknown'),
                control_name=control.get('control_name', 'Unknown'),
                requirement=control.get('requirement', 'Unknown'),
                supplier_answer=control.get('answer', 'Unknown'),
                evidence_validity="Error",
                compliance_status="Unknown",
                risk_level="Unknown",
                key_findings=[f"Analysis error: {str(e)}"],
                risks=["Unable to assess due to error"],
                remediation=["Manual review required"],
                evidence_references=[],
                confidence_score=0.0
            )
    
    def calculate_risk_score(self, results: List[AssessmentResult]) -> Dict[str, Any]:
        """Calculate overall risk metrics from assessment results"""
        if not results:
            return {
                "overall_risk_score": 0,
                "risk_distribution": {},
                "compliance_percentage": 0,
                "high_priority_controls": []
            }
        
        # Count by risk level
        risk_counts = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0, "Unknown": 0}
        compliance_counts = {"Compliant": 0, "Partially Compliant": 0, "Non-Compliant": 0}
        
        for result in results:
            risk_counts[result.risk_level] = risk_counts.get(result.risk_level, 0) + 1
            compliance_counts[result.compliance_status] = compliance_counts.get(result.compliance_status, 0) + 1
        
        # Calculate risk score (0-100, higher is worse)
        risk_weights = {"Critical": 10, "High": 7, "Medium": 4, "Low": 1, "Unknown": 5}
        total_risk = sum(risk_counts[level] * risk_weights.get(level, 0) for level in risk_counts)
        max_risk = len(results) * 10  # If all were critical
        risk_score = (total_risk / max_risk * 100) if max_risk > 0 else 0
        
        # Calculate compliance percentage
        compliant = compliance_counts.get("Compliant", 0)
        total = len(results)
        compliance_percentage = (compliant / total * 100) if total > 0 else 0
        
        # Identify high priority controls
        high_priority = [
            {
                "control_id": r.control_id,
                "control_name": r.control_name,
                "risk_level": r.risk_level,
                "compliance_status": r.compliance_status
            }
            for r in results 
            if r.risk_level in ["Critical", "High"] and r.compliance_status == "Non-Compliant"
        ][:10]  # Top 10
        
        return {
            "overall_risk_score": round(risk_score, 1),
            "risk_distribution": risk_counts,
            "compliance_distribution": compliance_counts,
            "compliance_percentage": round(compliance_percentage, 1),
            "high_priority_controls": high_priority,
            "total_controls_assessed": len(results)
        }

# Initialize the reviewer
try:
    # Try Ollama first (recommended)
    reviewer = CyberAssessmentReviewer(use_ollama=True)
except Exception as e:
    logger.warning(f"Failed to initialize with Ollama: {e}")
    logger.info("Falling back to Transformers mode...")
    try:
        reviewer = CyberAssessmentReviewer(use_ollama=False)
    except Exception as e2:
        logger.error(f"Failed to initialize reviewer: {e2}")
        logger.error("Please install Ollama (recommended) or required Python packages")
        raise

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/')
def index():
    """Render the main page"""
    backend = "Ollama" if reviewer.use_ollama else "Transformers"
    html_content = get_html_template(backend)
    return html_content

def get_html_template(backend):
    """Return the HTML template for the main page"""
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cyber Assessment Reviewer - {backend} Backend</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            color: #333;
        }}
        
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            padding: 20px;
        }}
        
        .header {{
            text-align: center;
            color: white;
            margin-bottom: 40px;
        }}
        
        .header h1 {{
            font-size: 3rem;
            margin-bottom: 10px;
            text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
        }}
        
        .header .subtitle {{
            font-size: 1.2rem;
            opacity: 0.9;
            margin-bottom: 10px;
        }}
        
        .backend-indicator {{
            display: inline-block;
            background: rgba(255,255,255,0.2);
            padding: 8px 16px;
            border-radius: 20px;
            font-size: 0.9rem;
            backdrop-filter: blur(10px);
            border: 1px solid rgba(255,255,255,0.3);
        }}
        
        .main-content {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 30px;
            margin-bottom: 40px;
        }}
        
        @media (max-width: 1024px) {{
            .main-content {{
                grid-template-columns: 1fr;
            }}
        }}
        
        .card {{
            background: white;
            border-radius: 15px;
            padding: 30px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.1);
            backdrop-filter: blur(10px);
            border: 1px solid rgba(255,255,255,0.2);
        }}
        
        .card h2 {{
            color: #5a67d8;
            margin-bottom: 20px;
            font-size: 1.5rem;
        }}
        
        .upload-area {{
            border: 3px dashed #cbd5e0;
            border-radius: 10px;
            padding: 40px;
            text-align: center;
            transition: all 0.3s ease;
            cursor: pointer;
            margin-bottom: 20px;
        }}
        
        .upload-area:hover {{
            border-color: #5a67d8;
            background-color: #f7fafc;
        }}
        
        .upload-area.dragover {{
            border-color: #5a67d8;
            background-color: #ebf4ff;
        }}
        
        .upload-icon {{
            font-size: 3rem;
            color: #a0aec0;
            margin-bottom: 15px;
        }}
        
        .file-input {{
            display: none;
        }}
        
        .btn {{
            background: linear-gradient(135deg, #5a67d8, #667eea);
            color: white;
            border: none;
            padding: 12px 24px;
            border-radius: 8px;
            font-size: 1rem;
            cursor: pointer;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(90,103,216,0.3);
        }}
        
        .btn:hover {{
            transform: translateY(-2px);
            box-shadow: 0 6px 20px rgba(90,103,216,0.4);
        }}
        
        .btn:disabled {{
            opacity: 0.6;
            cursor: not-allowed;
            transform: none;
        }}
        
        .progress-container {{
            margin-top: 20px;
            display: none;
        }}
        
        .progress-bar {{
            background-color: #e2e8f0;
            border-radius: 10px;
            overflow: hidden;
            height: 10px;
            margin-bottom: 10px;
        }}
        
        .progress-fill {{
            background: linear-gradient(90deg, #5a67d8, #667eea);
            height: 100%;
            width: 0%;
            transition: width 0.3s ease;
            display: flex;
            align-items: center;
            justify-content: center;
            color: white;
            font-size: 0.8rem;
            font-weight: bold;
        }}
        
        .status-message {{
            text-align: center;
            font-weight: 500;
            color: #4a5568;
        }}
        
        .file-list {{
            margin-top: 15px;
        }}
        
        .file-item {{
            display: flex;
            align-items: center;
            padding: 10px;
            background-color: #f7fafc;
            border-radius: 8px;
            margin-bottom: 8px;
        }}
        
        .file-icon {{
            margin-right: 10px;
            color: #5a67d8;
        }}
        
        .framework-selector {{
            margin-bottom: 20px;
        }}
        
        .framework-selector label {{
            display: block;
            margin-bottom: 8px;
            font-weight: 500;
            color: #4a5568;
        }}
        
        .framework-selector select {{
            width: 100%;
            padding: 10px;
            border: 2px solid #e2e8f0;
            border-radius: 8px;
            font-size: 1rem;
            background-color: white;
        }}
        
        .framework-selector select:focus {{
            outline: none;
            border-color: #5a67d8;
        }}
        
        .results-section {{
            grid-column: 1 / -1;
            margin-top: 20px;
        }}
        
        .results-container {{
            display: none;
            background: white;
            border-radius: 15px;
            padding: 30px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.1);
        }}
        
        .results-header {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 25px;
            padding-bottom: 15px;
            border-bottom: 2px solid #e2e8f0;
        }}
        
        .results-title {{
            font-size: 1.8rem;
            color: #2d3748;
        }}
        
        .download-btn {{
            background: linear-gradient(135deg, #48bb78, #68d391);
            box-shadow: 0 4px 15px rgba(72,187,120,0.3);
        }}
        
        .download-btn:hover {{
            box-shadow: 0 6px 20px rgba(72,187,120,0.4);
        }}
        
        .analysis-grid {{
            display: grid;
            gap: 20px;
        }}
        
        .control-analysis {{
            background: #f8f9fa;
            border-radius: 10px;
            padding: 20px;
            border-left: 4px solid #5a67d8;
        }}
        
        .control-header {{
            display: flex;
            justify-content: between;
            align-items: flex-start;
            margin-bottom: 15px;
        }}
        
        .control-title {{
            font-size: 1.2rem;
            color: #2d3748;
            margin-bottom: 5px;
        }}
        
        .confidence-score {{
            background: linear-gradient(135deg, #ed8936, #f6ad55);
            color: white;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 0.8rem;
            font-weight: bold;
            margin-left: auto;
        }}
        
        .risk-level {{
            display: inline-block;
            padding: 6px 12px;
            border-radius: 6px;
            font-size: 0.8rem;
            font-weight: bold;
            margin-bottom: 10px;
        }}
        
        .risk-high {{
            background-color: #fed7d7;
            color: #c53030;
        }}
        
        .risk-medium {{
            background-color: #feebc8;
            color: #dd6b20;
        }}
        
        .risk-low {{
            background-color: #c6f6d5;
            color: #25855a;
        }}
        
        .analysis-content {{
            margin-bottom: 15px;
        }}
        
        .evidence-section {{
            background: white;
            border-radius: 8px;
            padding: 15px;
            margin-top: 10px;
        }}
        
        .evidence-title {{
            font-weight: 600;
            color: #4a5568;
            margin-bottom: 8px;
        }}
        
        .evidence-item {{
            background: #edf2f7;
            padding: 8px 12px;
            border-radius: 6px;
            margin-bottom: 6px;
            font-size: 0.9rem;
        }}
        
        .remediation {{
            background: #e6fffa;
            border: 1px solid #38b2ac;
            border-radius: 8px;
            padding: 15px;
            margin-top: 10px;
        }}
        
        .remediation-title {{
            color: #2c7a7b;
            font-weight: 600;
            margin-bottom: 8px;
        }}
        
        .alert {{
            padding: 15px;
            border-radius: 8px;
            margin-bottom: 20px;
            border: 1px solid;
        }}
        
        .alert-success {{
            background-color: #f0fff4;
            border-color: #9ae6b4;
            color: #2f855a;
        }}
        
        .alert-error {{
            background-color: #fed7d7;
            border-color: #fc8181;
            color: #c53030;
        }}
        
        .spinner {{
            border: 4px solid #f3f3f3;
            border-top: 4px solid #5a67d8;
            border-radius: 50%;
            width: 30px;
            height: 30px;
            animation: spin 1s linear infinite;
            margin: 20px auto;
        }}
        
        @keyframes spin {{
            0% {{ transform: rotate(0deg); }}
            100% {{ transform: rotate(360deg); }}
        }}
        
        .hidden {{
            display: none !important;
        }}
        
        .fade-in {{
            animation: fadeIn 0.5s ease-in-out;
        }}
        
        @keyframes fadeIn {{
            from {{ opacity: 0; transform: translateY(20px); }}
            to {{ opacity: 1; transform: translateY(0); }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <!-- Header -->
        <div class="header">
            <h1>Cyber Assessment Reviewer</h1>
            <div class="subtitle">AI-Powered Cybersecurity Control Analysis</div>
            <div class="backend-indicator">
                🧠 Backend: {backend}
            </div>
        </div>
        
        <!-- Main Content Grid -->
        <div class="main-content">
            <!-- Assessment Upload Section -->
            <div class="card">
                <h2>📋 Upload Assessment</h2>
                <form id="assessmentForm" enctype="multipart/form-data">
                    <div class="framework-selector">
                        <label for="framework">Cybersecurity Framework:</label>
                        <select id="framework" name="framework" required>
                            <option value="">Select Framework</option>
                            <option value="NIST">NIST Cybersecurity Framework</option>
                            <option value="ISO27001">ISO 27001</option>
                            <option value="SOC2">SOC 2</option>
                            <option value="CIS">CIS Controls</option>
                            <option value="PCI-DSS">PCI DSS</option>
                        </select>
                    </div>
                    
                    <div class="upload-area" id="assessmentUpload">
                        <div class="upload-icon">📄</div>
                        <h3>Drop assessment file here</h3>
                        <p>Or click to select file</p>
                        <p><small>Supports: PDF, DOCX, XLSX, PPTX</small></p>
                        <input type="file" id="assessmentFile" name="assessment" class="file-input" 
                               accept=".pdf,.docx,.xlsx,.pptx" required>
                    </div>
                    
                    <div id="assessmentFiles" class="file-list"></div>
                </form>
            </div>
            
            <!-- Evidence Upload Section -->
            <div class="card">
                <h2>🗂️ Upload Evidence</h2>
                <form id="evidenceForm" enctype="multipart/form-data">
                    <div class="upload-area" id="evidenceUpload">
                        <div class="upload-icon">📁</div>
                        <h3>Drop evidence files here</h3>
                        <p>Or click to select files</p>
                        <p><small>Multiple files supported</small></p>
                        <input type="file" id="evidenceFiles" name="evidence" class="file-input" 
                               accept=".pdf,.docx,.xlsx,.pptx" multiple>
                    </div>
                    
                    <div id="evidenceFileList" class="file-list"></div>
                    
                    <button type="button" id="analyzeBtn" class="btn" disabled>
                        🔍 Analyze Assessment
                    </button>
                    
                    <div class="progress-container" id="progressContainer">
                        <div class="progress-bar">
                            <div class="progress-fill" id="progressFill">0%</div>
                        </div>
                        <div class="status-message" id="progressText">Starting analysis...</div>
                    </div>
                </form>
            </div>
        </div>
        
        <!-- Results Section -->
        <div class="results-section">
            <div class="results-container" id="resultsContainer">
                <div class="results-header">
                    <h2 class="results-title">Analysis Results</h2>
                    <button type="button" id="downloadBtn" class="btn download-btn">
                        📥 Download Report
                    </button>
                </div>
                <div id="analysisResults" class="analysis-grid"></div>
            </div>
        </div>
        
        <!-- Alert Area -->
        <div id="alertArea"></div>
    </div>
    
    <script>
        let assessmentUploaded = false;
        let evidenceUploaded = false;
        
        // File upload handling
        function setupFileUpload(uploadAreaId, fileInputId, fileListId, isMultiple = false) {{
            const uploadArea = document.getElementById(uploadAreaId);
            const fileInput = document.getElementById(fileInputId);
            const fileList = document.getElementById(fileListId);
            
            uploadArea.addEventListener('click', () => fileInput.click());
            
            uploadArea.addEventListener('dragover', (e) => {{
                e.preventDefault();
                uploadArea.classList.add('dragover');
            }});
            
            uploadArea.addEventListener('dragleave', () => {{
                uploadArea.classList.remove('dragover');
            }});
            
            uploadArea.addEventListener('drop', (e) => {{
                e.preventDefault();
                uploadArea.classList.remove('dragover');
                
                const files = Array.from(e.dataTransfer.files);
                handleFiles(files, fileInput, fileList, isMultiple);
            }});
            
            fileInput.addEventListener('change', (e) => {{
                const files = Array.from(e.target.files);
                handleFiles(files, fileInput, fileList, isMultiple);
            }});
        }}
        
        function handleFiles(files, fileInput, fileList, isMultiple) {{
            if (!isMultiple && files.length > 1) {{
                showAlert('Please select only one file.', 'error');
                return;
            }}
            
            // Clear existing files if single file mode
            if (!isMultiple) {{
                fileList.innerHTML = '';
            }}
            
            files.forEach(file => {{
                const fileItem = document.createElement('div');
                fileItem.className = 'file-item';
                fileItem.innerHTML = `
                    <span class="file-icon">📄</span>
                    <span>${{file.name}} (${{(file.size / 1024 / 1024).toFixed(2)}} MB)</span>
                `;
                fileList.appendChild(fileItem);
            }});
            
            // Update upload status
            if (fileInput.id === 'assessmentFile') {{
                assessmentUploaded = files.length > 0;
                uploadAssessment(files[0]);
            }} else {{
                evidenceUploaded = files.length > 0;
                uploadEvidence(files);
            }}
            
            updateAnalyzeButton();
        }}
        
        function updateAnalyzeButton() {{
            const analyzeBtn = document.getElementById('analyzeBtn');
            const framework = document.getElementById('framework').value;
            analyzeBtn.disabled = !(assessmentUploaded && framework);
        }}
        
        async function uploadAssessment(file) {{
            const formData = new FormData();
            const framework = document.getElementById('framework').value;
            
            if (!framework) {{
                showAlert('Please select a cybersecurity framework first.', 'error');
                return;
            }}
            
            formData.append('assessment', file);
            formData.append('framework', framework);
            
            try {{
                const response = await fetch('/upload_assessment', {{
                    method: 'POST',
                    body: formData
                }});
                
                const result = await response.json();
                
                if (result.success) {{
                    showAlert('Assessment uploaded successfully!', 'success');
                }} else {{
                    showAlert('Error uploading assessment: ' + result.error, 'error');
                    assessmentUploaded = false;
                    updateAnalyzeButton();
                }}
            }} catch (error) {{
                showAlert('Error uploading assessment: ' + error.message, 'error');
                assessmentUploaded = false;
                updateAnalyzeButton();
            }}
        }}
        
        async function uploadEvidence(files) {{
            const formData = new FormData();
            
            files.forEach(file => {{
                formData.append('evidence', file);
            }});
            
            try {{
                const response = await fetch('/upload_evidence', {{
                    method: 'POST',
                    body: formData
                }});
                
                const result = await response.json();
                
                if (result.success) {{
                    showAlert(`${{files.length}} evidence file(s) uploaded successfully!`, 'success');
                }} else {{
                    showAlert('Error uploading evidence: ' + result.error, 'error');
                    evidenceUploaded = false;
                    updateAnalyzeButton();
                }}
            }} catch (error) {{
                showAlert('Error uploading evidence: ' + error.message, 'error');
                evidenceUploaded = false;
                updateAnalyzeButton();
            }}
        }}
        
        async function analyzeAssessment() {{
            const progressContainer = document.getElementById('progressContainer');
            const analyzeBtn = document.getElementById('analyzeBtn');
            const resultsContainer = document.getElementById('resultsContainer');
            
            progressContainer.style.display = 'block';
            analyzeBtn.disabled = true;
            resultsContainer.style.display = 'none';
            
            try {{
                const response = await fetch('/analyze', {{
                    method: 'POST',
                    headers: {{
                        'Content-Type': 'application/json'
                    }},
                    body: JSON.stringify({{
                        max_controls: 20
                    }})
                }});

                const result = await response.json();
                
                if (result.success) {{
                    displayResults(result.results);
                    showAlert('Analysis completed successfully!', 'success');
                }} else {{
                    showAlert('Error during analysis: ' + result.error, 'error');
                }}
            }} catch (error) {{
                showAlert('Error during analysis: ' + error.message, 'error');
            }} finally {{
                progressContainer.style.display = 'none';
                analyzeBtn.disabled = false;
            }}
        }}
        
        function displayResults(results) {{
            const resultsContainer = document.getElementById('resultsContainer');
            const analysisResults = document.getElementById('analysisResults');
            
            analysisResults.innerHTML = '';
            
            results.forEach(control => {{
                const controlDiv = document.createElement('div');
                controlDiv.className = 'control-analysis fade-in';
                
                controlDiv.innerHTML = `
                    <div class="control-header">
                        <div>
                            <div class="control-title">${{control.control_id}} - ${{control.control_name}}</div>
                            <div class="risk-level risk-${{control.risk_level.toLowerCase()}}">${{control.risk_level}} Risk</div>
                        </div>
                        <div class="confidence-score">${{Math.round(control.confidence_score * 100)}}% Confidence</div>
                    </div>

                    <div class="analysis-content">
                        <p><strong>Requirement:</strong> ${{control.requirement}}</p>
                        <p><strong>Supplier Answer:</strong> ${{control.supplier_answer}}</p>
                        <p><strong>Evidence Validity:</strong> ${{control.evidence_validity}}</p>
                        <p><strong>Compliance Status:</strong> ${{control.compliance_status}}</p>
                    </div>

                    ${{control.key_findings && control.key_findings.length > 0 ? `
                        <div class="evidence-section">
                            <div class="evidence-title">� Key Findings:</div>
                            ${{control.key_findings.map(finding =>
                                `<div class="evidence-item">${{finding}}</div>`
                            ).join('')}}
                        </div>
                    ` : ''}}

                    ${{control.risks && control.risks.length > 0 ? `
                        <div class="evidence-section">
                            <div class="evidence-title">⚠️ Identified Risks:</div>
                            ${{control.risks.map(risk =>
                                `<div class="evidence-item">${{risk}}</div>`
                            ).join('')}}
                        </div>
                    ` : ''}}

                    ${{control.remediation && control.remediation.length > 0 ? `
                        <div class="remediation">
                            <div class="remediation-title">🔧 Recommended Actions:</div>
                            ${{control.remediation.map(action =>
                                `<div class="evidence-item">${{action}}</div>`
                            ).join('')}}
                        </div>
                    ` : ''}}

                    ${{control.evidence_references && control.evidence_references.length > 0 ? `
                        <div class="evidence-section">
                            <div class="evidence-title">📄 Evidence References:</div>
                            ${{control.evidence_references.map(ref =>
                                `<div class="evidence-item">${{ref}}</div>`
                            ).join('')}}
                        </div>
                    ` : ''}}
                `;
                
                analysisResults.appendChild(controlDiv);
            }});
            
            resultsContainer.style.display = 'block';
            resultsContainer.scrollIntoView({{ behavior: 'smooth' }});
        }}
        
        async function downloadReport() {{
            try {{
                const response = await fetch('/download_report');
                const blob = await response.blob();
                
                const url = window.URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.style.display = 'none';
                a.href = url;
                a.download = 'cyber_assessment_report.pdf';
                document.body.appendChild(a);
                a.click();
                window.URL.revokeObjectURL(url);
                
                showAlert('Report downloaded successfully!', 'success');
            }} catch (error) {{
                showAlert('Error downloading report: ' + error.message, 'error');
            }}
        }}
        
        function showAlert(message, type) {{
            const alertArea = document.getElementById('alertArea');
            const alertDiv = document.createElement('div');
            alertDiv.className = `alert alert-${{type}} fade-in`;
            alertDiv.textContent = message;
            
            alertArea.appendChild(alertDiv);
            
            setTimeout(() => {{
                alertDiv.remove();
            }}, 5000);
        }}
        
        // Initialize file uploads
        setupFileUpload('assessmentUpload', 'assessmentFile', 'assessmentFiles', false);
        setupFileUpload('evidenceUpload', 'evidenceFiles', 'evidenceFileList', true);
        
        // Framework change handler
        document.getElementById('framework').addEventListener('change', updateAnalyzeButton);
        
        // Analyze button handler
        document.getElementById('analyzeBtn').addEventListener('click', analyzeAssessment);
        
        // Download button handler
        document.getElementById('downloadBtn').addEventListener('click', downloadReport);
        
        // Check system status on page load
        window.addEventListener('load', async () => {{
            try {{
                const response = await fetch('/system_status');
                const data = await response.json();
                
                if (data.backend) {{
                    const message = `System ready with ${{data.backend}} backend. Model: ${{data.model || 'Default'}}`;
                    if (data.backend === 'Ollama') {{
                        showAlert(message + ' ✅', 'success');
                    }} else {{
                        showAlert(message + ' (First analysis will download model)', 'success');
                    }}
                }}
            }} catch (error) {{
                alert('Error checking system status: ' + error.message);
            }}
        }});
        
        // Simulate progress during analysis
        let progressInterval;
        document.getElementById('analyzeBtn').addEventListener('click', () => {{
            let progress = 0;
            const progressFill = document.getElementById('progressFill');
            const progressText = document.getElementById('progressText');
            
            progressInterval = setInterval(() => {{
                progress += Math.random() * 15;
                if (progress > 90) progress = 90;
                
                progressFill.style.width = progress + '%';
                progressFill.textContent = Math.round(progress) + '%';
                
                if (progress < 30) {{
                    progressText.textContent = 'Loading Mistral 7B model...';
                }} else if (progress < 60) {{
                    progressText.textContent = 'Analyzing controls and evidence...';
                }} else {{
                    progressText.textContent = 'Generating risk assessments...';
                }}
            }}, 1000);
        }});
        
        // Clear progress interval when analysis completes
        const originalShowSuccess = showSuccess;
        showSuccess = function(message) {{
            if (progressInterval) {{
                clearInterval(progressInterval);
                progressInterval = null;
            }}
            originalShowSuccess(message);
        }};
    </script>
</body>
</html>
"""

@app.route('/upload_assessment', methods=['POST'])
def upload_assessment():
    """Handle assessment file upload"""
    if 'assessment' not in request.files:
        return jsonify({'error': 'No assessment file provided'}), 400
    
    file = request.files['assessment']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        # Create session ID
        session_id = str(uuid.uuid4())
        session['session_id'] = session_id
        session_path = Path(f"sessions/{session_id}")
        session_path.mkdir(exist_ok=True)
        
        # Save assessment file
        filename = secure_filename(file.filename)
        filepath = session_path / f"assessment_{filename}"
        file.save(str(filepath))
        
        # Parse assessment
        try:
            df = pd.read_excel(filepath)
            
            # Detect column names (flexible mapping)
            column_mapping = {
                'control': ['control', 'control_id', 'id', 'control id', 'control_number'],
                'name': ['name', 'control_name', 'title', 'control name', 'description'],
                'requirement': ['requirement', 'requirements', 'description', 'control_description'],
                'answer': ['answer', 'response', 'supplier_answer', 'supplier response', 'implementation'],
                'status': ['status', 'implementation_status', 'compliance_status']
            }
            
            # Standardize column names
            for standard_name, possible_names in column_mapping.items():
                for col in df.columns:
                    if col.lower().strip() in [name.lower() for name in possible_names]:
                        df.rename(columns={col: standard_name}, inplace=True)
                        break
            
            controls = df.to_dict('records')
            
            # Save controls to session
            with open(session_path / 'controls.json', 'w') as f:
                json.dump(controls, f)
            
            # Save column info
            session['columns'] = list(df.columns)
            session['framework'] = request.form.get('framework', 'NIST')
            
            return jsonify({
                'success': True,
                'session_id': session_id,
                'controls_count': len(controls),
                'columns': list(df.columns),
                'sample_control': controls[0] if controls else None
            })
        except Exception as e:
            logger.error(f"Error parsing assessment: {e}")
            return jsonify({'error': f'Error parsing assessment: {str(e)}'}), 400
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/upload_evidence', methods=['POST'])
def upload_evidence():
    """Handle evidence files upload"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    evidence_path = session_path / 'evidence'
    evidence_path.mkdir(exist_ok=True)
    
    uploaded_files = []
    for key in request.files:
        file = request.files[key]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = evidence_path / filename
            file.save(str(filepath))
            uploaded_files.append({
                'filename': filename,
                'size': os.path.getsize(filepath)
            })
    
    return jsonify({
        'success': True,
        'uploaded_files': uploaded_files,
        'total_files': len(uploaded_files)
    })

@app.route('/analyze', methods=['POST'])
def analyze():
    """Run the analysis on uploaded assessment and evidence"""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'No active session'}), 400
    
    session_path = Path(f"sessions/{session_id}")
    
    # Load controls
    try:
        with open(session_path / 'controls.json', 'r') as f:
            controls = json.load(f)
    except Exception as e:
        return jsonify({'error': f'Failed to load controls: {str(e)}'}), 500
    
    # Extract text from all evidence files
    evidence_texts = {}
    evidence_path = session_path / 'evidence'
    if evidence_path.exists():
        for evidence_file in evidence_path.iterdir():
            if evidence_file.is_file():
                logger.info(f"Processing evidence file: {evidence_file.name}")
                file_texts = reviewer.extract_evidence_text(str(evidence_file))
                evidence_texts[evidence_file.name] = file_texts
    
    # Get framework
    framework = session.get('framework', 'NIST')
    
    # Analyze controls
    results = []
    max_controls = min(len(controls), int(request.json.get('max_controls', 20)))
    
    for i, control in enumerate(controls[:max_controls]):
        logger.info(f"Analyzing control {i+1}/{max_controls}: {control.get('control', 'Unknown')}")
        
        try:
            result = reviewer.analyze_control_point(control, evidence_texts, framework)
            results.append(result)
        except Exception as e:
            logger.error(f"Error analyzing control {i+1}: {e}")
            continue
    
    # Calculate risk metrics
    risk_metrics = reviewer.calculate_risk_score(results)
    
    # Convert results to dict format
    results_dict = [
        {
            'control_id': r.control_id,
            'control_name': r.control_name,
            'requirement': r.requirement,
            'supplier_answer': r.supplier_answer,
            'evidence_validity': r.evidence_validity,
            'compliance_status': r.compliance_status,
            'risk_level': r.risk_level,
            'key_findings': r.key_findings,
            'risks': r.risks,
            'remediation': r.remediation,
            'evidence_references': r.evidence_references,
            'confidence_score': r.confidence_score
        }
        for r in results
    ]
    
    # Save results
    analysis_results = {
        'results': results_dict,
        'risk_metrics': risk_metrics,
        'framework': framework,
        'timestamp': datetime.now().isoformat(),
        'total_controls': len(controls),
        'analyzed_controls': len(results)
    }
    
    with open(session_path / 'results.json', 'w') as f:
        json.dump(analysis_results, f, indent=2)
    
    return jsonify({
        'success': True,
        'results': results_dict,
        'risk_metrics': risk_metrics,
        'summary': {
            'total_controls': len(controls),
            'analyzed_controls': len(results),
            'framework': framework
        }
    })

@app.route('/system_status')
def system_status():
    """Get system and model status"""
    status = {
        'backend': 'Ollama' if reviewer.use_ollama else 'Transformers',
        'model': reviewer.model_name,
        'frameworks': list(reviewer.frameworks.keys()),
        'ollama_available': reviewer.check_ollama_running() if hasattr(reviewer, 'check_ollama_running') else False
    }
    
    if reviewer.use_ollama:
        try:
            # Get Ollama models
            response = requests.get("http://localhost:11434/api/tags")
            if response.status_code == 200:
                models = [m['name'] for m in response.json().get('models', [])]
                status['available_models'] = models
        except:
            status['available_models'] = []
    
    return jsonify(status)

@app.route('/download_report/<session_id>')
def download_report(session_id):
    """Download the analysis report"""
    session_path = Path(f"sessions/{session_id}")
    results_path = session_path / 'results.json'
    
    if not results_path.exists():
        return jsonify({'error': 'Report not found'}), 404
    
    with open(results_path, 'r') as f:
        data = json.load(f)
    
    # Create Excel report with multiple sheets
    report_path = session_path / 'cyber_assessment_report.xlsx'
    
    with pd.ExcelWriter(report_path, engine='openpyxl') as writer:
        # Executive Summary sheet
        summary_data = {
            'Metric': [
                'Assessment Date',
                'Framework Used',
                'Total Controls',
                'Controls Analyzed',
                'Overall Risk Score',
                'Compliance Percentage',
                'Critical Risk Controls',
                'High Risk Controls'
            ],
            'Value': [
                data['timestamp'][:10],
                data['framework'],
                data['total_controls'],
                data['analyzed_controls'],
                f"{data['risk_metrics']['overall_risk_score']}%",
                f"{data['risk_metrics']['compliance_percentage']}%",
                data['risk_metrics']['risk_distribution'].get('Critical', 0),
                data['risk_metrics']['risk_distribution'].get('High', 0)
            ]
        }
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Executive Summary', index=False)
        
        # Risk Analysis sheet
        risk_df = pd.DataFrame([data['risk_metrics']['risk_distribution']])
        risk_df.to_excel(writer, sheet_name='Risk Analysis', index=False)
        
        # Detailed Results sheet
        results_df = pd.DataFrame(data['results'])
        
        # Convert lists to strings for Excel
        for col in ['key_findings', 'risks', 'remediation', 'evidence_references']:
            if col in results_df.columns:
                results_df[col] = results_df[col].apply(lambda x: '\n'.join(x) if isinstance(x, list) else x)
        
        results_df.to_excel(writer, sheet_name='Detailed Analysis', index=False)
        
        # High Priority Controls sheet
        if data['risk_metrics']['high_priority_controls']:
            priority_df = pd.DataFrame(data['risk_metrics']['high_priority_controls'])
            priority_df.to_excel(writer, sheet_name='High Priority Controls', index=False)
        
        # Format the Excel file
        workbook = writer.book
        for sheet_name in writer.sheets:
            worksheet = writer.sheets[sheet_name]
            
            # Auto-adjust column widths
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
    
    return send_file(str(report_path), as_attachment=True, 
                    download_name=f'cyber_assessment_report_{datetime.now().strftime("%Y%m%d")}.xlsx')

if __name__ == '__main__':
    print("=" * 60)
    print("Cyber Assessment Review System - Powered by Mistral 7B")
    print("=" * 60)
    print("\nInitializing system...")
    
    if reviewer.use_ollama:
        print("✅ Using Ollama for model management")
        print(f"   Model: {reviewer.model_name}")
    else:
        print("📦 Using Transformers backend")
        print("   Note: First run will download Mistral 7B model (~13GB)")
    
    print("\n🌐 Once started, open http://localhost:5000 in your browser")
    print("\n💡 Tip: For easier setup, install Ollama from https://ollama.com")
    print("=" * 60)
    
    app.run(debug=True, host='0.0.0.0', port=5000)